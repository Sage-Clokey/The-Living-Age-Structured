"""
Build research poster for BME 129C Capstone.
48 x 36 inches, landscape. Dark theme matching slide deck.

Layout (left to right, top to bottom):
  HEADER: Title, author, affiliation, abstract
  ROW 1: Layer 1 (topology) | Layer 1b (single-cell) | Layer 2 (economy)
  ROW 2: Layer 2b (FBA/price system) | Layer 3 (trade) | Layers 4-7 (immune/genome/viral)
  FOOTER: Conclusion + key numbers | Methods | References

Usage: python paper/build_poster.py
"""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

PAPER_DIR = Path(__file__).resolve().parent.parent
FIGURES = PAPER_DIR / "figures"
INDIVIDUAL = FIGURES / "individual"
OUT = PAPER_DIR / "deliverables" / "capstone_poster.pptx"

# Colors
GREEN_MID = RGBColor(0x52, 0xb7, 0x88)
GREEN_LIGHT = RGBColor(0x95, 0xd5, 0xb2)
GREEN_DARK = RGBColor(0x1b, 0x43, 0x32)
GOLD = RGBColor(0xe9, 0xc4, 0x6a)
RED = RGBColor(0xe6, 0x39, 0x46)
WHITE = RGBColor(0xff, 0xff, 0xff)
OFF_WHITE = RGBColor(0xf5, 0xf5, 0xf5)
DARK_BG = RGBColor(0x0d, 0x11, 0x17)
PANEL_BG = RGBColor(0x15, 0x1c, 0x25)
GRAY = RGBColor(0x8b, 0x94, 0x9e)
LIGHT_GRAY = RGBColor(0xcc, 0xcc, 0xcc)

# Poster dimensions (48 x 36 inches)
POSTER_W = Inches(48)
POSTER_H = Inches(36)

# Margins and spacing
MARGIN = Inches(0.6)
COL_GAP = Inches(0.4)
ROW_GAP = Inches(0.4)

# 3-column layout
NUM_COLS = 3
USABLE_W = POSTER_W - 2 * MARGIN - (NUM_COLS - 1) * COL_GAP
COL_W = int(USABLE_W / NUM_COLS)

# Row heights
HEADER_H = Inches(5.5)
ROW_H = Inches(12.5)
FOOTER_H = Inches(4.0)


def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_panel(slide, left, top, width, height, color=PANEL_BG):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.color.rgb = GREEN_MID
    shape.line.width = Pt(1.5)
    shape.shadow.inherit = False
    return shape


def add_text(slide, left, top, width, height, text, font_size=18,
             color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = name
    p.alignment = alignment
    return txBox


def add_rich_text(slide, left, top, width, height, lines, default_size=14,
                  default_color=OFF_WHITE):
    """lines = [(text, size, color, bold), ...]"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(lines):
        text = item[0]
        size = item[1] if len(item) > 1 else default_size
        color = item[2] if len(item) > 2 else default_color
        bold = item[3] if len(item) > 3 else False

        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = text
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.font.name = "Calibri"
        p.space_after = Pt(4)
    return txBox


def try_add_picture(slide, *candidates, left=0, top=0, width=None, height=None):
    for path in candidates:
        p = Path(path)
        if p.exists():
            slide.shapes.add_picture(str(p), left, top, width, height)
            return True
    return False


def build():
    prs = Presentation()
    prs.slide_width = POSTER_W
    prs.slide_height = POSTER_H

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)

    # Top green bar
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, POSTER_W, Inches(0.12))
    bar.fill.solid()
    bar.fill.fore_color.rgb = GREEN_MID
    bar.line.fill.background()

    # =====================================================================
    # HEADER
    # =====================================================================
    header_top = Inches(0.3)

    # Title
    add_text(slide, MARGIN, header_top, Inches(32), Inches(1.8),
             "The Price System of the Cell:",
             font_size=72, color=GREEN_MID, bold=True)
    add_text(slide, MARGIN, header_top + Inches(1.6), Inches(32), Inches(1.2),
             "Living Systems as Decentralized Economies",
             font_size=48, color=GOLD, bold=True)

    # Author / affiliation
    add_text(slide, MARGIN, header_top + Inches(3.0), Inches(20), Inches(0.6),
             "Sage Clokey  |  Department of Bioengineering  |  UC Santa Cruz  |  BME 129C Capstone  |  Spring 2026  |  Advisor: R. Dubois",
             font_size=22, color=LIGHT_GRAY)

    # Abstract box (right side of header)
    abstract_left = Inches(33)
    abstract_w = POSTER_W - abstract_left - MARGIN
    add_panel(slide, abstract_left, header_top + Inches(0.2), abstract_w, Inches(4.8))
    add_text(slide, abstract_left + Inches(0.3), header_top + Inches(0.4),
             abstract_w - Inches(0.6), Inches(0.5),
             "Abstract", font_size=24, color=GOLD, bold=True)
    add_rich_text(slide, abstract_left + Inches(0.3), header_top + Inches(1.0),
                  abstract_w - Inches(0.6), Inches(3.8), [
        ("Living things are not machines. There are two types of order: "
         "decentralized and centralized. Centralized order requires uniformity. "
         "Decentralized order requires ordered diversity. Life is decentralized — "
         "the function is in the connections, not the parts.", 16, OFF_WHITE),
        ("", 8),
        ("This study presents quantitative evidence — from network topology, "
         "single-cell transcriptomics, metabolic simulation, genome-scale FBA, "
         "cross-species trade, and cancer genomics — that biology operates as a "
         "decentralized economy coordinated by prices.", 16, OFF_WHITE),
        ("", 8),
        ("Key findings: 19:1 robustness ratio over centralized architectures. "
         "Distributed allocation retains 71% vs 53% centralized under perturbation. "
         "The omniscient planner (FBA) still fails on 30% of real knockouts. "
         "Cancer mutations target price system components — the disease is the "
         "central planning, not the variation.", 16, GREEN_LIGHT),
    ])

    # =====================================================================
    # ROW 1: Layer 1 | Layer 1b | Layer 2
    # =====================================================================
    row1_top = header_top + HEADER_H + Inches(0.3)

    # --- PANEL: Layer 1 — Network Topology ---
    col1_left = MARGIN
    add_panel(slide, col1_left, row1_top, COL_W, ROW_H)

    add_text(slide, col1_left + Inches(0.3), row1_top + Inches(0.3),
             COL_W - Inches(0.6), Inches(0.5),
             "Layer 1: No Master Node", font_size=28, color=GOLD, bold=True)
    add_text(slide, col1_left + Inches(0.3), row1_top + Inches(0.9),
             COL_W - Inches(0.6), Inches(0.4),
             "Network Topology Analysis", font_size=18, color=GREEN_MID)

    try_add_picture(slide,
                    FIGURES / "layer1_topology_annotated.png",
                    FIGURES / "layer1_topology.png",
                    left=col1_left + Inches(0.3), top=row1_top + Inches(1.5),
                    width=COL_W - Inches(0.6), height=Inches(7.0))

    add_rich_text(slide, col1_left + Inches(0.3), row1_top + Inches(8.7),
                  COL_W - Inches(0.6), Inches(3.5), [
        ("No master node — the network has no central controller", 16, OFF_WHITE, True),
        ("Power-law degree distribution (alpha 2.0-2.5) — hubs exist but no hub dominates", 15, OFF_WHITE),
        ("19:1 robustness ratio: bio networks survive 37% node removal vs 1.9% for star graph", 15, GREEN_LIGHT),
        ("Feed-forward loops over-represented — evolved price signals, not accidental", 15, OFF_WHITE),
        ("Hub erosion (WBPA): network actively resists centralization", 15, OFF_WHITE),
        ("Data: RegulonDB, STRING v12, KEGG REST API", 13, GRAY),
    ])

    # --- PANEL: Layer 1b — Single-Cell Economy ---
    col2_left = col1_left + COL_W + COL_GAP
    add_panel(slide, col2_left, row1_top, COL_W, ROW_H)

    add_text(slide, col2_left + Inches(0.3), row1_top + Inches(0.3),
             COL_W - Inches(0.6), Inches(0.5),
             "Layer 1b: Cells as Entrepreneurs", font_size=28, color=GOLD, bold=True)
    add_text(slide, col2_left + Inches(0.3), row1_top + Inches(0.9),
             COL_W - Inches(0.6), Inches(0.4),
             "Single-Cell RNA-seq Economy", font_size=18, color=GREEN_MID)

    try_add_picture(slide,
                    FIGURES / "layer1b_single_cell_economy_annotated.png",
                    FIGURES / "layer1b_single_cell_economy.png",
                    left=col2_left + Inches(0.3), top=row1_top + Inches(1.5),
                    width=COL_W - Inches(0.6), height=Inches(7.0))

    add_rich_text(slide, col2_left + Inches(0.3), row1_top + Inches(8.7),
                  COL_W - Inches(0.6), Inches(3.5), [
        ("Same genome, different output — cells specialize by choice, not command", 16, OFF_WHITE, True),
        ("8 cell types, division of labor without a master cell", 15, OFF_WHITE),
        ("Communication Gini = 0.0 — every cell type signals directly, no gatekeeper", 15, GREEN_LIGHT),
        ("70-90% of communication survives any single cell type removal", 15, GREEN_LIGHT),
        ("Same ligand, different meaning by receiver — subjective value at the molecular level", 15, OFF_WHITE),
        ("Data: 10x Genomics PBMC3k, Scanpy", 13, GRAY),
    ])

    # --- PANEL: Layer 2 — Economic Simulation ---
    col3_left = col2_left + COL_W + COL_GAP
    add_panel(slide, col3_left, row1_top, COL_W, ROW_H)

    add_text(slide, col3_left + Inches(0.3), row1_top + Inches(0.3),
             COL_W - Inches(0.6), Inches(0.5),
             "Layer 2: Distributed Beats Centralized", font_size=28, color=GOLD, bold=True)
    add_text(slide, col3_left + Inches(0.3), row1_top + Inches(0.9),
             COL_W - Inches(0.6), Inches(0.4),
             "Agent-Based Metabolic Simulation", font_size=18, color=GREEN_MID)

    try_add_picture(slide,
                    FIGURES / "layer2_economy_annotated.png",
                    FIGURES / "layer2_economy.png",
                    left=col3_left + Inches(0.3), top=row1_top + Inches(1.5),
                    width=COL_W - Inches(0.6), height=Inches(7.0))

    add_rich_text(slide, col3_left + Inches(0.3), row1_top + Inches(8.7),
                  COL_W - Inches(0.6), Inches(3.5), [
        ("13 pathway agents, no central allocator — the invisible hand finds equilibrium", 16, OFF_WHITE, True),
        ("Production rates oscillate early (price discovery) then converge", 15, OFF_WHITE),
        ("Distributed retains 71% of GDP under perturbation", 15, GREEN_LIGHT),
        ("Centralized retains only 53% — the plan breaks when conditions change", 15, RED),
        ("Distributed wins on all 4 perturbation tests (Hayek, Mises, Kirzner)", 15, OFF_WHITE),
        ("Agents discover optimal rates through local feedback — no planner needed", 15, OFF_WHITE),
    ])

    # =====================================================================
    # ROW 2: Price System | Layer 3 | Immune/Genome/Viral
    # =====================================================================
    row2_top = row1_top + ROW_H + ROW_GAP

    # --- PANEL: Price System of the Cell ---
    add_panel(slide, col1_left, row2_top, COL_W, ROW_H)

    add_text(slide, col1_left + Inches(0.3), row2_top + Inches(0.3),
             COL_W - Inches(0.6), Inches(0.5),
             "The Price System of the Cell", font_size=28, color=GOLD, bold=True)
    add_text(slide, col1_left + Inches(0.3), row2_top + Inches(0.9),
             COL_W - Inches(0.6), Inches(0.4),
             "FBA Shadow Prices + Cancer Mutations", font_size=18, color=GREEN_MID)

    try_add_picture(slide,
                    FIGURES / "price_system_of_the_cell.png",
                    left=col1_left + Inches(0.3), top=row2_top + Inches(1.5),
                    width=COL_W - Inches(0.6), height=Inches(7.0))

    add_rich_text(slide, col1_left + Inches(0.3), row2_top + Inches(8.7),
                  COL_W - Inches(0.6), Inches(3.5), [
        ("The cell has a real price system — and cancer is its destruction", 16, OFF_WHITE, True),
        ("Tier 1: ATP/ADP, NAD+/NADH = cost of capital (intracellular)", 15, OFF_WHITE),
        ("Tier 2: Cytokines, morphogens, O2 = market prices (intercellular)", 15, OFF_WHITE),
        ("Tier 3: mTOR = the entrepreneur reading all prices at once", 15, OFF_WHITE),
        ("Shadow prices shift with context — NADH cheap on glucose, expensive on acetate", 15, GREEN_LIGHT),
        ("Cancer targets price receptors, integrators, and decision makers", 15, RED),
        ("TP53: 96% ovarian, 1% thyroid — the disease is the broken context, not the variant", 15, RED),
        ("Data: iML1515 FBA (COBRApy), TCGA PanCancerAtlas (10,967 samples)", 13, GRAY),
    ])

    # --- PANEL: Layer 3 — Cross-Species Trade ---
    add_panel(slide, col2_left, row2_top, COL_W, ROW_H)

    add_text(slide, col2_left + Inches(0.3), row2_top + Inches(0.3),
             COL_W - Inches(0.6), Inches(0.5),
             "Layer 3: The Tree of Life Is a Trade Network", font_size=28, color=GOLD, bold=True)
    add_text(slide, col2_left + Inches(0.3), row2_top + Inches(0.9),
             COL_W - Inches(0.6), Inches(0.4),
             "Cross-Species Gene Transfer", font_size=18, color=GREEN_MID)

    try_add_picture(slide,
                    FIGURES / "layer3_trade_network_annotated.png",
                    FIGURES / "layer3_trade_network.png",
                    left=col2_left + Inches(0.3), top=row2_top + Inches(1.5),
                    width=COL_W - Inches(0.6), height=Inches(7.0))

    add_rich_text(slide, col2_left + Inches(0.3), row2_top + Inches(8.7),
                  COL_W - Inches(0.6), Inches(3.5), [
        ("No organism does everything — Ricardian comparative advantage at the molecular level", 16, OFF_WHITE, True),
        ("Codon usage distance = trade friction (biological gravity model)", 15, OFF_WHITE),
        ("Within-kingdom: 0.17-0.38 (voluntary exchange succeeds)", 15, GREEN_LIGHT),
        ("Cross-kingdom: 0.65-0.83 (forced exchange destroys information)", 15, RED),
        ("Coral: biomineralization. Spider: silk. Bacteria: cellulose. Each specializes.", 15, OFF_WHITE),
        ("Trade blocs emerge spontaneously — compatible partners naturally trade more", 15, OFF_WHITE),
        ("Data: Kazusa Codon Usage DB, 8 organisms across 4 kingdoms", 13, GRAY),
    ])

    # --- PANEL: Layers 4-7 — Immune, Genome, Convergent, Viral ---
    add_panel(slide, col3_left, row2_top, COL_W, ROW_H)

    add_text(slide, col3_left + Inches(0.3), row2_top + Inches(0.3),
             COL_W - Inches(0.6), Inches(0.5),
             "Layers 4-7: The Pattern Repeats", font_size=28, color=GOLD, bold=True)
    add_text(slide, col3_left + Inches(0.3), row2_top + Inches(0.9),
             COL_W - Inches(0.6), Inches(0.4),
             "Immune, Genome, Convergent Evolution, Viral", font_size=18, color=GREEN_MID)

    # 2x2 grid of smaller figures
    fig_w = (COL_W - Inches(1.0)) / 2
    fig_h = Inches(3.2)
    fig_left1 = col3_left + Inches(0.3)
    fig_left2 = fig_left1 + fig_w + Inches(0.3)
    fig_top1 = row2_top + Inches(1.5)
    fig_top2 = fig_top1 + fig_h + Inches(0.2)

    try_add_picture(slide,
                    FIGURES / "immune_shm_hotspots_annotated.png",
                    FIGURES / "immune_shm_hotspots.png",
                    left=fig_left1, top=fig_top1,
                    width=fig_w, height=fig_h)

    try_add_picture(slide,
                    FIGURES / "genome_tissue_specialization_annotated.png",
                    FIGURES / "genome_tissue_specialization.png",
                    left=fig_left2, top=fig_top1,
                    width=fig_w, height=fig_h)

    try_add_picture(slide,
                    FIGURES / "genome_convergent_evolution_annotated.png",
                    FIGURES / "genome_convergent_evolution.png",
                    left=fig_left1, top=fig_top2,
                    width=fig_w, height=fig_h)

    try_add_picture(slide,
                    FIGURES / "viral_genome_composition.png",
                    left=fig_left2, top=fig_top2,
                    width=fig_w, height=fig_h)

    add_rich_text(slide, col3_left + Inches(0.3), row2_top + Inches(8.7),
                  COL_W - Inches(0.6), Inches(3.5), [
        ("Immune: SHM hotspot targeting at 19:1 enrichment", 15, OFF_WHITE),
        ("Public clonotypes at 10^15 above random chance", 15, OFF_WHITE),
        ("Genome: 20% of genes tissue-exclusive (tau > 0.95)", 15, OFF_WHITE),
        ("CpG>TpG mutations: 15-40x baseline", 15, OFF_WHITE),
        ("Convergent evolution: 35 events, 17 traits, up to 1.5 Byr apart", 15, GREEN_LIGHT),
        ("Viral: 8% of human genome is endogenous retrovirus", 15, OFF_WHITE),
        ("Syncytin: captured viral gene enables mammalian pregnancy", 15, GREEN_LIGHT),
    ])

    # =====================================================================
    # FOOTER: Conclusion | Methods | References
    # =====================================================================
    footer_top = row2_top + ROW_H + ROW_GAP

    # --- Conclusion ---
    concl_w = Inches(20)
    add_panel(slide, col1_left, footer_top, concl_w, FOOTER_H, color=GREEN_DARK)

    add_text(slide, col1_left + Inches(0.3), footer_top + Inches(0.3),
             concl_w - Inches(0.6), Inches(0.5),
             "Conclusion", font_size=28, color=GOLD, bold=True)

    add_rich_text(slide, col1_left + Inches(0.3), footer_top + Inches(1.0),
                  concl_w - Inches(0.6), Inches(2.8), [
        ("Living things are not machines. The function is in the connections, not the parts. "
         "The differences are how adaptation is possible — trial and error directed by distributed knowledge.", 20, WHITE, True),
        ("", 8),
        ("Life is the original decentralized network — the self-assembling internet, the adaptive computer "
         "that built itself. The question is not how to program life. It is how to join the computation. "
         "Cultivate conditions, don't command outcomes. Lead life to lead itself.", 18, GREEN_LIGHT),
    ])

    # --- Methods + References ---
    methods_left = col1_left + concl_w + COL_GAP
    methods_w = POSTER_W - methods_left - MARGIN
    add_panel(slide, methods_left, footer_top, methods_w, FOOTER_H)

    add_text(slide, methods_left + Inches(0.3), footer_top + Inches(0.3),
             methods_w - Inches(0.6), Inches(0.4),
             "Methods & Data Sources", font_size=22, color=GOLD, bold=True)

    add_rich_text(slide, methods_left + Inches(0.3), footer_top + Inches(0.9),
                  methods_w - Inches(0.6), Inches(2.8), [
        ("Python  |  NetworkX  |  Scanpy  |  COBRApy  |  matplotlib", 14, GREEN_LIGHT),
        ("RegulonDB v11  |  STRING v12  |  KEGG REST API  |  CellxGene", 14, OFF_WHITE),
        ("Kazusa Codon Usage DB  |  iML1515 genome-scale model  |  TCGA", 14, OFF_WHITE),
        ("", 6),
        ("Hayek, F.A. (1945) The Use of Knowledge in Society. AER 35(4).", 12, GRAY),
        ("Menger, C. (1871) Principles of Economics.", 12, GRAY),
        ("Mises, L. (1920) Economic Calculation in the Socialist Commonwealth.", 12, GRAY),
        ("Barabasi, A-L & Oltvai, Z. (2004) Network biology. Nature Rev Genetics 5.", 12, GRAY),
        ("Hayek, F.A. (1988) The Fatal Conceit. U Chicago Press.", 12, GRAY),
    ])

    # Bottom green bar
    bar2 = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, POSTER_H - Inches(0.12), POSTER_W, Inches(0.12))
    bar2.fill.solid()
    bar2.fill.fore_color.rgb = GREEN_MID
    bar2.line.fill.background()

    prs.save(str(OUT))
    print(f"Saved: {OUT}")
    print(f"  Poster size: 48 x 36 inches")
    print(f"  1 slide")


if __name__ == "__main__":
    build()
