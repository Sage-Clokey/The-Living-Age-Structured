"""
Build PowerPoint progress report slides for BME 129C Capstone — Week 5.
3 minutes: 0.5 min goal/approach, 1 min accomplishments, 1.5 min recent progress.

Week 5 (April 19–25, 2026):
  - Extended from 3 layers to 7 layers of evidence (immune + genome)
  - Wrote "The Living Architecture" — 55-page synthesis paper (715 lines, 13,625 words)
  - Created immune analysis pipeline: SHM hotspots, V(D)J bias, public clonotypes
  - Created genome analysis pipeline: mutation hotspots, tissue specialization, convergent evolution
  - Generated 16 annotated figures with callout boxes and key findings
  - Wrote "Viruses as Communication" standalone essay with 7 data-backed figures
  - Wrote "Vaccines as Central Planners" essay mapping 5 distortions to Austrian economics
  - Built full 22-slide and 3-minute presentation decks
  - Built DOCX exports for both distributed knowledge paper and Living Architecture

Usage: python paper/build_slides_week5.py
"""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

FIGURES = Path(__file__).resolve().parent.parent / "figures"
OUT = Path(__file__).resolve().parent.parent / "progress_report_slides_week5.pptx"

# Colors
GREEN_MID = RGBColor(0x52, 0xb7, 0x88)
GREEN_LIGHT = RGBColor(0x95, 0xd5, 0xb2)
GOLD = RGBColor(0xe9, 0xc4, 0x6a)
RED = RGBColor(0xe6, 0x39, 0x46)
WHITE = RGBColor(0xff, 0xff, 0xff)
OFF_WHITE = RGBColor(0xf5, 0xf5, 0xf5)
DARK_BG = RGBColor(0x0d, 0x11, 0x17)
GRAY = RGBColor(0x8b, 0x94, 0x9e)
BLUE = RGBColor(0x4e, 0xa8, 0xde)


def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text_box(slide, left, top, width, height, text, font_size=18,
                 color=WHITE, bold=False, alignment=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = "Calibri"
    p.alignment = alignment
    return txBox


def add_bullets(slide, bullets, left, top, width, height,
                font_size=16, color=WHITE, spacing=Pt(8)):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, (text, subs) in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = text
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Calibri"
        p.space_after = spacing
        p.level = 0
        for sub in subs:
            p2 = tf.add_paragraph()
            p2.text = sub
            p2.font.size = Pt(font_size - 2)
            p2.font.color.rgb = GRAY
            p2.font.name = "Calibri"
            p2.space_after = Pt(4)
            p2.level = 1
    return txBox


def add_green_bar(slide):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.33), Inches(0.06))
    shape.fill.solid()
    shape.fill.fore_color.rgb = GREEN_MID
    shape.line.fill.background()


def try_add_picture(slide, *candidates, left=0, top=0, width=None, height=None):
    for path in candidates:
        p = Path(path)
        if p.exists():
            slide.shapes.add_picture(str(p), left, top, width, height)
            return True
    return False


def build():
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    # =========================================================================
    # SLIDE 1: Title
    # =========================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)
    add_green_bar(slide)

    add_text_box(slide, Inches(1), Inches(1.2), Inches(11), Inches(1.2),
                 "The Living Architecture",
                 font_size=44, color=GREEN_MID, bold=True)
    add_text_box(slide, Inches(1), Inches(2.6), Inches(11), Inches(0.8),
                 "Distributed Knowledge as the Design Language of Life",
                 font_size=24, color=GOLD)
    add_text_box(slide, Inches(1), Inches(3.8), Inches(11), Inches(0.5),
                 "BME 129C Capstone  |  Week 5 Progress Report",
                 font_size=16, color=GRAY)
    add_text_box(slide, Inches(1), Inches(4.4), Inches(11), Inches(0.5),
                 "Sage Clokey   |   Spring 2026   |   UC Santa Cruz",
                 font_size=18, color=GREEN_LIGHT)
    add_text_box(slide, Inches(1), Inches(5.5), Inches(11), Inches(0.8),
                 "Seven layers of evidence  |  39 figures  |  3 papers  |  55-page synthesis",
                 font_size=14, color=GRAY, alignment=PP_ALIGN.CENTER)

    notes = slide.notes_slide
    notes.notes_text_frame.text = (
        "Title slide. Week 5 progress report.\n\n"
        "This week the project expanded from 3 computational layers to 7 layers of "
        "biological evidence, added immune system and genome-wide analyses, wrote the "
        "all-encompassing synthesis paper, and produced two standalone essays on vaccines "
        "and viral communication. The figure count went from 21 to 39."
    )

    # =========================================================================
    # SLIDE 2: Progress on Last Week's Next Steps
    # =========================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)
    add_green_bar(slide)

    add_text_box(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.5),
                 "Progress on Last Week's Next Steps",
                 font_size=28, color=GOLD, bold=True)

    bullets = [
        ("\u2714  Statistical validation", [
            "Completed — nonparametric tests integrated into all layer analyses",
            "Mann-Whitney U, bootstrap CIs, KS tests across all comparisons",
        ]),
        ("\u2714  Finalize paper for advisor review", [
            "Exceeded — wrote 55-page synthesis paper (The Living Architecture)",
            "Integrates all 7 layers, design principles, and theological framework",
        ]),
        ("\u25B6  Run full pipeline with live API data", [
            "Deferred — prioritized extending to 7 layers over re-running existing 3",
            "Immune + genome pipelines use published empirical data directly",
        ]),
        ("\u25B6  Begin poster layout", [
            "Deferred to Week 6 — scope expanded significantly",
            "22-slide and 3-minute presentation decks built as intermediate step",
        ]),
    ]
    add_bullets(slide, bullets, Inches(0.8), Inches(1.0), Inches(11.5), Inches(5.5),
                font_size=15)

    add_text_box(slide, Inches(0.8), Inches(6.5), Inches(11.5), Inches(0.5),
                 "\u2714 = completed     \u25B6 = deferred with rationale",
                 font_size=12, color=GRAY)

    notes = slide.notes_slide
    notes.notes_text_frame.text = (
        "PROGRESS ON LAST WEEK'S NEXT STEPS:\n\n"
        "Two of four completed, two deferred.\n\n"
        "Statistical validation: Done. Every comparison now has proper nonparametric "
        "testing — Mann-Whitney U for distributed vs centralized GDP, bootstrap CIs on "
        "robustness metrics.\n\n"
        "Paper finalization: Exceeded target. Instead of polishing the existing paper, "
        "I wrote a new 55-page synthesis that integrates all 7 layers.\n\n"
        "Live API pipeline: Deferred. The immune and genome analyses use published "
        "empirical data directly (IMGT, GTEx, ClinVar), so the API integration was "
        "less critical than extending the evidence base.\n\n"
        "Poster: Deferred to Week 6. The scope expansion this week required all "
        "available time. Built slide decks as intermediate deliverables."
    )

    # =========================================================================
    # SLIDE 3: Accomplishments — Layers 4-7 (1 min)
    # =========================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)
    add_green_bar(slide)

    add_text_box(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.5),
                 "Accomplished: 3 Layers \u2192 7 Layers of Evidence",
                 font_size=28, color=GOLD, bold=True)

    left_bullets = [
        ("Layers 1-3 (existing, from Weeks 1-4):", [
            "1. Network topology — 19:1 robustness ratio",
            "1b. Single-cell economy — Gini = 0.0",
            "2. Economic simulation — distributed 71% vs centralized 53%",
            "2b. FBA — omniscient planner: 70% accuracy ceiling",
            "3. Cross-species trade — codon distance = trade friction",
        ]),
        ("Layer 4: Immune System (NEW)", [
            "SHM hotspot analysis — 19:1 enrichment at WRC motifs",
            "V(D)J segment usage — IGHV3-23 at 10-20x rare segments",
            "Public clonotypes — random P ≈ 10⁻¹⁵, observed: 200 shared",
        ]),
    ]
    add_bullets(slide, left_bullets, Inches(0.5), Inches(1.0), Inches(6), Inches(5.5),
                font_size=13, color=OFF_WHITE)

    right_bullets = [
        ("Layer 5: Whole Genome (NEW)", [
            "CpG C>T mutations: 15-40x baseline from ~1% of contexts",
            "Ti/Tv ratio 2:1 — 4x random expectation",
            "20% of genes: tau > 0.95 — tissue-exclusive expression",
        ]),
        ("Layer 6: Convergent Evolution (NEW)", [
            "Bats + dolphins: 14 identical Prestin substitutions",
            "35 events, 17 traits, 8,000 yrs to 1.5 Byr",
        ]),
        ("Layer 7: Viral Communication (NEW)", [
            "8% of human genome is endogenous retrovirus",
            "Syncytin: captured viral gene enables mammalian pregnancy",
            "10³¹ phages, 10²⁵ gene transfers/day — the first internet",
        ]),
    ]
    add_bullets(slide, right_bullets, Inches(6.8), Inches(1.0), Inches(6), Inches(5.5),
                font_size=13, color=OFF_WHITE)

    add_text_box(slide, Inches(0.8), Inches(6.5), Inches(11.5), Inches(0.5),
                 "Every layer independently confirms: distributed knowledge, not central planning",
                 font_size=14, color=GREEN_MID, bold=True)

    notes = slide.notes_slide
    notes.notes_text_frame.text = (
        "LAYERS 4-7 (45 seconds):\n\n"
        "The project expanded from 3 computational layers to 7 layers of biological evidence.\n\n"
        "Layer 4 — Immune system: Created immune_distributed_knowledge.py. Three analyses: "
        "SHM hotspot targeting (AID targets WRC motifs at 19:1 over coldspots), V(D)J segment "
        "usage bias (IGHV3-23 dominates at 10-20x), and public clonotypes (identical TCR "
        "sequences in unrelated people at rates 10^15 above random chance).\n\n"
        "Layer 5 — Whole genome: Created genome_distributed_knowledge.py. CpG mutation "
        "hotspots, transition/transversion bias, and tissue-specific gene expression showing "
        "extreme division of labor.\n\n"
        "Layer 6 — Convergent evolution: 35 documented convergent evolution events across "
        "kingdoms, including 14 identical amino acid substitutions in echolocation proteins "
        "between bats and dolphins.\n\n"
        "Layer 7 — Viral communication: Entirely new analysis. Viruses as horizontal "
        "information transfer in the cellular economy. 7 new data-backed figures."
    )

    # =========================================================================
    # SLIDE 4: Key Figures — Immune + Genome
    # =========================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)
    add_green_bar(slide)

    add_text_box(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.5),
                 "New Figures: Immune System & Genome",
                 font_size=28, color=GOLD, bold=True)

    # Try to show immune and genome figures
    try_add_picture(slide,
                    FIGURES / "immune_shm_hotspots_annotated.png",
                    FIGURES / "immune_shm_hotspots.png",
                    left=Inches(0.3), top=Inches(1.0),
                    width=Inches(6.2), height=Inches(2.8))

    try_add_picture(slide,
                    FIGURES / "genome_convergent_evolution_annotated.png",
                    FIGURES / "genome_convergent_evolution.png",
                    left=Inches(6.8), top=Inches(1.0),
                    width=Inches(6.2), height=Inches(2.8))

    try_add_picture(slide,
                    FIGURES / "immune_public_clonotypes_annotated.png",
                    FIGURES / "immune_public_clonotypes.png",
                    left=Inches(0.3), top=Inches(4.0),
                    width=Inches(6.2), height=Inches(2.8))

    try_add_picture(slide,
                    FIGURES / "genome_tissue_specialization_annotated.png",
                    FIGURES / "genome_tissue_specialization.png",
                    left=Inches(6.8), top=Inches(4.0),
                    width=Inches(6.2), height=Inches(2.8))

    notes = slide.notes_slide
    notes.notes_text_frame.text = (
        "KEY FIGURES — IMMUNE + GENOME:\n\n"
        "Top left: SHM hotspot analysis — AID targets WRC motifs at 19:1 enrichment. "
        "Every individual mutation position shown as a scatter point.\n\n"
        "Top right: Convergent evolution — 35 events across kingdoms. Bats and dolphins "
        "arriving at identical Prestin substitutions across 95 million years.\n\n"
        "Bottom left: Public clonotypes — identical TCR sequences in unrelated people. "
        "Random probability ~10^-15 per sequence, yet hundreds are shared.\n\n"
        "Bottom right: Tissue specialization — 20% of genes with tau > 0.95, expressed "
        "almost exclusively in one tissue. Division of labor at the molecular level."
    )

    # =========================================================================
    # SLIDE 5: Key Figures — Viral Communication
    # =========================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)
    add_green_bar(slide)

    add_text_box(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.5),
                 "New Figures: Viruses as the Communication Layer",
                 font_size=28, color=GOLD, bold=True)

    try_add_picture(slide,
                    FIGURES / "viral_genome_composition.png",
                    left=Inches(0.3), top=Inches(1.0),
                    width=Inches(6.2), height=Inches(2.8))

    try_add_picture(slide,
                    FIGURES / "viral_syncytin_conservation.png",
                    left=Inches(6.8), top=Inches(1.0),
                    width=Inches(6.2), height=Inches(2.8))

    try_add_picture(slide,
                    FIGURES / "viral_phage_network.png",
                    left=Inches(0.3), top=Inches(4.0),
                    width=Inches(6.2), height=Inches(2.8))

    try_add_picture(slide,
                    FIGURES / "viral_autoimmune_rise.png",
                    left=Inches(6.8), top=Inches(4.0),
                    width=Inches(6.2), height=Inches(2.8))

    notes = slide.notes_slide
    notes.notes_text_frame.text = (
        "VIRAL COMMUNICATION FIGURES:\n\n"
        "Top left: Genome composition — 8% ERV vs 1.5% protein-coding. The genome is "
        "an open ledger, not a sealed vault.\n\n"
        "Top right: Syncytin conservation — captured viral genes under purifying selection "
        "for 85+ Mya. 8 independent capture events across mammalian orders.\n\n"
        "Bottom left: Global phage network — 10^31 particles, 10^25 gene transfers per "
        "day. More daily gene transfers than stars in the observable universe.\n\n"
        "Bottom right: Autoimmune/allergic disease rise — inverse correlation with "
        "infectious disease burden across 18 countries. GABRIELA study showing 45-67% "
        "reduction in farm children with microbial exposure."
    )

    # =========================================================================
    # SLIDE 6: Papers & Deliverables
    # =========================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)
    add_green_bar(slide)

    add_text_box(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.5),
                 "Written Deliverables This Week",
                 font_size=28, color=GOLD, bold=True)

    bullets = [
        ("The Living Architecture — 55-page synthesis paper", [
            "715 lines, 13,625 words, 16 annotated figures, 45 citations",
            "All 7 layers unified: topology → economy → trade → immune → genome → convergence → viral",
            "Includes: design principles, vaccines as application, architecture of creation",
            "Full DOCX with embedded annotated figures (6.6 MB)",
        ]),
        ("Vaccines as Central Planners — standalone essay", [
            "5 distortion mechanisms mapped to Austrian economics",
            "Germinal center monopolization, original antigenic sin, tissue misallocation",
        ]),
        ("Viruses as Communication — standalone essay + 7 figures", [
            "Viruses as horizontal information transfer, not invasion",
            "8% ERV genome, syncytin conservation, phage network scale",
            "Autoimmune/allergic disease rise as communication disruption",
        ]),
        ("Distributed Knowledge Paper — full slide decks", [
            "22-slide comprehensive deck + 3-minute condensed version",
            "16 annotated figures with callout boxes, arrows, key findings",
        ]),
    ]
    add_bullets(slide, bullets, Inches(0.8), Inches(0.9), Inches(11.5), Inches(5.5),
                font_size=14)

    add_text_box(slide, Inches(0.8), Inches(6.5), Inches(11.5), Inches(0.5),
                 "Total figures generated: 39  |  Total analysis scripts: 6  |  Total papers: 3",
                 font_size=13, color=GREEN_MID)

    notes = slide.notes_slide
    notes.notes_text_frame.text = (
        "WRITTEN DELIVERABLES:\n\n"
        "1. The Living Architecture — the synthesis paper. 55 pages, 16 annotated figures, "
        "45 citations. This is the all-encompassing document that ties every layer together "
        "from network topology through viral communication. Includes the vaccines section "
        "as an application of the framework and the theological section as the culmination.\n\n"
        "2. Vaccines as Central Planners — standalone essay mapping five vaccine-induced "
        "distortions to Austrian economic concepts.\n\n"
        "3. Viruses as Communication — standalone essay arguing viruses are the horizontal "
        "information transfer layer of the cellular economy. Backed by 7 new data figures.\n\n"
        "4. Slide decks — 22-slide comprehensive and 3-minute condensed versions of the "
        "distributed knowledge paper, plus the weekly progress slides.\n\n"
        "Total figure count went from 21 last week to 39 this week."
    )

    # =========================================================================
    # SLIDE 7: Figure Annotation System
    # =========================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)
    add_green_bar(slide)

    add_text_box(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.5),
                 "Figure Annotation Pipeline",
                 font_size=28, color=GOLD, bold=True)

    bullets = [
        ("Built annotate_figures.py — automated callout overlay system", [
            "Loads each PNG, overlays callout boxes with arrows + key findings",
            "Three styles: teal (findings), red (warnings), gold (highlights)",
            "All 16 core figures annotated with quantitative callouts",
        ]),
        ("Every figure shows key numbers directly on the plot:", [
            "\"19:1 robustness ratio\" — topology figure",
            "\"Distributed: 71% / Centralized: 53%\" — economy figure",
            "\"Random probability: ~10⁻¹⁵\" — public clonotypes figure",
            "\"Spearman rho ≈ 1.0\" — V(D)J bias figure",
        ]),
        ("Principle: the figure should tell its own story", [
            "Reader should not need to find the paragraph to understand the data",
            "Every data point visible + every key finding annotated on the figure",
        ]),
    ]
    add_bullets(slide, bullets, Inches(0.8), Inches(0.9), Inches(5.5), Inches(5.5),
                font_size=14)

    # Show annotated figure examples
    try_add_picture(slide,
                    FIGURES / "layer1_topology_annotated.png",
                    FIGURES / "layer1_topology.png",
                    left=Inches(7), top=Inches(0.9),
                    width=Inches(5.8), height=Inches(2.8))

    try_add_picture(slide,
                    FIGURES / "layer2_economy_annotated.png",
                    FIGURES / "layer2_economy.png",
                    left=Inches(7), top=Inches(4.0),
                    width=Inches(5.8), height=Inches(2.8))

    notes = slide.notes_slide
    notes.notes_text_frame.text = (
        "ANNOTATION PIPELINE:\n\n"
        "Built an automated figure annotation system. Each of the 16 core figures now "
        "has an annotated version with callout boxes pointing to key data features. "
        "Three visual styles: teal border for findings, red for warnings/critical points, "
        "gold for highlights.\n\n"
        "The design principle: every figure should tell its own story. The reader should "
        "be able to understand the key finding by looking at the figure alone, without "
        "hunting through paragraphs of text.\n\n"
        "This is the same 'every data point visible' philosophy from Week 4, extended: "
        "not only show every data point, but also annotate the key findings directly on "
        "the visualization."
    )

    # =========================================================================
    # SLIDE 8: Next Steps
    # =========================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)
    add_green_bar(slide)

    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.6),
                 "Next Steps", font_size=32, color=GOLD, bold=True)

    nexts = [
        ("Poster layout — Week 6 priority", [
            "7-layer visual hierarchy with key figures",
            "One-line finding per layer + synthesis conclusion",
        ]),
        ("Advisor review of Living Architecture paper", [
            "55-page draft ready for feedback",
            "Key question: which layers are strongest for final presentation?",
        ]),
        ("Annotate viral communication figures", [
            "Apply same callout overlay system to the 7 new viral figures",
        ]),
        ("Refine experimental spread for feasibility", [
            "6 violacein experiments — need advisor input on wet lab constraints",
            "Prioritize experiments 1-3 if resources are limited",
        ]),
        ("Begin final presentation preparation", [
            "Select 10-12 strongest figures from the 39 available",
            "Draft 10-minute presentation narrative",
        ]),
    ]
    add_bullets(slide, nexts, Inches(0.8), Inches(1.3), Inches(5.5), Inches(5.5),
                font_size=15)

    # Summary box on right
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(7.5), Inches(1.5), Inches(5), Inches(4.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0x1b, 0x43, 0x32)
    shape.line.color.rgb = GREEN_MID
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Week 5 Summary"
    p.font.size = Pt(18)
    p.font.color.rgb = GOLD
    p.font.bold = True
    p.font.name = "Calibri"
    p.alignment = PP_ALIGN.CENTER

    p2 = tf.add_paragraph()
    p2.text = (
        "\n3 layers \u2192 7 layers\n"
        "21 figures \u2192 39 figures\n"
        "1 paper \u2192 3 papers + synthesis\n\n"
        "Every layer of biological\n"
        "organization tells the same story:\n\n"
        "The knowledge is distributed\n"
        "because the system was designed\n"
        "to distribute it.\n\n"
        "The first job was gardener,\n"
        "not king."
    )
    p2.font.size = Pt(14)
    p2.font.color.rgb = GREEN_LIGHT
    p2.font.name = "Calibri"
    p2.alignment = PP_ALIGN.CENTER

    notes = slide.notes_slide
    notes.notes_text_frame.text = (
        "NEXT STEPS:\n\n"
        "1. Poster layout — the primary Week 6 deliverable. The 7-layer structure maps "
        "naturally to a visual hierarchy. Each layer gets its strongest figure and a "
        "one-line finding.\n\n"
        "2. Advisor review — the Living Architecture paper is ready for feedback. Key "
        "question: with 7 layers of evidence, which are strongest? Should the final "
        "presentation focus on 3-4 layers deeply or touch all 7?\n\n"
        "3. Annotate viral figures — apply the same callout box system to the 7 new "
        "viral communication figures.\n\n"
        "4. Experimental spread — need advisor input on which of the 6 violacein "
        "experiments are feasible with available resources.\n\n"
        "5. Final presentation — begin selecting the 10-12 strongest figures from "
        "the 39 available and drafting the narrative arc."
    )

    prs.save(str(OUT))
    print(f"Saved: {OUT}")
    print(f"  {len(prs.slides)} slides")


if __name__ == "__main__":
    build()
