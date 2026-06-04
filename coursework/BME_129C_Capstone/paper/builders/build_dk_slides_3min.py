#!/usr/bin/env python3
"""
build_dk_slides_3min.py
3-minute presentation version of Distributed Knowledge paper.
~6-7 slides, ~25-30 seconds each.

Usage:
    python build_dk_slides_3min.py
"""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

SCRIPT_DIR = Path(__file__).resolve().parent
FIGURES_DIR = SCRIPT_DIR / "figures"
OUTPUT = SCRIPT_DIR / "distributed_knowledge_3min.pptx"

BG = RGBColor(0x1A, 0x1A, 0x2E)
TITLE_COLOR = RGBColor(0x00, 0xD4, 0xAA)
SUBTITLE_COLOR = RGBColor(0xAA, 0xBB, 0xCC)
BODY_COLOR = RGBColor(0xE0, 0xE0, 0xE0)
ACCENT = RGBColor(0xFF, 0xAA, 0x00)
DIM_COLOR = RGBColor(0x88, 0x99, 0xAA)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def set_bg(slide):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = BG


def add_text(slide, left, top, width, height, text, size=18,
             color=BODY_COLOR, bold=False, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = "Calibri"
    p.alignment = align
    return tf


def add_bullets(slide, left, top, width, height, items, size=18, color=BODY_COLOR):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"\u2022  {item}"
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.name = "Calibri"
        p.space_after = Pt(8)
    return tf


def two_fig_slide(prs, title, fig1, fig2, bullets):
    """Two figures stacked on left, bullets on right."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_text(slide, Inches(0.3), Inches(0.15), Inches(12.5), Inches(0.6),
             title, size=24, color=TITLE_COLOR, bold=True)
    f1 = FIGURES_DIR / fig1
    f2 = FIGURES_DIR / fig2
    if f1.exists():
        slide.shapes.add_picture(str(f1), Inches(0.3), Inches(0.9), width=Inches(6.2))
    if f2.exists():
        slide.shapes.add_picture(str(f2), Inches(0.3), Inches(3.9), width=Inches(6.2))
    add_bullets(slide, Inches(6.8), Inches(0.9), Inches(6.2), Inches(6.0),
                bullets, size=16)


def one_fig_slide(prs, title, fig, bullets):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_text(slide, Inches(0.3), Inches(0.15), Inches(12.5), Inches(0.6),
             title, size=24, color=TITLE_COLOR, bold=True)
    f = FIGURES_DIR / fig
    if f.exists():
        slide.shapes.add_picture(str(f), Inches(0.3), Inches(0.9), width=Inches(6.8))
    add_bullets(slide, Inches(7.3), Inches(0.9), Inches(5.7), Inches(6.0),
                bullets, size=16)


def build():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # ── Slide 1: Title (15 sec) ──
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_text(slide, Inches(1), Inches(1.2), Inches(11), Inches(1.5),
             "Distributed Knowledge as the\nArchitecture of Life",
             size=40, color=TITLE_COLOR, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, Inches(1), Inches(3.0), Inches(11), Inches(0.8),
             "Why Living Systems Cannot Be Centrally Planned",
             size=24, color=SUBTITLE_COLOR, align=PP_ALIGN.CENTER)
    add_text(slide, Inches(1), Inches(4.3), Inches(11), Inches(0.5),
             "Sage Clokey  |  BME 129C Capstone  |  UC Santa Cruz  |  Spring 2026",
             size=16, color=DIM_COLOR, align=PP_ALIGN.CENTER)
    add_text(slide, Inches(2), Inches(5.5), Inches(9), Inches(0.8),
             "Living systems coordinate billions of nodes without a central planner.\n"
             "Seven layers of evidence show this is not a bug — it is the design.",
             size=18, color=BODY_COLOR, align=PP_ALIGN.CENTER)

    # ── Slide 2: The Question + Topology (30 sec) ──
    one_fig_slide(prs,
        "Is There a Master Node? No. — 19:1 Robustness Advantage",
        "layer1_topology.png",
        [
            "Biological networks: heavy-tailed, no single node dominates",
            "Star graph collapses at 1.9% node removal",
            "PPI networks survive removing 37% of top nodes",
            "19:1 robustness ratio — the structural cost of centralization",
            "Distribute knowledge across many nodes → survive loss of any one",
            "Concentrate knowledge in one node → die when it fails",
        ])

    # ── Slide 3: Economy — Distributed vs Centralized (30 sec) ──
    two_fig_slide(prs,
        "The Planner Wins in Theory — The Market Wins Under Stress",
        "layer2_economy.png",
        "layer2_fba_analysis.png",
        [
            "Centralized achieves 1.68x GDP under stable conditions",
            "Under perturbation: distributed retains 71%, centralized 53%",
            "Even with PERFECT knowledge (FBA): only 70% accuracy",
            "30% failure is structural — regulatory knowledge can't be centralized",
            "The omniscient planner must compute prices to optimize",
            "Prices emerge free from distributed exchange — Hayek was right",
        ])

    # ── Slide 4: Immune System — Not Random (30 sec) ──
    two_fig_slide(prs,
        "The Immune System: Directed Knowledge, Not Random Generation",
        "immune_shm_hotspots.png",
        "immune_vdj_bias.png",
        [
            "SHM targets hotspot motifs at 19:1 over coldspots — AID knows where to write",
            "V(D)J: IGHV3-23 used 10-20x rare segments — not uniform, not random",
            "Same bias in unrelated individuals — Spearman rho ≈ 1.0",
            "Knowledge embedded in chromatin, RSS strength, locus architecture",
            "Every data point shown — no bins, no averages",
            "The mutation machinery reads the landscape before it acts",
        ])

    # ── Slide 5: Convergent Discovery (30 sec) ──
    two_fig_slide(prs,
        "Convergent Discovery: Independent Systems Find the Same Answers",
        "immune_public_clonotypes.png",
        "genome_convergent_evolution.png",
        [
            "Public clonotypes: same TCR sequences in unrelated people",
            "Random probability: 10^-15. Observed: hundreds shared",
            "Convergent evolution: same amino acid substitutions across kingdoms",
            "Bats + dolphins: 14 identical Prestin substitutions (95 Mya apart)",
            "35 events, 17 traits, 8,000 yrs to 1.5 billion yrs",
            "Independent systems converge → solution landscape is structured",
        ])

    # ── Slide 6: Genome-Wide Confirmation (30 sec) ──
    two_fig_slide(prs,
        "Genome-Wide: Same Three Patterns at Every Scale",
        "genome_mutation_hotspots.png",
        "genome_tissue_specialization.png",
        [
            "CpG C>T mutations: 15-40x baseline from ~1% of contexts",
            "Ti/Tv ratio 2:1 — 4x random expectation — conservative bias",
            "20% of genes: tau > 0.95 — expressed almost exclusively in one tissue",
            "Fold enrichment 100-1,000x — division of labor at the molecular level",
            "Directed mutation + specialized expression + convergent discovery",
            "Immune scale and genome scale: same architecture, same conclusion",
        ])

    # ── Slide 7: Conclusion (15 sec) ──
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_text(slide, Inches(1), Inches(0.8), Inches(11), Inches(1),
             "Seven Layers, One Conclusion",
             size=36, color=TITLE_COLOR, bold=True, align=PP_ALIGN.CENTER)

    # Summary table
    rows = [
        ("Topology", "No master node — 19:1 robustness"),
        ("Single-Cell", "No master cell — Gini = 0.0"),
        ("Economy", "Distributed 71% vs centralized 53%"),
        ("FBA", "Omniscient planner: 70% — 30% structural failure"),
        ("Trade", "Forced exchange destroys information"),
        ("Immune", "Directed mutation, biased recombination, convergent clonotypes"),
        ("Genome", "CpG hotspots, tissue specialization, convergent evolution"),
    ]
    for i, (layer, finding) in enumerate(rows):
        y = Inches(2.0) + Inches(0.6) * i
        add_text(slide, Inches(1.5), y, Inches(2.8), Inches(0.5),
                 layer, size=17, color=ACCENT, bold=True)
        add_text(slide, Inches(4.5), y, Inches(7.5), Inches(0.5),
                 finding, size=16, color=BODY_COLOR)

    add_text(slide, Inches(1.5), Inches(6.3), Inches(10), Inches(0.8),
             "The knowledge is distributed because the system was designed to distribute it.\n"
             "The first job was gardener, not king.",
             size=18, color=SUBTITLE_COLOR, align=PP_ALIGN.CENTER)

    prs.save(str(OUTPUT))
    print(f"Saved: {OUTPUT}  ({len(prs.slides)} slides)")


if __name__ == "__main__":
    build()
