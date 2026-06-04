"""
Generate 6-Minute Presentation Slides
======================================
Builds a PowerPoint deck aligned to the capstone presentation speech.
Uses existing figures from paper/figures/.

Usage:
    python paper/builders/generate_speech_slides.py

BME 129C Capstone — Sage Clokey — Spring 2026
"""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

PROJECT_ROOT = Path(__file__).resolve().parent.parent.parent
FIG_DIR = PROJECT_ROOT / "paper" / "figures"
INDIVIDUAL_DIR = FIG_DIR / "individual"

# Colors matching existing presentation style
BACKGROUND = RGBColor(0x0d, 0x11, 0x17)
SPIRAL_GREEN = RGBColor(0x2d, 0x6a, 0x4f)
SPIRAL_MID = RGBColor(0x52, 0xb7, 0x88)
GOLD = RGBColor(0xe9, 0xc4, 0x6a)
RED = RGBColor(0xe6, 0x39, 0x46)
TEXT_MAIN = RGBColor(0xe6, 0xed, 0xf3)
TEXT_DIM = RGBColor(0x8b, 0x94, 0x9e)
WHITE = RGBColor(0xff, 0xff, 0xff)


def set_slide_bg(slide, color=BACKGROUND):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_title_slide(prs, title, subtitle, notes=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_slide_bg(slide)

    # Title
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(11.5), Inches(1.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.color.rgb = SPIRAL_MID
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    # Subtitle
    txBox2 = slide.shapes.add_textbox(Inches(1.5), Inches(3.8), Inches(10.0), Inches(1.5))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.text = subtitle
    p2.font.size = Pt(20)
    p2.font.color.rgb = GOLD
    p2.alignment = PP_ALIGN.CENTER

    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = notes
    return slide


def add_statement_slide(prs, statement, sub="", notes=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    txBox = slide.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(11.0), Inches(2.0))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = statement
    p.font.size = Pt(32)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    if sub:
        txBox2 = slide.shapes.add_textbox(Inches(1.5), Inches(4.2), Inches(10.0), Inches(1.5))
        tf2 = txBox2.text_frame
        tf2.word_wrap = True
        p2 = tf2.paragraphs[0]
        p2.text = sub
        p2.font.size = Pt(18)
        p2.font.color.rgb = TEXT_DIM
        p2.alignment = PP_ALIGN.CENTER

    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = notes
    return slide


def add_figure_slide(prs, title, fig_name, bullets, notes="", fig_dir=None):
    """Figure on left ~60%, bullets on right ~40%."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    # Title bar
    txBox = slide.shapes.add_textbox(Inches(0.3), Inches(0.15), Inches(12.7), Inches(0.6))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(22)
    p.font.color.rgb = SPIRAL_MID
    p.font.bold = True

    # Figure
    search_dirs = [fig_dir] if fig_dir else [FIG_DIR, INDIVIDUAL_DIR]
    img_path = None
    for d in search_dirs:
        if d is None:
            continue
        candidate = d / f"{fig_name}.png"
        if candidate.exists():
            img_path = candidate
            break

    if img_path and img_path.exists():
        slide.shapes.add_picture(
            str(img_path), Inches(0.3), Inches(0.9), Inches(7.8), Inches(5.9)
        )
    else:
        txBox_ph = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(7.5), Inches(2))
        tf_ph = txBox_ph.text_frame
        tf_ph.word_wrap = True
        p_ph = tf_ph.paragraphs[0]
        p_ph.text = f"[Figure: {fig_name}.png]"
        p_ph.font.size = Pt(14)
        p_ph.font.color.rgb = TEXT_DIM
        p_ph.alignment = PP_ALIGN.CENTER

    # Bullets
    txBox_b = slide.shapes.add_textbox(Inches(8.3), Inches(1.0), Inches(4.7), Inches(5.8))
    tf_b = txBox_b.text_frame
    tf_b.word_wrap = True

    for i, bullet in enumerate(bullets):
        if i == 0:
            p_b = tf_b.paragraphs[0]
        else:
            p_b = tf_b.add_paragraph()
        p_b.text = f"\u2022  {bullet}"
        p_b.font.size = Pt(13)
        p_b.font.color.rgb = TEXT_MAIN
        p_b.space_after = Pt(8)

    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = notes
    return slide


def add_table_slide(prs, title, headers, rows, notes=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    # Title
    txBox = slide.shapes.add_textbox(Inches(0.3), Inches(0.15), Inches(12.7), Inches(0.6))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(22)
    p.font.color.rgb = SPIRAL_MID
    p.font.bold = True

    # Table
    n_rows = len(rows) + 1
    n_cols = len(headers)
    table_shape = slide.shapes.add_table(
        n_rows, n_cols, Inches(0.4), Inches(0.9), Inches(12.3), Inches(5.8)
    )
    table = table_shape.table

    # Header row
    for j, h in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = h
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.size = Pt(11)
            paragraph.font.bold = True
            paragraph.font.color.rgb = GOLD
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(0x16, 0x1b, 0x22)

    # Data rows
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            cell = table.cell(i + 1, j)
            cell.text = val
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(10)
                paragraph.font.color.rgb = TEXT_MAIN
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(0x0d, 0x11, 0x17)

    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = notes
    return slide


def add_bullets_slide(prs, title, bullets, notes=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    # Title
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.0), Inches(0.8))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(26)
    p.font.color.rgb = SPIRAL_MID
    p.font.bold = True

    # Bullets
    txBox_b = slide.shapes.add_textbox(Inches(0.8), Inches(1.4), Inches(11.5), Inches(5.5))
    tf_b = txBox_b.text_frame
    tf_b.word_wrap = True

    for i, bullet in enumerate(bullets):
        if i == 0:
            p_b = tf_b.paragraphs[0]
        else:
            p_b = tf_b.add_paragraph()
        p_b.text = f"\u2022  {bullet}"
        p_b.font.size = Pt(18)
        p_b.font.color.rgb = TEXT_MAIN
        p_b.space_after = Pt(12)

    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = notes
    return slide


def build():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # ---------------------------------------------------------------
    # SLIDE 1: Title
    # ---------------------------------------------------------------
    add_title_slide(
        prs,
        "The Living Architecture:\nDistributed Knowledge as the Design Language of Life",
        "Sage Clokey  \u2022  BME 129C Capstone  \u2022  Spring 2026\nAdvisor: R. Dubois  \u2022  UC Santa Cruz",
        notes="Living things are not machines. That is the thesis. That is the finding. "
              "That is what seven layers of quantitative evidence confirm."
    )

    # ---------------------------------------------------------------
    # SLIDE 2: Thesis
    # ---------------------------------------------------------------
    add_statement_slide(
        prs,
        "Living things are not machines.",
        "Does life operate like a centrally planned economy \u2014 or a free market?",
        notes="My name is Sage Clokey, and my capstone asks a question that sounds simple "
              "but has enormous consequences for how we do molecular biology."
    )

    # ---------------------------------------------------------------
    # SLIDE 3: The Problem
    # ---------------------------------------------------------------
    add_bullets_slide(
        prs,
        "The Central Planner\u2019s Promise",
        [
            "Synthetic biology assumes the engineer can design from above",
            "Pick the promoter. Optimize codons. Calculate ratios. Assemble parts.",
            "Most designed genetic circuits fail on the first attempt (Kwok, 2010)",
            "Complexity hits a wall beyond a handful of components (Purnick & Weiss, 2009)",
            "Standard explanation: biology is \u201ccomplex\u201d and models are \u201cincomplete\u201d",
            "This paper argues: the problem is structural, not informational",
        ],
        notes="Modern synthetic biology operates on an implicit assumption: that the engineer "
              "can gather enough knowledge to design a living system from above. "
              "This is central planning applied to molecular biology."
    )

    # ---------------------------------------------------------------
    # SLIDE 4: Austrian Economics Framework
    # ---------------------------------------------------------------
    add_bullets_slide(
        prs,
        "The Knowledge Problem \u2014 From Vienna to the Genome",
        [
            "Hayek (1945): knowledge for coordination never exists in concentrated form",
            "Mises (1920): rational calculation impossible without prices",
            "Menger (1871): complex institutions arise spontaneously from individual action",
            "Rothbard (1962): every act of central planning destroys information",
            "Kirzner (1973): agents discover information through competitive process",
            "Every principle they discovered about human cooperation, life discovered first \u2014 and wrote it in DNA",
        ],
        notes="The insight that coordination can depend on knowledge that cannot be centralized "
              "was first articulated not by a biologist but by an economist."
    )

    # ---------------------------------------------------------------
    # SLIDE 5: Layer 1 — Network Topology
    # ---------------------------------------------------------------
    add_figure_slide(
        prs,
        "Layer 1: No Master Node \u2014 Network Topology",
        "layer1_topology_annotated",
        [
            "5 biological networks vs 5 synthetic architectures",
            "Heavy-tailed degree distributions (\u03b1 = 2.0\u20132.5)",
            "Hubs exist but no hub dominates",
            "PPI: 36.8% removal to fragment",
            "Star graph: 1.9% removal to collapse",
            "19:1 robustness ratio",
            "Feed-forward loops over-represented (Z > 10)",
            "Networks self-regulate against centralization (WBPA)",
        ],
        notes="I built five biological networks from public databases and compared them against "
              "five synthetic architectures including a star graph. The result: there is no master node. "
              "19-to-1 robustness ratio favoring distributed architecture."
    )

    # ---------------------------------------------------------------
    # SLIDE 6: Layer 1b — Single-Cell Economy
    # ---------------------------------------------------------------
    add_figure_slide(
        prs,
        "Layer 1b: No Master Cell \u2014 The Single-Cell Economy",
        "layer1b_single_cell_economy_annotated",
        [
            "2,638 human PBMCs, 8 cell types",
            "Shannon entropy 0.852\u20130.915",
            "Division of labor without hierarchy",
            "Communication betweenness Gini = 0.000",
            "Perfectly distributed \u2014 no gatekeeper",
            "75% of communication survives any single removal",
            "Menger\u2019s spontaneous order made visible",
        ],
        notes="I analyzed 2,638 human immune cells across 8 cell types. Each cell type specializes. "
              "The communication network has a betweenness Gini of zero. "
              "Remove any single cell type and 75% of communication survives."
    )

    # ---------------------------------------------------------------
    # SLIDE 7: Layer 2 — Metabolic Simulation
    # ---------------------------------------------------------------
    add_figure_slide(
        prs,
        "Layer 2: Distributed Outperforms Centralized Under Perturbation",
        "layer2_economy_annotated",
        [
            "13 metabolic pathway agents",
            "Stable conditions: planner wins (1.68x GDP)",
            "Under perturbation:",
            "  Distributed retains 71.1% GDP",
            "  Centralized retains 53.0% GDP",
            "18.1 percentage point advantage",
            "Planner\u2019s fixed allocation has no self-correction",
            "Agents discover equilibrium via price feedback",
        ],
        notes="Under stable conditions, the central planner achieves higher output. "
              "But when I removed an agent, the distributed system retained 71% of GDP "
              "while the centralized system retained only 53%."
    )

    # ---------------------------------------------------------------
    # SLIDE 8: Layer 2b — FBA
    # ---------------------------------------------------------------
    add_figure_slide(
        prs,
        "Layer 2b: The Omniscient Planner Achieves 70%",
        "layer2_fba_analysis_annotated",
        [
            "iML1515: 2,712 reactions, 1,877 metabolites, 1,516 genes",
            "The strongest possible central planner",
            "70% accuracy on gene knockouts vs Keio collection",
            "30% failure is structural, not noise",
            "Allosteric feedback, folding, timing, chaperones",
            "Knowledge exists only in local molecular state",
            "The LP must compute shadow prices (= Hayekian prices)",
            "The planner computes what the market produces for free",
        ],
        notes="I gave the planner perfect information. iML1515 genome-scale model. "
              "It achieves 70% accuracy. The 30% failure is structural. "
              "And the irony: the LP must compute shadow prices to solve its optimization."
    )

    # ---------------------------------------------------------------
    # SLIDE 9: Layer 3 — Cross-Species Trade
    # ---------------------------------------------------------------
    add_figure_slide(
        prs,
        "Layer 3: Cross-Species Gene Transfer Follows Trade Rules",
        "layer3_trade_network_annotated",
        [
            "8 organisms across 4 kingdoms",
            "Trade cost scales with evolutionary distance",
            "Within-kingdom: 0.17\u20130.38",
            "Cross-kingdom: 0.65\u20130.83",
            "Forced codon optimization destroys information",
            "Codon harmonization preserves local knowledge",
            "Trade blocs emerge spontaneously (Louvain)",
            "Ricardian comparative advantage in codons",
        ],
        notes="Gene transfer between species follows trade network rules. "
              "Costs scale with evolutionary distance. "
              "Forced codon optimization destroys information. "
              "Trade blocs emerge spontaneously from shared evolutionary history."
    )

    # ---------------------------------------------------------------
    # SLIDE 10: Layer 4 — Immune System
    # ---------------------------------------------------------------
    add_figure_slide(
        prs,
        "Layer 4: The Immune System \u2014 Directed, Not Random",
        "immune_distributed_summary_annotated",
        [
            "SHM: AID targets WRC/GYW hotspots at 19:1 over coldspots",
            "V(D)J: IGHV3-23 used at 10\u201320x rare segments",
            "Same bias in unrelated individuals (Spearman \u03c1 \u2248 1.0)",
            "IGHJ4 + IGHJ6 = 65% of usage (random: 17%)",
            "Public clonotypes: identical TCR CDR3 in unrelated people",
            "Random probability: ~10\u207b\xb9\u2075",
            "Observed: 200 shared across 10 individuals",
            "Convergent distributed discovery",
        ],
        notes="Textbooks call V(D)J recombination and somatic hypermutation random. They are not. "
              "AID targets hotspot motifs at 19-to-1 over coldspots. "
              "Public clonotypes appear at rates exceeding random expectation by 10-to-the-15th-fold."
    )

    # ---------------------------------------------------------------
    # SLIDE 11: Layer 5 — Whole Genome
    # ---------------------------------------------------------------
    add_figure_slide(
        prs,
        "Layer 5: The Whole Genome \u2014 Scale-Invariant Distributed Knowledge",
        "genome_distributed_summary_annotated",
        [
            "CpG C>T: 15\u201340x baseline from ~1% of contexts",
            "Ti/Tv ratio 2:1 (random: 0.5:1)",
            "Mutation machinery is a biased editor, not a random generator",
            "20% of genes: tissue specificity \u03c4 > 0.95",
            "100\u20131,000x fold enrichment in primary tissue",
            "35 convergent evolution events across 17 traits",
            "Same amino acid substitutions across 1.5 billion years",
            "The solution landscape is structured",
        ],
        notes="CpG mutations occur at 15 to 40x baseline from 1% of sequence contexts. "
              "20% of genes show tissue specificity above tau 0.95. "
              "35 convergent evolution events reproduce the same amino acid substitutions "
              "in lineages separated by up to 1.5 billion years."
    )

    # ---------------------------------------------------------------
    # SLIDE 12: Evidence Summary Table
    # ---------------------------------------------------------------
    add_table_slide(
        prs,
        "Seven Layers \u2014 Same Answer",
        ["Layer", "Question", "Finding", "Principle"],
        [
            ["1: Topology", "Is there a master node?",
             "No. 19:1 robustness advantage", "Hayek: knowledge is dispersed"],
            ["1b: Single-cell", "Is there a master cell?",
             "No. Gini = 0.0, 75% survival", "Menger: spontaneous order"],
            ["2: Economy", "Distributed > centralized?",
             "Yes. 71% vs 53% GDP", "Mises: calculation problem"],
            ["2b: FBA", "Perfect knowledge solves it?",
             "No. 70% accuracy, 30% structural", "Hayek: structural limit"],
            ["3: Trade", "Forced exchange works?",
             "No. Destroys information", "Rothbard: coercion destroys value"],
            ["4: Immune", "Immune generation random?",
             "No. 19:1 hotspots, \u03c1 \u2248 1.0, public clonotypes", "Distributed knowledge"],
            ["5: Genome", "Mutation random?",
             "No. CpG hotspots, convergent evolution", "Knowledge in the machinery"],
        ],
        notes="Seven layers of evidence answer the same question from different angles. "
              "The convergence across seven independent lines of evidence is the finding."
    )

    # ---------------------------------------------------------------
    # SLIDE 13: Design Principles
    # ---------------------------------------------------------------
    add_bullets_slide(
        prs,
        "Six Design Principles for the Molecular Biologist",
        [
            "1. Read the economy before entering it \u2014 use FBA shadow prices",
            "2. Build feedback loops, not fixed rates \u2014 71% vs 53% GDP retention",
            "3. Distribute control across the pathway \u2014 19:1 robustness ratio",
            "4. Harmonize codons, don\u2019t optimize them \u2014 preserve local knowledge",
            "5. Let the system evolve \u2014 Kirznerian discovery at the molecular level",
            "6. Design consortia, not monoliths \u2014 Gini = 0.0, 75% fault tolerance",
        ],
        notes="What does this mean for the practicing molecular biologist? Six design principles. "
              "Read the economy. Build feedback loops. Distribute control. Harmonize codons. "
              "Let it evolve. Design consortia."
    )

    # ---------------------------------------------------------------
    # SLIDE 14: Closing
    # ---------------------------------------------------------------
    add_statement_slide(
        prs,
        "The garden is still growing.",
        "The question is not how to program life.\n"
        "It is how to join it.\n\n"
        "Cultivate conditions, don\u2019t command outcomes.\n"
        "Read the price system, don\u2019t override it.\n"
        "Be a node, not a planner.",
        notes="The cell is not a chassis waiting to be programmed. "
              "It is a running economy. The garden is still growing. "
              "The question is not how to program life. It is how to join it. Thank you."
    )

    # ---------------------------------------------------------------
    # SLIDE 15: Thank You
    # ---------------------------------------------------------------
    add_title_slide(
        prs,
        "Thank You",
        "Sage Clokey  \u2022  sage-clokey.github.io\n"
        "Code: github.com/Sage-Clokey/Living-works-by-the-word\n"
        "Advisor: R. Dubois  \u2022  UC Santa Cruz",
        notes="Thank you. Questions?"
    )

    # Save
    out_path = PROJECT_ROOT / "paper" / "presentations" / "capstone_6min_speech.pptx"
    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out_path))
    print(f"Presentation saved to: {out_path}")
    return out_path


if __name__ == "__main__":
    build()
