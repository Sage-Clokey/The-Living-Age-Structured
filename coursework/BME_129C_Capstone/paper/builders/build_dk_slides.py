#!/usr/bin/env python3
"""
build_dk_slides.py
Generates a PowerPoint presentation for the Distributed Knowledge paper figures.

Usage:
    python build_dk_slides.py
"""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

SCRIPT_DIR = Path(__file__).resolve().parent
FIGURES_DIR = SCRIPT_DIR / "figures"
OUTPUT = SCRIPT_DIR / "distributed_knowledge_slides.pptx"

# Dark theme colors
BG = RGBColor(0x1A, 0x1A, 0x2E)
TITLE_COLOR = RGBColor(0x00, 0xD4, 0xAA)
SUBTITLE_COLOR = RGBColor(0xAA, 0xBB, 0xCC)
BODY_COLOR = RGBColor(0xE0, 0xE0, 0xE0)
ACCENT = RGBColor(0xFF, 0xAA, 0x00)
DIM_COLOR = RGBColor(0x88, 0x99, 0xAA)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def set_bg(slide, color=BG):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text_box(slide, left, top, width, height, text, font_size=18,
                 color=BODY_COLOR, bold=False, alignment=PP_ALIGN.LEFT,
                 font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return tf


def add_bullet_list(slide, left, top, width, height, items, font_size=16,
                    color=BODY_COLOR):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Calibri"
        p.space_after = Pt(6)
        p.level = 0
    return tf


def title_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_bg(slide)
    add_text_box(slide, Inches(1), Inches(1.5), Inches(11), Inches(1.5),
                 "Distributed Knowledge as the\nArchitecture of Life",
                 font_size=36, color=TITLE_COLOR, bold=True,
                 alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(1), Inches(3.2), Inches(11), Inches(0.8),
                 "Why Living Systems Cannot Be Centrally Planned",
                 font_size=24, color=SUBTITLE_COLOR, bold=False,
                 alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(1), Inches(4.5), Inches(11), Inches(0.5),
                 "Sage Clokey  |  BME 129C Capstone  |  UC Santa Cruz  |  Spring 2026",
                 font_size=16, color=DIM_COLOR, bold=False,
                 alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(1.5), Inches(5.5), Inches(10), Inches(1),
                 '"The curious task of economics is to demonstrate to men how little they\n'
                 'really know about what they imagine they can design." — F. A. Hayek',
                 font_size=14, color=DIM_COLOR, bold=False,
                 alignment=PP_ALIGN.CENTER)


def thesis_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_text_box(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.7),
                 "Core Thesis", font_size=28, color=TITLE_COLOR, bold=True)
    items = [
        "Living systems coordinate across billions of nodes without a central planner",
        "The knowledge required for coordination is distributed — it cannot be centralized",
        "This is not a limitation of technology. It is a structural feature of life",
        "Seven layers of evidence confirm: distributed architecture outperforms centralized planning",
        "The same pattern repeats at every scale: trinucleotides, genes, cells, pathways, organisms, species",
        "Austrian economics (Hayek, Mises, Menger, Rothbard, Kirzner) predicted every finding",
    ]
    add_bullet_list(slide, Inches(0.8), Inches(1.2), Inches(11.5), Inches(5.5),
                    [f"\u2022  {item}" for item in items], font_size=20)


def evidence_overview_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_text_box(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.7),
                 "Seven Layers of Evidence", font_size=28, color=TITLE_COLOR, bold=True)
    layers = [
        ("Layer 1: Topology", "No master node. 19:1 robustness advantage for distributed networks"),
        ("Layer 1b: Single-Cell", "No master cell. Gini = 0.0. 75% survives any single removal"),
        ("Layer 2: Economy", "Distributed retains 71% GDP under perturbation; centralized retains 53%"),
        ("Layer 2b: FBA", "Omniscient planner achieves only 70% accuracy — 30% structural failure"),
        ("Layer 3: Trade", "Forced gene transfer destroys information; voluntary exchange preserves it"),
        ("Layer 4: Immune", "SHM hotspot targeting, V(D)J bias, public clonotypes — not random"),
        ("Layer 5: Genome", "CpG hotspots, tissue specialization, convergent evolution — knowledge is embedded"),
    ]
    for i, (layer, finding) in enumerate(layers):
        y = Inches(1.2) + Inches(0.8) * i
        add_text_box(slide, Inches(0.8), y, Inches(3.5), Inches(0.7),
                     layer, font_size=18, color=ACCENT, bold=True)
        add_text_box(slide, Inches(4.5), y, Inches(8), Inches(0.7),
                     finding, font_size=16, color=BODY_COLOR)


def figure_slide(prs, title, fig_filename, bullets):
    """Create a slide with figure on left, bullet points on right."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    # Title
    add_text_box(slide, Inches(0.3), Inches(0.15), Inches(12.5), Inches(0.6),
                 title, font_size=22, color=TITLE_COLOR, bold=True)
    # Figure
    fig_path = FIGURES_DIR / fig_filename
    if fig_path.exists():
        slide.shapes.add_picture(str(fig_path), Inches(0.3), Inches(0.85),
                                 width=Inches(7.0))
    else:
        add_text_box(slide, Inches(0.3), Inches(2.5), Inches(7), Inches(1),
                     f"[Image not found: {fig_filename}]", font_size=14,
                     color=RGBColor(0xFF, 0x44, 0x44))
    # Bullets
    add_bullet_list(slide, Inches(7.5), Inches(0.85), Inches(5.5), Inches(6.2),
                    [f"\u2022  {b}" for b in bullets], font_size=15)


def synthesis_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_text_box(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.7),
                 "Synthesis: Why Distributed?", font_size=28, color=TITLE_COLOR, bold=True)
    items = [
        "At every scale — the same pattern: no master node, no master cell, no master pathway",
        "Knowledge is distributed because the system was designed to distribute it",
        "Mutations are directed information, not random error filtered by death",
        "Selection is the editor, not the engine — the engine is directed creative growth",
        "The Creator chose the garden over the kingdom, the gradient over the decree",
        "The sagent reads the economy before entering it; the planner overrides it",
        "Distributed systems sacrifice peak efficiency for survival under perturbation",
        "Four billion years of evidence: the gardener outperforms the central planner",
    ]
    add_bullet_list(slide, Inches(0.8), Inches(1.2), Inches(11.5), Inches(5.5),
                    [f"\u2022  {item}" for item in items], font_size=18)


def design_principles_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_text_box(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.7),
                 "Design Principles: Engineering as Gardening",
                 font_size=28, color=TITLE_COLOR, bold=True)
    principles = [
        ("Read the economy before entering it",
         "Run FBA. Read shadow prices. Know the competitive landscape before inserting your pathway."),
        ("Build feedback loops, not fixed rates",
         "Metabolite-responsive biosensors > constitutive promoters. 71% vs 53% GDP retention."),
        ("Distribute control across the pathway",
         "Each enzyme gets its own sensor. 19x longer survival than single-master-switch design."),
        ("Harmonize codons, don't optimize them",
         "Preserve rare codon information. Reduce barriers without destroying local knowledge."),
        ("Let it evolve",
         "Adaptive lab evolution. Kirznerian discovery at the molecular level."),
        ("Design consortia, not monoliths",
         "Split pathways across strains. Division of labor. Voluntary metabolite exchange."),
    ]
    for i, (title, desc) in enumerate(principles):
        y = Inches(1.15) + Inches(0.95) * i
        add_text_box(slide, Inches(0.8), y, Inches(11.5), Inches(0.4),
                     title, font_size=18, color=ACCENT, bold=True)
        add_text_box(slide, Inches(1.0), y + Inches(0.35), Inches(11.3), Inches(0.5),
                     desc, font_size=14, color=DIM_COLOR)


def build():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # --- Slide 1: Title ---
    title_slide(prs)

    # --- Slide 2: Core Thesis ---
    thesis_slide(prs)

    # --- Slide 3: Seven Layers Overview ---
    evidence_overview_slide(prs)

    # --- Slide 4: Figure 1 — Network Topology ---
    figure_slide(prs,
        "Layer 1: No Master Node — Network Topology",
        "layer1_topology.png",
        [
            "Biological networks follow heavy-tailed degree distributions (alpha 2.0-2.5)",
            "Hubs exist (CRP regulates 43 genes) but no hub dominates",
            "Betweenness Gini: star graph 0.998 vs biological 0.72-0.94",
            "19:1 robustness ratio: distributed networks survive removing 37% of top nodes",
            "Star graph collapses at 1.9% — one node removal destroys everything",
            "Hayek's knowledge problem in graph theory: distribute knowledge, survive loss",
        ])

    # --- Slide 5: Figure 2 — Single-Cell Economy ---
    figure_slide(prs,
        "Layer 1b: No Master Cell — Single-Cell Economy",
        "layer1b_single_cell_economy.png",
        [
            "2,638 PBMCs, 8 cell types from one genome — no central coordinator",
            "Division of labor: each cell type specializes via local signals",
            "Communication betweenness Gini = 0.0 — perfectly distributed",
            "No gatekeeper: every cell type signals directly to every other",
            "75% of edges survive any single cell type removal",
            "Menger's spontaneous order: complex organization from local decisions",
        ])

    # --- Slide 6: Figure 3 — Price Signals ---
    figure_slide(prs,
        "Layer 1b: Subjective Value — Same Signal, Different Meaning",
        "layer1b_price_signals.png",
        [
            "TNF-alpha activates inflammation in monocytes but triggers apoptosis in T cells",
            "Same molecule, different meaning depending on receiver's context",
            "Value is not intrinsic to the signal — it is determined by the receiver",
            "Menger's subjective value theory at the molecular level",
            "A central planner cannot assign fixed meaning to a cytokine",
            "The 'price' depends on local state no external observer can access",
        ])

    # --- Slide 7: Figure 4 — Distributed vs Centralized Economy ---
    figure_slide(prs,
        "Layer 2: Distributed vs Centralized Resource Allocation",
        "layer2_economy.png",
        [
            "13 metabolic pathways simulated as economic agents",
            "Centralized achieves 1.68x higher GDP under stable conditions",
            "Under perturbation (HIF1-alpha removal): distributed retains 71%, centralized 53%",
            "18.1 percentage point advantage for distributed coordination under stress",
            "Agents discover their own production rates through local feedback",
            "Mises' calculation problem quantified: the plan breaks when conditions change",
        ])

    # --- Slide 8: Figure 5 — FBA Analysis ---
    figure_slide(prs,
        "Layer 2b: The Omniscient Planner Still Fails — FBA",
        "layer2_fba_analysis.png",
        [
            "iML1515 model: 2,712 reactions, 1,877 metabolites, 1,516 genes",
            "This IS the omniscient planner — complete stoichiometric knowledge",
            "Achieves 70% accuracy on gene knockouts vs Keio collection",
            "30% failure is structural: regulatory knowledge cannot be encoded in LP",
            "Shadow prices prove Hayek's point: the planner must compute prices to plan",
            "Prices emerge naturally from distributed exchange — the planner reinvents the wheel",
        ])

    # --- Slide 9: Figure 6 — FBA Perturbation ---
    figure_slide(prs,
        "Layer 2b: The Planner Re-Solves Instantly; Biology Must Discover",
        "layer2_fba_perturbation.png",
        [
            "Glucose-to-acetate: FBA re-solves instantly; E. coli shows diauxic lag",
            "Aerobic-to-anaerobic: FBA re-solves; E. coli activates ArcAB/FNR cascades",
            "Nitrogen limitation: FBA predicts optimum; E. coli upregulates via NtrBC",
            "Biology pays the real cost of distributed discovery",
            "But biology can discover answers to questions the planner never anticipated",
            "Kirznerian entrepreneurial alertness at the molecular level",
        ])

    # --- Slide 10: Figure 7 — Trade Network ---
    figure_slide(prs,
        "Layer 3: The Tree of Life Is a Trade Network",
        "layer3_trade_network.png",
        [
            "8 organisms across 4 kingdoms analyzed for gene transferability",
            "Trade costs mirror evolutionary distance (biological gravity model)",
            "Lowest cost: human-to-axolotl (0.169) — shared vertebrate machinery",
            "Highest cost: prokaryote-to-eukaryote (0.65-0.83) — regulatory divergence",
            "Trade blocs emerge spontaneously from shared evolutionary history",
            "Menger's spontaneous order at the inter-species level",
        ])

    # --- Slide 11: Figure 8 — Voluntary Exchange ---
    figure_slide(prs,
        "Layer 3: Voluntary Exchange Succeeds; Forced Exchange Fails",
        "layer3_voluntary_exchange.png",
        [
            "Full codon optimization = forced trade: replaces every codon with host's preference",
            "Destroys information encoded in rare codons: folding signals, pausing sites, mRNA structure",
            "Codon harmonization = voluntary exchange: preserves usage pattern, reduces barriers",
            "Success scales with compatibility, not with how aggressively you rewrite",
            "The planner who forces trade destroys local knowledge",
            "Rothbard's prediction confirmed: coercion destroys value",
        ])

    # --- Slide 12: Figure 9 — SHM Hotspots ---
    figure_slide(prs,
        "Layer 4: Somatic Hypermutation Is Directed, Not Random",
        "immune_shm_hotspots.png",
        [
            "AID targets WRC/GYW hotspot motifs at 5x coldspot rate",
            "200 sequences: 285 hotspot mutations vs 15 coldspot — 19:1 enrichment",
            "Every mutation shown as individual data point — no bins, no averages",
            "If random: uniform scatter. Observed: visible clustering at hotspots",
            "AID carries local knowledge of where beneficial mutations are most likely",
            "Not random dice rolls filtered by selection — directed information",
        ])

    # --- Slide 13: Figure 10 — V(D)J Bias ---
    figure_slide(prs,
        "Layer 4: V(D)J Recombination Is Biased — Knowledge, Not Lottery",
        "immune_vdj_bias.png",
        [
            "Random prediction: all 50 V segments used at ~2% each",
            "Observed: IGHV3-23 and IGHV4-34 used at 10-20x rare segments",
            "Same bias across unrelated individuals — Spearman rho approaching 1.0",
            "J segments: IGHJ4 + IGHJ6 = 65% of usage (random: 17% each)",
            "Bias encoded in chromatin accessibility, RSS strength, locus architecture",
            "The machinery carries knowledge about which segments are most useful",
        ])

    # --- Slide 14: Figure 11 — Public Clonotypes ---
    figure_slide(prs,
        "Layer 4: Public Clonotypes — Convergent Distributed Discovery",
        "immune_public_clonotypes.png",
        [
            "TCR diversity: ~10^15 possible CDR3 sequences",
            "Random probability of same sequence in 2 people: ~10^-15",
            "Observed: 200 public clonotypes shared across 10 individuals",
            "Some sequences shared by all 10/10 — random expectation: zero",
            "Median ~40 shared clonotypes per pair of unrelated individuals",
            "Independent immune systems converge on the same molecular solutions",
            "Kirznerian entrepreneurial discovery: same opportunity found independently",
        ])

    # --- Slide 15: Figure 12 — Immune Summary ---
    figure_slide(prs,
        "Layer 4: Immune System Summary — Three Rejections of Randomness",
        "immune_distributed_summary.png",
        [
            "Hotspot/coldspot ratio: observed 19:1, random prediction 1:1",
            "V segment Gini: observed >> 0, random prediction 0 (uniform)",
            "Public clonotype fraction: observed substantial, random prediction 0%",
            "Every individual data point shown — no bins, no smoothing",
            "Three independent measurements, three rejections of random model",
            "The immune system carries distributed knowledge at every level",
        ])

    # --- Slide 16: Figure 13 — Genome Mutation Hotspots ---
    figure_slide(prs,
        "Layer 5: Genome-Wide Mutations Are Context-Dependent, Not Random",
        "genome_mutation_hotspots.png",
        [
            "96 trinucleotide mutation contexts — 40-fold variation in rate",
            "CpG C>T mutations: 15-40x baseline rate (methylation-driven)",
            "CpG transitions = 48.7% of all mutations from ~1% of genome contexts",
            "Ti/Tv ratio 2:1 (random expectation: 0.5:1) — 4x enrichment",
            "Mutation machinery preferentially produces conservative substitutions",
            "Information encoded in DNA chemistry itself — not random noise",
        ])

    # --- Slide 17: Figure 14 — Tissue Specialization ---
    figure_slide(prs,
        "Layer 5: Tissue Expression — Division of Labor at the Molecular Level",
        "genome_tissue_specialization.png",
        [
            "5,000 genes across 20 tissues analyzed",
            "20% of genes show tissue specificity tau > 0.95",
            "Machine model prediction: tau = 0 (no specialization) — falsified",
            "Fold enrichment in top tissue: 100-1,000x over mean",
            "Expression Gini > 0.8 for tissue-specific genes — extreme inequality",
            "Menger's spontaneous order: division of labor written in transcription",
        ])

    # --- Slide 18: Figure 15 — Convergent Evolution ---
    figure_slide(prs,
        "Layer 5: Convergent Evolution — Same Solutions Across Kingdoms",
        "genome_convergent_evolution.png",
        [
            "35 convergent events, 17 traits, divergence 8,000 to 1.5 billion years",
            "Bats + dolphins: same 14 amino acid substitutions in Prestin (95 Mya divergence)",
            "C4 photosynthesis: evolved independently 60+ times, same core enzymes",
            "EPAS1 altitude adaptation: same variants in Tibetans, Andeans, Ethiopians",
            "Convergence does not decline with divergence time — persists across all scales",
            "Solution landscape is structured, not flat — directed, not random",
        ])

    # --- Slide 19: Figure 16 — Genome Summary ---
    figure_slide(prs,
        "Layer 5: Genome-Wide Summary — Distributed Knowledge at Every Scale",
        "genome_distributed_summary.png",
        [
            "CpG enrichment ~10-25x over non-CpG — random prediction: equal rates",
            "Tissue specificity waterfall: machine model (tau=0) falsified entirely",
            "35 convergent evolution events across 1.5 billion years",
            "Three genome-wide patterns match three immune-level patterns",
            "Directed mutation, specialized expression, convergent discovery",
            "Same architecture at immune scale and genome scale — scale-invariant",
        ])

    # --- Slide 20: Synthesis ---
    synthesis_slide(prs)

    # --- Slide 21: Design Principles ---
    design_principles_slide(prs)

    # --- Slide 22: Closing ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_text_box(slide, Inches(1), Inches(1.5), Inches(11), Inches(1.2),
                 "The Garden Is Still Growing",
                 font_size=36, color=TITLE_COLOR, bold=True,
                 alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(1.5), Inches(3.0), Inches(10), Inches(2.5),
                 "The cell is not a chassis waiting to be programmed.\n"
                 "It is a running economy — 4,400 genes coordinating through distributed feedback.\n\n"
                 "The sagent reads the economy before entering it.\n"
                 "The planner overrides it.\n\n"
                 "The first job was gardener, not king.",
                 font_size=20, color=BODY_COLOR, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(1.5), Inches(5.8), Inches(10), Inches(0.8),
                 "Sage Clokey  |  BME 129C  |  UC Santa Cruz  |  Spring 2026",
                 font_size=14, color=DIM_COLOR, alignment=PP_ALIGN.CENTER)

    prs.save(str(OUTPUT))
    print(f"Saved: {OUTPUT}  ({len(prs.slides)} slides)")


if __name__ == "__main__":
    build()
