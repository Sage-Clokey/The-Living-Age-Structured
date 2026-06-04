"""
Generate Individual Panel PNGs + PowerPoint Presentation
=========================================================
Runs all analyses, saves each panel as its own PNG, then builds
a PowerPoint deck with figure + bullet points + speaker notes.

Usage:
    python paper/generate_presentation.py              # full run
    python paper/generate_presentation.py --quick       # skip API calls
    python paper/generate_presentation.py --pptx-only   # skip figure gen, just build PPTX from existing PNGs

BME 129C Capstone — Sage Clokey — Spring 2026
"""

from __future__ import annotations
import argparse
import sys
from pathlib import Path
import warnings
warnings.filterwarnings("ignore")

PROJECT_ROOT = Path(__file__).resolve().parent.parent.parent
sys.path.insert(0, str(PROJECT_ROOT))

import numpy as np
import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt

# ---------------------------------------------------------------------------
# Shared style
# ---------------------------------------------------------------------------
SPIRAL_GREEN = "#2d6a4f"
SPIRAL_MID   = "#52b788"
SPIRAL_LIGHT = "#95d5b2"
GOLD         = "#e9c46a"
RED          = "#e63946"
BLUE         = "#4361ee"
BACKGROUND   = "#0d1117"
PANEL_BG     = "#161b22"
TEXT_MAIN     = "#e6edf3"
TEXT_DIM      = "#8b949e"

FIG_DIR = PROJECT_ROOT / "paper" / "figures" / "individual"


def _style_ax(ax):
    ax.set_facecolor(PANEL_BG)
    ax.tick_params(colors=TEXT_DIM)
    for spine in ax.spines.values():
        spine.set_color(TEXT_DIM)


def _save(fig, name: str):
    FIG_DIR.mkdir(parents=True, exist_ok=True)
    out = FIG_DIR / f"{name}.png"
    fig.savefig(out, dpi=200, bbox_inches="tight", facecolor=BACKGROUND)
    plt.close(fig)
    print(f"  -> {out.name}")
    return out


# ===================================================================
# LAYER 1: TOPOLOGY — Individual panels
# ===================================================================

def gen_layer1_panels(bio_reports, ref_reports):
    """Split layer1_topology 3-panel into 3 individual PNGs."""

    # ---- Panel 1: Degree Distribution ----
    fig, ax = plt.subplots(figsize=(8, 6), facecolor=BACKGROUND)
    _style_ax(ax)

    bio_colors = [SPIRAL_GREEN, SPIRAL_MID, BLUE, SPIRAL_LIGHT, "#a8dadc"]
    ref_colors = [RED, GOLD, TEXT_DIM, "#f4a261", "#264653"]

    for i, r in enumerate(bio_reports):
        degrees = sorted(r.degree_sequence, reverse=True)
        if not degrees:
            continue
        unique, counts = np.unique(degrees, return_counts=True)
        prob = counts / counts.sum()
        c = bio_colors[i % len(bio_colors)]
        ax.scatter(unique, prob, color=c, s=30, alpha=0.8, label=r.name, zorder=3)

    if ref_reports:
        for i, r in enumerate(ref_reports):
            degrees = sorted(r.degree_sequence, reverse=True)
            if not degrees:
                continue
            unique, counts = np.unique(degrees, return_counts=True)
            prob = counts / counts.sum()
            c = ref_colors[i % len(ref_colors)]
            ax.scatter(unique, prob, color=c, s=20, alpha=0.6, marker="x", label=r.name, zorder=2)

    ax.set_xscale("log")
    ax.set_yscale("log")
    ax.set_xlabel("Degree (k)", color=TEXT_MAIN, fontsize=12)
    ax.set_ylabel("P(k)", color=TEXT_MAIN, fontsize=12)
    ax.set_title("Degree Distribution", color=TEXT_MAIN, fontsize=14, fontweight="bold")
    ax.legend(facecolor=PANEL_BG, edgecolor=TEXT_DIM, labelcolor=TEXT_MAIN, fontsize=8)
    _save(fig, "L1_degree_distribution")

    # ---- Panel 2: Betweenness Gini ----
    fig, ax = plt.subplots(figsize=(8, 6), facecolor=BACKGROUND)
    _style_ax(ax)

    all_reports = bio_reports + (ref_reports or [])
    names = [r.name.replace("_", " ")[:20] for r in all_reports]
    ginis = [r.betweenness_gini for r in all_reports]
    n_bio = len(bio_reports)
    colors = [bio_colors[i % len(bio_colors)] for i in range(n_bio)]
    colors += [ref_colors[i % len(ref_colors)] for i in range(len(ref_reports or []))]

    x = np.arange(len(names))
    ax.bar(x, ginis, color=colors, edgecolor="white", linewidth=0.5)
    ax.set_xticks(x)
    ax.set_xticklabels(names, rotation=45, ha="right", fontsize=8, color=TEXT_MAIN)
    ax.set_ylabel("Betweenness Gini Coefficient", color=TEXT_MAIN, fontsize=12)
    ax.set_title("Centrality Inequality", color=TEXT_MAIN, fontsize=14, fontweight="bold")
    ax.axhline(y=0.5, color=GOLD, linestyle="--", linewidth=1, alpha=0.5)
    ax.text(len(names) - 0.5, 0.52, "high centralization", color=GOLD, fontsize=8, ha="right")
    _save(fig, "L1_centrality_gini")

    # ---- Panel 3: Robustness Curves ----
    fig, ax = plt.subplots(figsize=(8, 6), facecolor=BACKGROUND)
    _style_ax(ax)

    for i, r in enumerate(bio_reports):
        c = bio_colors[i % len(bio_colors)]
        if r.robustness_random:
            fracs = np.linspace(0, 0.95, len(r.robustness_random))
            ax.plot(fracs, r.robustness_random, color=c, linewidth=2, alpha=0.6,
                    linestyle="--", label=f"{r.name} (random)")
        if r.robustness_targeted:
            fracs = np.linspace(0, 0.95, len(r.robustness_targeted))
            ax.plot(fracs, r.robustness_targeted, color=c, linewidth=2,
                    label=f"{r.name} (targeted)")

    if ref_reports:
        for i, r in enumerate(ref_reports):
            c = ref_colors[i % len(ref_colors)]
            if r.robustness_targeted:
                fracs = np.linspace(0, 0.95, len(r.robustness_targeted))
                ax.plot(fracs, r.robustness_targeted, color=c, linewidth=1.5, alpha=0.5,
                        linestyle=":", label=f"{r.name}")

    ax.axhline(y=0.5, color=TEXT_DIM, linestyle="--", linewidth=0.8, alpha=0.5)
    ax.set_xlabel("Fraction of Nodes Removed", color=TEXT_MAIN, fontsize=12)
    ax.set_ylabel("Giant Component / Original", color=TEXT_MAIN, fontsize=12)
    ax.set_title("Robustness Under Attack", color=TEXT_MAIN, fontsize=14, fontweight="bold")
    ax.legend(facecolor=PANEL_BG, edgecolor=TEXT_DIM, labelcolor=TEXT_MAIN, fontsize=7, ncol=2)
    ax.set_xlim(0, 0.95)
    ax.set_ylim(0, 1.05)
    _save(fig, "L1_robustness_curves")


# ===================================================================
# LAYER 1 MOTIFS — Individual panels
# ===================================================================

def gen_motif_panels(report):
    """Split motif_analysis 2-panel into 2 individual PNGs."""

    # ---- Panel 1: Z-score bar chart ----
    fig, ax = plt.subplots(figsize=(10, 6), facecolor=BACKGROUND)
    _style_ax(ax)

    triad_types = sorted(report.z_scores.keys(), key=lambda t: report.z_scores[t], reverse=True)
    z_vals = [report.z_scores[t] for t in triad_types]
    bar_colors = []
    for z in z_vals:
        if z > 2:
            bar_colors.append(SPIRAL_MID)
        elif z < -2:
            bar_colors.append(RED)
        else:
            bar_colors.append(TEXT_DIM)

    x = np.arange(len(triad_types))
    ax.bar(x, z_vals, color=bar_colors, edgecolor="white", linewidth=0.5)
    ax.set_xticks(x)
    ax.set_xticklabels(triad_types, rotation=45, ha="right", fontsize=8, color=TEXT_MAIN)
    ax.axhline(y=2.0, color=SPIRAL_MID, linestyle="--", linewidth=1, alpha=0.5, label="Motif (Z>2)")
    ax.axhline(y=-2.0, color=RED, linestyle="--", linewidth=1, alpha=0.5, label="Anti-motif (Z<-2)")
    ax.axhline(y=0, color=TEXT_DIM, linewidth=0.5)
    ax.set_xlabel("Triad Type", color=TEXT_MAIN, fontsize=12)
    ax.set_ylabel("Z-score", color=TEXT_MAIN, fontsize=12)
    ax.set_title("Network Motif Detection (Triadic Census vs Null Model)",
                 color=TEXT_MAIN, fontsize=14, fontweight="bold")
    ax.legend(facecolor=PANEL_BG, edgecolor=TEXT_DIM, labelcolor=TEXT_MAIN, fontsize=9)
    _save(fig, "L1_motif_zscore")

    # ---- Panel 2: FFL diagram ----
    fig, ax = plt.subplots(figsize=(8, 7), facecolor=BACKGROUND)
    ax.set_facecolor(PANEL_BG)
    ax.set_xlim(-0.5, 1.5)
    ax.set_ylim(-0.3, 1.3)
    ax.set_aspect("equal")
    ax.axis("off")

    # Triangle: A (top), B (bottom-left), C (bottom-right)
    positions = {"A": (0.5, 1.1), "B": (0.0, 0.3), "C": (1.0, 0.3)}
    labels = {
        "A": "Transcription\nFactor (A)\n'Price Setter'",
        "B": "Intermediate (B)\n'Local\nEntrepreneur'",
        "C": "Target Gene (C)\n'Market Actor'",
    }

    for node, (x, y) in positions.items():
        circle = plt.Circle((x, y), 0.12, color=SPIRAL_MID, ec="white", lw=2, zorder=3)
        ax.add_patch(circle)
        ax.text(x, y, node, ha="center", va="center", color="white",
                fontsize=16, fontweight="bold", zorder=4)
        if node == "A":
            ax.text(x, y + 0.2, labels[node], ha="center", va="bottom",
                    color=TEXT_MAIN, fontsize=9)
        elif node == "B":
            ax.text(x - 0.18, y, labels[node], ha="right", va="center",
                    color=TEXT_MAIN, fontsize=9)
        else:
            ax.text(x + 0.18, y, labels[node], ha="left", va="center",
                    color=TEXT_MAIN, fontsize=9)

    # Arrows
    from matplotlib.patches import FancyArrowPatch
    for (start, end, label) in [("A", "B", "regulates"), ("B", "C", "activates"), ("A", "C", "shortcut")]:
        sx, sy = positions[start]
        ex, ey = positions[end]
        dx, dy = ex - sx, ey - sy
        dist = np.sqrt(dx**2 + dy**2)
        r = 0.13
        sx2, sy2 = sx + r * dx / dist, sy + r * dy / dist
        ex2, ey2 = ex - r * dx / dist, ey - r * dy / dist
        style = "dashed" if label == "shortcut" else "solid"
        color = GOLD if label == "shortcut" else SPIRAL_LIGHT
        arrow = FancyArrowPatch((sx2, sy2), (ex2, ey2), arrowstyle="-|>",
                                mutation_scale=15, color=color, linewidth=2,
                                linestyle=style, zorder=2)
        ax.add_patch(arrow)
        mx, my = (sx2 + ex2) / 2, (sy2 + ey2) / 2
        offset = 0.08 if label != "shortcut" else -0.08
        ax.text(mx + offset, my + offset, label, color=color, fontsize=8,
                ha="center", va="center", style="italic")

    ax.set_title("Feed-Forward Loop (030T) — Hayekian Price Signal",
                 color=TEXT_MAIN, fontsize=14, fontweight="bold", pad=10)

    ffl_z = report.ffl_z_score if hasattr(report, "ffl_z_score") else report.z_scores.get("030T", 0)
    ffl_n = report.ffl_count if hasattr(report, "ffl_count") else report.triad_counts.get("030T", 0)
    stats_text = f"FFL count: {ffl_n}\nZ-score: {ffl_z:.1f}"
    ax.text(0.5, -0.15, stats_text, ha="center", va="top", color=GOLD,
            fontsize=11, fontweight="bold", transform=ax.transAxes)

    _save(fig, "L1_motif_ffl_diagram")


# ===================================================================
# LAYER 1 SELF-REGULATION — Individual panels
# ===================================================================

def gen_self_regulation_panels(bio_reports, ref_reports):
    """Split self_regulation 3-panel into individual PNGs."""

    bio_colors = [SPIRAL_GREEN, SPIRAL_MID, SPIRAL_LIGHT, BLUE, "#a8dadc"]
    ref_colors = [RED, GOLD, TEXT_DIM, "#f4a261", "#264653"]

    # ---- Panel 1: Spearman Correlation Bars ----
    fig, ax = plt.subplots(figsize=(8, 6), facecolor=BACKGROUND)
    _style_ax(ax)

    all_reports = bio_reports + ref_reports
    names = [r.name.replace("_", " ")[:20] for r in all_reports]
    spearman_vals = [r.degree_betweenness_spearman for r in all_reports]
    n_bio = len(bio_reports)
    colors = [bio_colors[i % len(bio_colors)] for i in range(n_bio)]
    colors += [ref_colors[i % len(ref_colors)] for i in range(len(ref_reports))]

    x = np.arange(len(names))
    ax.bar(x, spearman_vals, color=colors, edgecolor="white", linewidth=0.5)
    ax.set_xticks(x)
    ax.set_xticklabels(names, rotation=45, ha="right", fontsize=8, color=TEXT_MAIN)
    ax.set_ylabel("Spearman(degree, betweenness)", color=TEXT_MAIN, fontsize=11)
    ax.set_title("Degree-Betweenness Correlation", color=TEXT_MAIN, fontsize=14, fontweight="bold")
    ax.axhline(y=0.8, color=RED, linestyle="--", linewidth=1, alpha=0.5)
    ax.text(len(names) - 0.5, 0.82, "high correlation =\nno self-regulation", color=RED, fontsize=7, ha="right")
    ax.axhline(y=0.5, color=GOLD, linestyle=":", linewidth=1, alpha=0.5)
    ax.text(len(names) - 0.5, 0.52, "self-regulating zone", color=GOLD, fontsize=7, ha="right")
    _save(fig, "L1_self_reg_spearman")

    # ---- Panel 2: Gini Ratio Bars ----
    fig, ax = plt.subplots(figsize=(8, 6), facecolor=BACKGROUND)
    _style_ax(ax)

    gini_ratios = [r.gini_ratio for r in all_reports]
    ax.bar(x, gini_ratios, color=colors, edgecolor="white", linewidth=0.5)
    ax.set_xticks(x)
    ax.set_xticklabels(names, rotation=45, ha="right", fontsize=8, color=TEXT_MAIN)
    ax.set_ylabel("Gini Ratio (BC_gini / Deg_gini)", color=TEXT_MAIN, fontsize=11)
    ax.set_title("Self-Regulation Index", color=TEXT_MAIN, fontsize=14, fontweight="bold")
    ax.axhline(y=1.0, color=GOLD, linestyle="--", linewidth=1, alpha=0.7)
    ax.text(len(names) - 0.5, 1.03, "ratio=1.0: no self-regulation", color=GOLD, fontsize=7, ha="right")
    ax.text(0.5, min(gini_ratios) - 0.05, "<1.0 = self-regulating (invisible hand active)",
            color=SPIRAL_LIGHT, fontsize=8, ha="left")
    _save(fig, "L1_self_reg_gini_ratio")


# ===================================================================
# LAYER 2: ECONOMY — Individual panels
# ===================================================================

def gen_economy_panels(dist, cent, robustness):
    """Split layer2_economy 3-panel into 3 individual PNGs."""

    # ---- Panel 1: GDP Over Time ----
    fig, ax = plt.subplots(figsize=(8, 6), facecolor=BACKGROUND)
    _style_ax(ax)

    steps = range(len(dist.gdp_history))
    ax.plot(steps, dist.gdp_history, color=SPIRAL_MID, linewidth=2, label="Distributed (biological)")
    ax.plot(steps, cent.gdp_history, color=RED, linewidth=2, alpha=0.8, label="Centralized (planned)")
    ax.set_xlabel("Time Step", color=TEXT_MAIN, fontsize=12)
    ax.set_ylabel("GDP (total metabolite output)", color=TEXT_MAIN, fontsize=12)
    ax.set_title("GDP Over Time: Invisible Hand vs Central Planner",
                 color=TEXT_MAIN, fontsize=14, fontweight="bold")
    ax.legend(facecolor=PANEL_BG, edgecolor=TEXT_DIM, labelcolor=TEXT_MAIN)
    _save(fig, "L2_gdp_over_time")

    # ---- Panel 2: Perturbation Robustness ----
    fig, ax = plt.subplots(figsize=(8, 6), facecolor=BACKGROUND)
    _style_ax(ax)

    if robustness["distributed"]:
        pathways = list(robustness["distributed"].keys())
        short_names = [n.replace(" pathway", "").replace(" synthesis", "")[:20] for n in pathways]
        dist_scores = [robustness["distributed"][p] for p in pathways]
        cent_scores = [robustness["centralized"][p] for p in pathways]

        y = np.arange(len(pathways))
        width = 0.35
        ax.barh(y - width / 2, dist_scores, width, color=SPIRAL_MID, label="Distributed")
        ax.barh(y + width / 2, cent_scores, width, color=RED, alpha=0.7, label="Centralized")
        ax.set_yticks(y)
        ax.set_yticklabels(short_names, fontsize=7, color=TEXT_MAIN)
        ax.axvline(x=1.0, color=TEXT_DIM, linestyle="--", linewidth=0.8, alpha=0.5)
        ax.set_xlabel("GDP Fraction Retained", color=TEXT_MAIN, fontsize=12)
        ax.legend(facecolor=PANEL_BG, edgecolor=TEXT_DIM, labelcolor=TEXT_MAIN, fontsize=9)

    ax.set_title("Robustness: GDP After Removing One Agent",
                 color=TEXT_MAIN, fontsize=14, fontweight="bold")
    _save(fig, "L2_perturbation_robustness")

    # ---- Panel 3: Production Rate Convergence ----
    fig, ax = plt.subplots(figsize=(8, 6), facecolor=BACKGROUND)
    _style_ax(ax)

    agent_colors = plt.cm.Set2(np.linspace(0, 1, len(dist.agents)))
    for i, agent in enumerate(dist.agents):
        if agent.production_history:
            ax.plot(agent.production_history, color=agent_colors[i], linewidth=1, alpha=0.7,
                    label=agent.name[:18] if i < 8 else None)
    ax.set_xlabel("Time Step", color=TEXT_MAIN, fontsize=12)
    ax.set_ylabel("Production Rate", color=TEXT_MAIN, fontsize=12)
    ax.set_title("Distributed Rate Convergence (Price Discovery)",
                 color=TEXT_MAIN, fontsize=14, fontweight="bold")
    ax.legend(facecolor=PANEL_BG, edgecolor=TEXT_DIM, labelcolor=TEXT_MAIN, fontsize=6,
              loc="upper right", ncol=2)
    _save(fig, "L2_rate_convergence")


# ===================================================================
# LAYER 2 FBA — Individual panels
# ===================================================================

def gen_fba_panels(report):
    """Split layer2_fba_analysis 3-panel into 3 individual PNGs."""

    # ---- Panel 1: FBA vs Keio Scatter ----
    fig, ax = plt.subplots(figsize=(8, 7), facecolor=BACKGROUND)
    _style_ax(ax)

    fba_vals, keio_vals, colors = [], [], []
    for gene_id, info in report.knockout_results.items():
        fba_vals.append(info["fba_growth"])
        keio_vals.append(info["keio_growth"])
        colors.append(SPIRAL_MID if info["essential_fba"] == info["essential_keio"] else RED)

    ax.scatter(keio_vals, fba_vals, c=colors, s=60, alpha=0.7, edgecolors="white", linewidth=0.5)
    ax.plot([0, 1], [0, 1], color=TEXT_DIM, linestyle="--", linewidth=1, alpha=0.5)
    ax.axhline(y=0.05, color=GOLD, linestyle=":", linewidth=1, alpha=0.7, label="Essentiality threshold")
    ax.axvline(x=0.05, color=GOLD, linestyle=":", linewidth=1, alpha=0.7)
    ax.set_xlabel("Keio Growth (actual)", color=TEXT_MAIN, fontsize=12)
    ax.set_ylabel("FBA Growth (predicted)", color=TEXT_MAIN, fontsize=12)
    ax.set_title(f"Omniscient Planner vs Reality (acc={report.accuracy:.0%})",
                 color=TEXT_MAIN, fontsize=14, fontweight="bold")
    ax.legend(facecolor=PANEL_BG, edgecolor=TEXT_DIM, labelcolor=TEXT_MAIN, fontsize=9)
    _save(fig, "L2_fba_knockout_scatter")

    # ---- Panel 2: Confusion Matrix ----
    fig, ax = plt.subplots(figsize=(7, 6), facecolor=BACKGROUND)
    _style_ax(ax)

    cm = np.array([
        [report.true_positives, report.false_negatives],
        [report.false_positives, report.true_negatives],
    ])
    im = ax.imshow(cm, cmap="YlGn", aspect="auto")
    labels = [
        [f"TP\n{report.true_positives}", f"FN\n{report.false_negatives}"],
        [f"FP\n{report.false_positives}", f"TN\n{report.true_negatives}"],
    ]
    for i in range(2):
        for j in range(2):
            tc = "white" if cm[i, j] > cm.max() / 2 else "black"
            ax.text(j, i, labels[i][j], ha="center", va="center",
                    color=tc, fontsize=16, fontweight="bold")

    ax.set_xticks([0, 1])
    ax.set_xticklabels(["Essential\n(Keio)", "Viable\n(Keio)"], color=TEXT_MAIN, fontsize=11)
    ax.set_yticks([0, 1])
    ax.set_yticklabels(["Essential\n(FBA)", "Viable\n(FBA)"], color=TEXT_MAIN, fontsize=11)
    ax.set_title("Mises' Calculation Problem\nWhere the Planner Fails",
                 color=TEXT_MAIN, fontsize=14, fontweight="bold")
    _save(fig, "L2_fba_confusion_matrix")

    # ---- Panel 3: Shadow Prices ----
    fig, ax = plt.subplots(figsize=(8, 7), facecolor=BACKGROUND)
    _style_ax(ax)

    if report.key_shadow_prices:
        mets = list(report.key_shadow_prices.keys())[:15]
        prices = [report.key_shadow_prices[m] for m in mets]
        short_mets = [m.replace("_c", "").replace("_e", "").replace("_p", "")[:25] for m in mets]
        bar_colors = [SPIRAL_MID if p > 0 else RED for p in prices]
        y_pos = np.arange(len(mets))
        ax.barh(y_pos, prices, color=bar_colors, alpha=0.8)
        ax.set_yticks(y_pos)
        ax.set_yticklabels(short_mets, fontsize=8, color=TEXT_MAIN)
        ax.invert_yaxis()

    ax.set_xlabel("Shadow Price (marginal growth value)", color=TEXT_MAIN, fontsize=12)
    ax.set_title("Hayekian Price Signals in FBA\n(marginal value of each metabolite)",
                 color=TEXT_MAIN, fontsize=14, fontweight="bold")
    _save(fig, "L2_fba_shadow_prices")


def gen_fba_perturbation_panels(perturbation_results):
    """Split fba_perturbation 3-panel into 3 individual PNGs."""

    austrian_map = {
        "carbon_switch": ("Hayek's Knowledge Problem\nGlucose -> Acetate", "L2_fba_pert_carbon"),
        "anaerobic": ("Mises' Calculation Under Crisis\nAerobic -> Anaerobic", "L2_fba_pert_anaerobic"),
        "nitrogen_limit": ("Kirzner's Entrepreneurial Alertness\nNitrogen Limitation", "L2_fba_pert_nitrogen"),
    }

    for pname, pdata in perturbation_results.items():
        title, fname = austrian_map.get(pname, (pname, f"L2_fba_pert_{pname}"))

        fig, ax = plt.subplots(figsize=(7, 6), facecolor=BACKGROUND)
        _style_ax(ax)

        if "glucose_growth" in pdata:
            vals = [pdata["glucose_growth"], pdata["acetate_growth"]]
            labels = ["Glucose\n(baseline)", "Acetate\n(perturbed)"]
        elif "aerobic_growth" in pdata:
            vals = [pdata["aerobic_growth"], pdata["anaerobic_growth"]]
            labels = ["Aerobic\n(baseline)", "Anaerobic\n(perturbed)"]
        else:
            vals = [pdata["normal_growth"], pdata["limited_growth"]]
            labels = ["Normal N\n(baseline)", "Limited N\n(perturbed)"]

        ax.bar([0, 1], vals, color=[SPIRAL_MID, GOLD], width=0.6,
               edgecolor="white", linewidth=0.5)
        ax.set_xticks([0, 1])
        ax.set_xticklabels(labels, color=TEXT_MAIN, fontsize=11)
        ax.set_ylabel("FBA Growth Rate (h\u207b\u00b9)", color=TEXT_MAIN, fontsize=12)
        ax.set_title(title, color=TEXT_MAIN, fontsize=13, fontweight="bold")

        ratio = pdata.get("growth_ratio", 0)
        ax.text(0.5, 0.92, f"Ratio: {ratio:.1%}", transform=ax.transAxes,
                ha="center", va="top", color=GOLD, fontsize=14, fontweight="bold")

        _save(fig, fname)


# ===================================================================
# LAYER 3: TRADE — Individual panels
# ===================================================================

def gen_trade_panels():
    """Split layer3_trade_network 3-panel into 3 individual PNGs."""
    import networkx as nx
    from layer3_trade.trade_network import (
        build_trade_network, compute_trade_matrix, comparative_advantage_table,
    )

    G = build_trade_network()
    names, matrix = compute_trade_matrix()

    # ---- Panel 1: Trade Network Graph ----
    fig, ax = plt.subplots(figsize=(9, 7), facecolor=BACKGROUND)
    _style_ax(ax)
    ax.axis("off")

    kingdom_colors = {}
    for node in G.nodes():
        data = G.nodes[node]
        kingdom = data.get("kingdom", "unknown")
        if kingdom == "prokaryote":
            kingdom_colors[node] = RED
        elif kingdom == "fungal":
            kingdom_colors[node] = SPIRAL_GREEN
        elif kingdom == "plant":
            kingdom_colors[node] = GOLD
        else:
            kingdom_colors[node] = BLUE

    pos = nx.spring_layout(G, seed=42, k=2)
    node_colors = [kingdom_colors.get(n, TEXT_DIM) for n in G.nodes()]
    node_sizes = [400 + 200 * len(G.nodes[n].get("capabilities", [])) for n in G.nodes()]

    max_ease = max((d.get("trade_ease", 1) for _, _, d in G.edges(data=True)), default=1)
    edge_widths = [1 + 4 * (d.get("trade_ease", 0.5) / max_ease) for _, _, d in G.edges(data=True)]

    nx.draw_networkx_edges(G, pos, width=edge_widths, edge_color=SPIRAL_LIGHT, alpha=0.4, ax=ax)
    nx.draw_networkx_nodes(G, pos, node_color=node_colors, node_size=node_sizes,
                           edgecolors="white", linewidths=1, ax=ax)
    nx.draw_networkx_labels(G, pos, font_size=8, font_color=TEXT_MAIN, ax=ax)

    ax.set_title("Cross-Species Trade Network", color=TEXT_MAIN, fontsize=14, fontweight="bold")
    _save(fig, "L3_trade_network_graph")

    # ---- Panel 2: Trade Cost Heatmap ----
    fig, ax = plt.subplots(figsize=(8, 7), facecolor=BACKGROUND)
    _style_ax(ax)

    im = ax.imshow(matrix, cmap="YlOrRd", aspect="auto")
    ax.set_xticks(range(len(names)))
    ax.set_xticklabels(names, rotation=45, ha="right", fontsize=8, color=TEXT_MAIN)
    ax.set_yticks(range(len(names)))
    ax.set_yticklabels(names, fontsize=8, color=TEXT_MAIN)
    for i in range(len(names)):
        for j in range(len(names)):
            ax.text(j, i, f"{matrix[i, j]:.2f}", ha="center", va="center",
                    fontsize=6, color="white" if matrix[i, j] > matrix.max() / 2 else "black")
    plt.colorbar(im, ax=ax, label="Trade Cost", shrink=0.8)
    ax.set_title("Trade Barriers Heatmap", color=TEXT_MAIN, fontsize=14, fontweight="bold")
    _save(fig, "L3_trade_cost_heatmap")

    # ---- Panel 3: Comparative Advantage Table ----
    fig, ax = plt.subplots(figsize=(8, 6), facecolor=BACKGROUND)
    ax.set_facecolor(PANEL_BG)
    ax.axis("off")

    table_data = comparative_advantage_table()
    y = 0.9
    ax.text(0.05, 0.95, "Organism", color=GOLD, fontsize=12, fontweight="bold",
            transform=ax.transAxes, va="top")
    ax.text(0.35, 0.95, "Comparative Advantage (Exports)", color=GOLD, fontsize=12,
            fontweight="bold", transform=ax.transAxes, va="top")

    for org, caps in table_data.items():
        y -= 0.08
        ax.text(0.05, y, org, color=TEXT_MAIN, fontsize=10, transform=ax.transAxes, va="top")
        ax.text(0.35, y, ", ".join(caps), color=SPIRAL_LIGHT, fontsize=9,
                transform=ax.transAxes, va="top")

    ax.set_title("Comparative Advantage", color=TEXT_MAIN, fontsize=14,
                 fontweight="bold", pad=20)
    _save(fig, "L3_comparative_advantage_table")


# ===================================================================
# POWERPOINT BUILDER
# ===================================================================

# Slide definitions: (filename, title, bullets, notes)
SLIDE_DEFS = [
    # --- Title slide ---
    (None, "Living Systems as Decentralized Economies",
     ["BME 129C Capstone — Sage Clokey — Spring 2026",
      "UC Santa Cruz, Biomolecular Engineering"],
     "This capstone uses biological data across three layers of analysis to test "
     "whether living systems operate as decentralized economies. Each layer maps "
     "to specific Austrian economic principles from Hayek, Mises, Menger, Rothbard, "
     "and Kirzner. The thesis: biology is not a machine — it is a market."),

    # --- Layer 1: Topology ---
    ("L1_degree_distribution", "Degree Distribution — No Master Node",
     ["E. coli GRN follows a heavy-tailed (scale-free-like) distribution",
      "No single node controls the network",
      "A few coordinating hubs, many small players — like a market with large firms but no central planner",
      "Reference networks (star, lattice) show sharply different distributions"],
     "The degree distribution of the E. coli gene regulatory network is plotted on log-log axes. "
     "Biological networks show a heavy-tailed distribution where most nodes have few connections "
     "and a small number of hubs have many. This is the hallmark of a decentralized system: "
     "influence is distributed, not concentrated in a single master node. Comparison with reference "
     "architectures (star = centralized, random = no structure, BA = scale-free model) shows biology "
     "sits between pure hierarchy and pure randomness — spontaneous order, not design."),

    ("L1_centrality_gini", "Centrality Inequality — Betweenness Gini",
     ["Gini coefficient measures inequality in information flow control",
      "Biological networks: moderate Gini (hubs exist but don't monopolize)",
      "Star network: Gini near 1.0 (one node controls everything)",
      "Biology self-regulates against centralization"],
     "The betweenness Gini coefficient measures how unequally information flow is distributed "
     "across the network. A Gini of 1.0 means one node controls all shortest paths (star topology). "
     "A Gini of 0 means perfectly equal control. Biological networks sit in the moderate range — "
     "hubs exist (division of labor) but they don't monopolize information flow. This is the "
     "invisible hand in topology: the network self-regulates to prevent any single node from "
     "becoming a bottleneck."),

    ("L1_robustness_curves", "Robustness Under Attack — Albert et al. 2000",
     ["Random failure: network degrades gracefully (like a market losing random firms)",
      "Targeted attack: vulnerable to hub removal (like removing key entrepreneurs)",
      "Centralized networks collapse faster under targeted attack",
      "Biology trades some targeted vulnerability for massive random resilience"],
     "Following the Albert et al. (2000) method, nodes are removed either randomly or by "
     "targeting the highest-degree node at each step (recalculated after each removal). "
     "The y-axis shows the fraction of the network still connected. Biological networks survive "
     "random failure far better than centralized architectures, but are vulnerable to targeted "
     "hub removal — exactly the tradeoff Austrian economists predict: a free market is resilient "
     "to random shocks but vulnerable to removing key entrepreneurs. A planned economy collapses "
     "faster because its central bureau is a single point of failure."),

    # --- Layer 1: Motifs ---
    ("L1_motif_zscore", "Network Motifs — Alon 2002 Triadic Census",
     ["Z-scores compare real E. coli GRN to 1,000 degree-preserving randomizations",
      "Feed-forward loop (030T) is massively over-represented",
      "These are evolved coordination motifs — the 'grammar' of the network",
      "Anti-motifs (negative Z) are patterns biology actively avoids"],
     "The triadic census counts every possible 3-node subgraph pattern in the real E. coli GRN "
     "and compares to 1,000 degree-preserving random networks. Z-scores above 2 indicate the "
     "pattern appears far more than expected by chance (motif = evolved function). Z-scores below "
     "-2 indicate the pattern is actively suppressed (anti-motif). The feed-forward loop (030T) "
     "has the highest Z-score, confirming Alon et al. (2002): this is the recurring regulatory "
     "grammar that enables fast, local information propagation without central coordination."),

    ("L1_motif_ffl_diagram", "Feed-Forward Loop — Hayekian Price Signal",
     ["Structure: A regulates B, B activates C, A also shortcuts directly to C",
      "Maps to Hayek: TF A is the 'price setter', B is the 'local entrepreneur', C is the 'market actor'",
      "Fast regulatory shortcut without waiting for central command",
      "Most over-represented motif in E. coli — evolved for decentralized coordination"],
     "The feed-forward loop (FFL) is the most statistically significant network motif in E. coli. "
     "In the FFL, transcription factor A regulates an intermediate B, which in turn activates "
     "target gene C. But A also directly regulates C via a shortcut path. This creates a "
     "Hayekian price signal: local information (metabolite levels sensed by A) propagates "
     "rapidly through two parallel channels. The 'price setter' (A) broadcasts the signal, "
     "the 'local entrepreneur' (B) amplifies it with additional context, and the 'market actor' "
     "(C) integrates both streams. No central planner needed — just local feedback loops."),

    # --- Layer 1: Self-Regulation ---
    ("L1_self_reg_spearman", "Degree-Betweenness Correlation — Self-Regulation",
     ["Low Spearman rho: high degree does NOT guarantee high betweenness",
      "Biological networks show lower correlation than star or hub-spoke",
      "Hubs don't monopolize information flow — the invisible hand in topology",
      "Based on Topirceanu's WBPA mechanism (2018)"],
     "For each network, we compute the Spearman rank correlation between node degree (how many "
     "connections) and betweenness centrality (how much information flows through it). In a "
     "star network, this correlation is near 1.0 — the central hub has both the most connections "
     "and the most information control. In biological networks, the correlation is significantly "
     "lower: a node can have many connections without monopolizing information flow. This is "
     "Topirceanu's WBPA mechanism: biology grows by betweenness-preferential attachment, which "
     "automatically distributes control even as hubs form. The network self-regulates against "
     "centralization — the invisible hand in topology."),

    ("L1_self_reg_gini_ratio", "Self-Regulation Index — Gini Ratio",
     ["Gini ratio = BC_gini / Degree_gini",
      "Below 1.0: betweenness is MORE equal than degree would predict",
      "Biological networks score below 1.0 — they actively resist centralization",
      "Star/hub-spoke networks score at or above 1.0 — no self-regulation"],
     "The self-regulation index is the ratio of betweenness Gini to degree Gini. If this ratio "
     "is below 1.0, it means betweenness centrality is distributed MORE equally than degree "
     "alone would predict. In other words, even though some nodes have many connections (high "
     "degree), the network routes information around them so they don't become bottlenecks. "
     "Biological networks consistently score below 1.0 — they actively resist centralization. "
     "Engineered hierarchies (star, hub-spoke) score at or above 1.0 — they have no such "
     "self-regulation. This is the quantitative proof that biology employs an invisible hand "
     "mechanism to prevent monopoly."),

    # --- Layer 2: Economy ---
    ("L2_gdp_over_time", "GDP Over Time: Invisible Hand vs Central Planner",
     ["Distributed (green): each pathway adjusts production based on local feedback",
      "Centralized (red): a single allocator assigns production rates top-down",
      "Both converge, but distributed uses price discovery (Hayekian local signals)",
      "12 metabolic pathways modeled as economic agents with ATP as currency"],
     "Each of the 12 metabolic pathways from the adaptive_Automation genome design system is "
     "modeled as an economic agent with inputs (substrates consumed), outputs (products), and "
     "an ATP budget (currency). In the distributed regime, each agent adjusts its own production "
     "rate based on local metabolite concentrations — if products pile up, slow down; if substrates "
     "are abundant, speed up. No planner coordinates this. In the centralized regime, a single "
     "allocator ranks agents by ATP efficiency and assigns production rates top-down. Both reach "
     "equilibrium, but the distributed system finds it through price discovery — the same mechanism "
     "Hayek described for market economies."),

    ("L2_perturbation_robustness", "Perturbation Robustness — The Market Adapts",
     ["Remove one pathway agent at a time, measure GDP fraction retained",
      "Distributed retains more GDP across ALL knockouts",
      "The market adapts through local rerouting — no re-planning needed",
      "The plan breaks because it can't redistribute without its allocator"],
     "For each of the 12 pathway agents, we remove it from the simulation and measure the "
     "fraction of GDP that the remaining agents can maintain. The distributed system retains "
     "more GDP in every case because the remaining agents independently discover new equilibria "
     "through local feedback. The centralized system loses more because its fixed allocation "
     "formula can't adapt to the loss — it was optimized for the pre-perturbation state. This is "
     "the core Austrian prediction: distributed systems are more robust to unexpected disruptions "
     "because each agent carries local knowledge about its own environment."),

    ("L2_rate_convergence", "Price Discovery — Spontaneous Order in Action",
     ["Each line is one pathway agent's production rate over time",
      "Start chaotic, converge to stable rates WITHOUT any coordinator",
      "This IS spontaneous order — the invisible hand finding equilibrium",
      "Rate stability emerges from local feedback, not central planning"],
     "Each colored line traces a single pathway agent's production rate over 200 time steps. "
     "Initially the rates are chaotic as agents react to rapidly changing metabolite pools. "
     "Over time, they converge to stable, differentiated rates — each agent producing at exactly "
     "the rate that matches the economy's demand for its products. No coordinator sets these "
     "rates. They emerge from individual agents responding to local price signals (metabolite "
     "concentrations). This IS Hayek's spontaneous order: stable, efficient allocation arising "
     "from decentralized decisions, not central planning."),

    # --- Layer 2: FBA ---
    ("L2_fba_knockout_scatter", "FBA vs Keio Collection — The Omniscient Planner Tested",
     ["Each dot = one gene knockout (40 genes tested)",
      "X-axis = actual growth (Keio collection, Baba et al. 2006)",
      "Y-axis = FBA-predicted growth (iML1515, 2,712 reactions)",
      "Green = correct prediction, Red = planner failure",
      "70% accuracy — the planner gets most right but NOT all"],
     "Flux Balance Analysis (FBA) is the strongest possible central planner: it solves a linear "
     "program across all 2,712 reactions in E. coli's genome-scale metabolic model (iML1515, "
     "Monk et al. 2017) simultaneously, with perfect knowledge of every stoichiometric constraint. "
     "We compare its gene knockout predictions against the Keio collection — systematic single-gene "
     "deletions with measured growth rates (Baba et al. 2006). The planner achieves 70% accuracy: "
     "high, but not 100%. Each red dot is a gene where the omniscient planner fails — Mises was "
     "right that perfect economic calculation is impossible, even with complete information."),

    ("L2_fba_confusion_matrix", "Mises' Calculation Problem — Where the Planner Fails",
     ["TP: Both FBA and Keio agree the gene is essential",
      "TN: Both agree the gene is non-essential",
      "FN: FBA says 'viable' but cell dies — Hayek's local knowledge gap",
      "FP: FBA says 'essential' but cell reroutes — Kirzner's adaptation",
      "Each error type maps to a specific Austrian failure mode"],
     "The confusion matrix reveals two distinct types of planner failure. False negatives (FN) = "
     "FBA predicts the cell will survive the knockout but it actually dies. These genes are essential "
     "because of REGULATORY requirements that the LP cannot encode — allosteric feedback, gene "
     "expression timing, protein folding dependencies. This is Hayek's knowledge problem: the "
     "planner has all the stoichiometric data but misses the local, tacit knowledge embedded in "
     "regulatory networks. False positives (FP) = FBA predicts essential but the cell reroutes "
     "through alternative pathways. This is Kirzner's entrepreneurial discovery: the cell finds "
     "solutions the planner didn't anticipate."),

    ("L2_fba_shadow_prices", "Hayekian Price Signals in FBA",
     ["Shadow prices = dual variables = marginal value of each metabolite",
      "High shadow price = scarce and valuable (Menger's subjective value)",
      "Zero shadow price = abundant, not constraining (no economic profit)",
      "The irony: FBA MUST compute these prices to solve the optimization",
      "This proves Hayek's point that prices contain essential information"],
     "Shadow prices are the dual variables of the linear program — they tell you how much one "
     "additional unit of each metabolite would increase the objective (growth rate). This IS "
     "the Hayekian price signal computed by the planner. Metabolites with high absolute shadow "
     "prices are the bottlenecks — scarce resources that limit growth. Metabolites with zero "
     "shadow prices are abundant and not constraining. The deep irony: FBA must compute these "
     "prices internally to solve its optimization. Hayek argued that market prices are essential "
     "because they compress dispersed information into a single signal. FBA proves his point: "
     "even an omniscient planner needs price signals to allocate resources."),

    ("L2_fba_pert_carbon", "FBA Perturbation: Glucose to Acetate (Hayek)",
     ["FBA instantly re-optimizes for the new carbon source",
      "Real E. coli shows a diauxic lag phase before switching",
      "Planner assumes instant knowledge; biology must DISCOVER the change",
      "Growth ratio: ~24% — massive efficiency drop on acetate"],
     "When glucose is removed and acetate is provided as the sole carbon source, FBA instantly "
     "re-solves the LP and finds the optimal flux distribution for acetate metabolism. But real "
     "E. coli doesn't switch instantly — it undergoes a diauxic shift with a lag phase as it "
     "upregulates acetate utilization genes (aceA, aceB, acs) through the glyoxylate shunt. "
     "The planner assumes instant, perfect knowledge of the new environment. The cell must "
     "discover it through distributed sensing — CRP-cAMP signaling, inducer exclusion, "
     "catabolite repression. This IS Hayek's knowledge problem: the relevant information "
     "exists only in local context, not in any central database."),

    ("L2_fba_pert_anaerobic", "FBA Perturbation: Aerobic to Anaerobic (Mises)",
     ["FBA re-solves the LP for anaerobic conditions",
      "Real E. coli uses ArcAB and FNR regulatory cascades",
      "The planner computes the answer; biology GROWS toward it",
      "Growth ratio: ~18% — fermentation is far less efficient"],
     "When oxygen is removed, FBA re-optimizes and finds the best anaerobic flux distribution, "
     "routing carbon through fermentation pathways. Real E. coli uses the ArcAB two-component "
     "system and FNR regulon to sense oxygen levels and gradually shift gene expression. "
     "The planner computes the answer in one step; the cell's distributed regulatory network "
     "grows toward the new equilibrium through iterative local feedback. This tests Mises' "
     "calculation problem: can a central authority compute the optimal allocation as efficiently "
     "as a distributed system? FBA finds the optimum but misses the transition dynamics."),

    ("L2_fba_pert_nitrogen", "FBA Perturbation: Nitrogen Limitation (Kirzner)",
     ["FBA predicts the new optimum under nitrogen scarcity",
      "Real E. coli upregulates high-affinity transporters via NtrBC",
      "Entrepreneurial alertness: discovering new routes to scarce resources",
      "Growth ratio: ~11% — nitrogen is critical"],
     "Under severe nitrogen limitation (NH4 uptake reduced to -1 mmol/gDW/h), FBA predicts "
     "the optimal growth rate by redistributing flux. Real E. coli responds through the NtrBC "
     "two-component system, which senses nitrogen scarcity and upregulates high-affinity "
     "ammonium transporters and alternative nitrogen sources (amino acid catabolism, nitrogen "
     "scavenging). This is Kirzner's entrepreneurial alertness at the molecular level: the "
     "cell doesn't just optimize within existing pathways — it discovers and activates new "
     "routes that weren't being used before. FBA can't predict which alternative pathways "
     "will be activated because that depends on regulatory context."),

    # --- Layer 3: Trade ---
    ("L3_trade_network_graph", "Cross-Species Trade Network",
     ["5-7 organisms connected by trade links (heterologous expression potential)",
      "Edge weight = trade ease (inverse of codon distance + regulatory barrier)",
      "Node color = kingdom (prokaryote, fungal, plant, animal)",
      "Closely related species trade more cheaply — Ricardian comparative advantage"],
     "The cross-species trade network models gene transfer potential between organisms. Each "
     "edge is weighted by 'trade cost' — the sum of codon usage distance (RSCU-based), "
     "regulatory barrier (can the host process the gene's signals?), and compartment "
     "compatibility. Low trade cost = easy heterologous expression. High trade cost = the gene "
     "product will misfold, be mistranslated, or lack necessary chaperones. This maps directly "
     "to Ricardo's comparative advantage: organisms specialize in what they produce most "
     "efficiently and 'trade' genetic capabilities where the cost is lowest."),

    ("L3_trade_cost_heatmap", "Trade Barriers Heatmap",
     ["Each cell = trade cost between two organisms",
      "Low cost (yellow) = compatible expression systems",
      "High cost (red) = incompatible codon usage, missing chaperones",
      "Diagonal = zero (self-trade is free)",
      "Prokaryote-to-eukaryote has highest barriers"],
     "The trade cost matrix quantifies the difficulty of expressing one organism's genes in "
     "another. It combines codon usage distance (RSCU), regulatory compatibility, and "
     "compartment barriers. The matrix is symmetric — the cost of E. coli expressing a yeast "
     "gene equals the cost of yeast expressing an E. coli gene (in this simplified model). "
     "Prokaryote-to-eukaryote transfers have the highest barriers: different ribosome structure, "
     "no intron processing, different post-translational modifications. Within-kingdom transfers "
     "are cheaper. This is the biological equivalent of trade barriers: tariffs, regulations, "
     "and infrastructure differences that increase the cost of exchange."),

    ("L3_comparative_advantage_table", "Comparative Advantage — Who Exports What",
     ["Each organism specializes in specific molecular capabilities",
      "E. coli: protein expression, rapid growth, simple genetics",
      "Yeast: post-translational modification, secretion, eukaryotic folding",
      "Human cells: complex signaling, tissue-specific expression",
      "Specialization IS division of labor (Menger's spontaneous order)"],
     "Each organism has evolved to specialize in specific molecular capabilities — not because "
     "a planner assigned roles, but because evolutionary pressure selected for efficiency in "
     "particular niches. E. coli is the 'factory worker' of molecular biology: fast growth, "
     "cheap protein production, simple genetics. Yeast handles eukaryotic complexity: "
     "glycosylation, disulfide bonds, secretion. Human cells do what nothing else can: "
     "tissue-specific expression, complex signaling cascades, mammalian post-translational "
     "modifications. This division of labor emerged spontaneously through billions of years "
     "of evolution — Menger's spontaneous order at the species level."),
]


def build_pptx():
    """Build PowerPoint from individual PNGs and slide definitions."""
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Colors
    bg_color = RGBColor(0x0d, 0x11, 0x17)
    text_color = RGBColor(0xe6, 0xed, 0xf3)
    accent_color = RGBColor(0x52, 0xb7, 0x88)
    gold_color = RGBColor(0xe9, 0xc4, 0x6a)
    dim_color = RGBColor(0x8b, 0x94, 0x9e)

    for fig_name, title, bullets, notes in SLIDE_DEFS:
        slide_layout = prs.slide_layouts[6]  # blank
        slide = prs.slides.add_slide(slide_layout)

        # Set background
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = bg_color

        if fig_name is None:
            # Title slide
            # Title
            txBox = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(11), Inches(1.5))
            tf = txBox.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = title
            p.font.size = Pt(36)
            p.font.color.rgb = text_color
            p.font.bold = True
            p.alignment = PP_ALIGN.CENTER

            # Subtitle bullets
            for i, bullet in enumerate(bullets):
                txBox2 = slide.shapes.add_textbox(Inches(1), Inches(3.8 + i * 0.5), Inches(11), Inches(0.5))
                tf2 = txBox2.text_frame
                tf2.word_wrap = True
                p2 = tf2.paragraphs[0]
                p2.text = bullet
                p2.font.size = Pt(18)
                p2.font.color.rgb = accent_color
                p2.alignment = PP_ALIGN.CENTER
        else:
            # Content slide: figure on left, bullets on right
            img_path = FIG_DIR / f"{fig_name}.png"

            # Title bar at top
            txBox = slide.shapes.add_textbox(Inches(0.3), Inches(0.2), Inches(12.7), Inches(0.6))
            tf = txBox.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = title
            p.font.size = Pt(22)
            p.font.color.rgb = accent_color
            p.font.bold = True

            # Figure on left (60% width)
            if img_path.exists():
                slide.shapes.add_picture(
                    str(img_path), Inches(0.3), Inches(1.0), Inches(7.8), Inches(5.8)
                )
            else:
                # Placeholder box
                txBox_ph = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(7.5), Inches(2))
                tf_ph = txBox_ph.text_frame
                tf_ph.word_wrap = True
                p_ph = tf_ph.paragraphs[0]
                p_ph.text = f"[Figure: {fig_name}.png]\nRun: python run_all.py to generate"
                p_ph.font.size = Pt(14)
                p_ph.font.color.rgb = dim_color
                p_ph.alignment = PP_ALIGN.CENTER

            # Bullet points on right (40% width)
            bullet_top = Inches(1.2)
            txBox_b = slide.shapes.add_textbox(Inches(8.3), bullet_top, Inches(4.7), Inches(5.5))
            tf_b = txBox_b.text_frame
            tf_b.word_wrap = True

            for i, bullet in enumerate(bullets):
                if i == 0:
                    p_b = tf_b.paragraphs[0]
                else:
                    p_b = tf_b.add_paragraph()

                p_b.text = f"\u2022  {bullet}"
                p_b.font.size = Pt(13)
                p_b.font.color.rgb = text_color
                p_b.space_after = Pt(8)

        # Speaker notes
        notes_slide = slide.notes_slide
        notes_tf = notes_slide.notes_text_frame
        notes_tf.text = notes

    # Save
    out_path = PROJECT_ROOT / "paper" / "capstone_presentation.pptx"
    prs.save(str(out_path))
    print(f"\nPowerPoint saved to: {out_path}")
    return out_path


# ===================================================================
# MAIN
# ===================================================================

def main():
    parser = argparse.ArgumentParser(description="Generate presentation figures and PPTX")
    parser.add_argument("--quick", action="store_true", help="Quick mode (skip API calls)")
    parser.add_argument("--pptx-only", action="store_true", help="Skip figure generation, just build PPTX")
    args = parser.parse_args()

    print("=" * 60)
    print("  GENERATING INDIVIDUAL PANEL FIGURES + POWERPOINT")
    print("=" * 60)

    if not args.pptx_only:
        FIG_DIR.mkdir(parents=True, exist_ok=True)

        # --- Layer 1: Topology ---
        print("\n[Layer 1] Topology panels...")
        try:
            from layer1_topology.network_fetcher import load_all_networks, _builtin_ecoli_grn
            from layer1_topology.topology_analysis import analyze_topology
            from layer1_topology.centralized_comparison import build_comparison_networks

            if args.quick:
                networks = {"ecoli_grn": _builtin_ecoli_grn()}
            else:
                networks = load_all_networks()

            bio_reports = []
            for name, G in networks.items():
                report = analyze_topology(G, name=name)
                bio_reports.append(report)

            ref_n = max(G.number_of_nodes() for G in networks.values())
            ref_m = max(G.number_of_edges() for G in networks.values())
            ref_networks = build_comparison_networks(n_nodes=ref_n, n_edges=ref_m)
            ref_reports = []
            for name, G in ref_networks.items():
                report = analyze_topology(G, name=name)
                ref_reports.append(report)

            gen_layer1_panels(bio_reports, ref_reports)
        except Exception as e:
            print(f"  [SKIP] Layer 1 topology: {e}")

        # --- Layer 1: Motifs ---
        print("\n[Layer 1] Motif panels...")
        try:
            from layer1_topology.network_fetcher import _builtin_ecoli_grn, fetch_regulondb_grn
            from layer1_topology.motif_analysis import run_motif_analysis

            G = _builtin_ecoli_grn() if args.quick else fetch_regulondb_grn()
            n_random = 100 if args.quick else 1000
            motif_report = run_motif_analysis(G, n_random=n_random)
            gen_motif_panels(motif_report)
        except Exception as e:
            print(f"  [SKIP] Layer 1 motifs: {e}")

        # --- Layer 1: Self-Regulation ---
        print("\n[Layer 1] Self-regulation panels...")
        try:
            from layer1_topology.self_regulation import compute_self_regulation

            sr_bio = []
            for name, G in networks.items():
                sr = compute_self_regulation(G, name=name)
                sr_bio.append(sr)

            sr_ref = []
            for name, G in ref_networks.items():
                sr = compute_self_regulation(G, name=name)
                sr_ref.append(sr)

            gen_self_regulation_panels(sr_bio, sr_ref)
        except Exception as e:
            print(f"  [SKIP] Layer 1 self-regulation: {e}")

        # --- Layer 2: Economy ---
        print("\n[Layer 2] Economy panels...")
        try:
            from layer2_economy.metabolic_economy import (
                run_distributed, run_centralized, run_perturbation_test,
            )

            base_supply = {
                "UDP-glucose": 1.0, "UDP-GlcNAc": 0.5, "Ca2+": 0.5,
                "O2": 2.0, "luciferin": 0.3, "glycine": 1.0,
                "alanine": 1.0, "glutamine": 0.5, "piRNA_precursors": 0.2,
                "Zn2+": 0.2, "cholesterol": 0.3,
            }
            dist = run_distributed(n_steps=200, base_supply=base_supply)
            cent = run_centralized(n_steps=200, base_supply=base_supply)
            robustness = run_perturbation_test(n_steps=100)
            gen_economy_panels(dist, cent, robustness)
        except Exception as e:
            print(f"  [SKIP] Layer 2 economy: {e}")

        # --- Layer 2: FBA ---
        print("\n[Layer 2] FBA panels...")
        try:
            from layer2_economy.fba_analysis import run_fba_analysis, run_fba_perturbation
            fba_report = run_fba_analysis()
            gen_fba_panels(fba_report)

            fba_pert = run_fba_perturbation()
            gen_fba_perturbation_panels(fba_pert)
        except Exception as e:
            print(f"  [SKIP] Layer 2 FBA: {e}")

        # --- Layer 3: Trade ---
        print("\n[Layer 3] Trade panels...")
        try:
            gen_trade_panels()
        except Exception as e:
            print(f"  [SKIP] Layer 3 trade: {e}")

    # --- Build PowerPoint ---
    print("\n[PPTX] Building PowerPoint...")
    pptx_path = build_pptx()

    print("\n" + "=" * 60)
    print("  COMPLETE")
    print(f"  Individual figures: {FIG_DIR}/")
    print(f"  PowerPoint: {pptx_path}")
    print("=" * 60)


if __name__ == "__main__":
    main()
