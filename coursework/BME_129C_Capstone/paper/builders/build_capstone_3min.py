"""
Build 3-minute capstone presentation for BME 129C.
Based on capstone_summary_3page.md — the core argument in 7 slides.

~25-30 seconds per slide:
  1. Title + thesis (15 sec)
  2. Layer 1: No Master Node — 19:1 robustness (30 sec)
  3. Layer 2: Distributed beats centralized — 71% vs 53% (30 sec)
  4. The Price System of the Cell — three tiers (30 sec)
  5. Cancer breaks the price system — TCGA data (30 sec)
  6. Layer 3: Trade + cells as agents (30 sec)
  7. Throughline — one conclusion (15 sec)

Usage: python paper/builders/build_capstone_3min.py
"""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

PAPER_DIR = Path(__file__).resolve().parent.parent
FIGURES = PAPER_DIR / "figures"
INDIVIDUAL = FIGURES / "individual"
OUT = PAPER_DIR / "presentations" / "capstone_3min.pptx"

# Colors — matching capstone dark theme
DARK_BG = RGBColor(0x0d, 0x11, 0x17)
GREEN_MID = RGBColor(0x52, 0xb7, 0x88)
GREEN_LIGHT = RGBColor(0x95, 0xd5, 0xb2)
GREEN_DARK = RGBColor(0x1b, 0x43, 0x32)
GOLD = RGBColor(0xe9, 0xc4, 0x6a)
RED = RGBColor(0xe6, 0x39, 0x46)
WHITE = RGBColor(0xff, 0xff, 0xff)
OFF_WHITE = RGBColor(0xf5, 0xf5, 0xf5)
GRAY = RGBColor(0x8b, 0x94, 0x9e)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)


def set_bg(slide):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = DARK_BG


def add_green_bar(slide):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_W, Inches(0.06))
    shape.fill.solid()
    shape.fill.fore_color.rgb = GREEN_MID
    shape.line.fill.background()


def add_text(slide, left, top, width, height, text, size=18,
             color=WHITE, bold=False, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = "Calibri"
    p.alignment = align
    return tf


def add_bullets(slide, left, top, width, height, items, size=16, color=OFF_WHITE):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"\u2022  {item}"
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.name = "Calibri"
        p.space_after = Pt(8)
    return tf


def try_pic(slide, *candidates, left=0, top=0, width=None, height=None):
    for c in candidates:
        p = Path(c)
        if p.exists():
            slide.shapes.add_picture(str(p), left, top, width, height)
            return True
    return False


def add_notes(slide, text):
    notes = slide.notes_slide
    notes.notes_text_frame.text = text


def build():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # =========================================================================
    # SLIDE 1: Title (15 sec)
    # =========================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_green_bar(slide)

    add_text(slide, Inches(1), Inches(1.0), Inches(11), Inches(1.2),
             "Living Systems as\nDecentralized Economies",
             size=44, color=GREEN_MID, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, Inches(1), Inches(2.8), Inches(11), Inches(0.8),
             "The Price System of the Cell",
             size=26, color=GOLD, align=PP_ALIGN.CENTER)
    add_text(slide, Inches(1), Inches(3.9), Inches(11), Inches(0.5),
             "Sage Clokey  |  BME 129C Capstone  |  UC Santa Cruz  |  Spring 2026",
             size=16, color=GRAY, align=PP_ALIGN.CENTER)

    add_text(slide, Inches(1.5), Inches(5.2), Inches(10), Inches(1.2),
             "Living things are not machines. The differences are a feature, not a bug.\n"
             "The distributed knowledge is coordinated by prices \u2014\n"
             "the ratio between the voluntary exchange of anything between nodes.",
             size=17, color=OFF_WHITE, align=PP_ALIGN.CENTER)

    add_notes(slide,
        "Title. 15 seconds.\n\n"
        "Living systems coordinate billions of nodes without a central planner. "
        "No master gene runs your body. No forester runs a forest. This capstone asks: "
        "are living systems literally decentralized economies? Not as metaphor \u2014 as "
        "measurable, quantifiable, structurally equivalent systems."
    )

    # =========================================================================
    # SLIDE 2: Layer 1 — No Master Node (30 sec)
    # =========================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_green_bar(slide)

    add_text(slide, Inches(0.5), Inches(0.2), Inches(12), Inches(0.6),
             "No Master Node \u2014 19:1 Robustness Advantage",
             size=28, color=GOLD, bold=True)

    try_pic(slide,
            INDIVIDUAL / "L1_degree_distribution.png",
            FIGURES / "layer1_topology.png",
            left=Inches(0.3), top=Inches(1.0),
            width=Inches(6.0), height=Inches(2.8))

    try_pic(slide,
            INDIVIDUAL / "L1_robustness_curves.png",
            FIGURES / "layer1_topology_annotated.png",
            left=Inches(0.3), top=Inches(4.0),
            width=Inches(6.0), height=Inches(2.8))

    add_bullets(slide, Inches(6.8), Inches(1.0), Inches(6.2), Inches(6.0), [
        "E. coli GRN: power-law degree distribution (\u03b1 2.0\u20132.5)",
        "Hubs exist (CRP \u2192 43 genes) but no hub dominates",
        "Targeted removal: biological networks survive 37% node loss",
        "Star graph (central planning) collapses at 1.9%",
        "19:1 robustness ratio \u2014 the structural cost of centralization",
        "Hub erosion: network actively resists centralization",
        "Feed-forward loops over-represented \u2014 evolved price signals",
    ], size=15)

    add_notes(slide,
        "Layer 1 \u2014 30 seconds.\n\n"
        "The E. coli gene regulatory network has no central controller. "
        "Under targeted removal of top nodes, biological networks survive "
        "37% loss. The star graph collapses at 2%. That's a 19:1 ratio. "
        "The distributed network routes around damage. The centralized "
        "network has a single point of failure."
    )

    # =========================================================================
    # SLIDE 3: Layer 2 — Distributed Beats Centralized (30 sec)
    # =========================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_green_bar(slide)

    add_text(slide, Inches(0.5), Inches(0.2), Inches(12), Inches(0.6),
             "The Market Wins Under Stress \u2014 71% vs 53%",
             size=28, color=GOLD, bold=True)

    try_pic(slide,
            INDIVIDUAL / "L2_gdp_over_time.png",
            FIGURES / "layer2_economy.png",
            left=Inches(0.3), top=Inches(1.0),
            width=Inches(6.0), height=Inches(2.8))

    try_pic(slide,
            INDIVIDUAL / "L2_perturbation_robustness.png",
            FIGURES / "layer2_fba_analysis.png",
            left=Inches(0.3), top=Inches(4.0),
            width=Inches(6.0), height=Inches(2.8))

    add_bullets(slide, Inches(6.8), Inches(1.0), Inches(6.2), Inches(6.0), [
        "13 metabolic agents share a pool \u2014 no central allocator",
        "Rates oscillate early (price discovery) then converge",
        "Under perturbation: distributed retains 71% GDP",
        "Centralized retains 53% \u2014 18-point market advantage",
        "The planner\u2019s fixed plan is wrong after the shock",
        "Distributed agents self-correct through local price signals",
        "FBA (omniscient planner): 70% accuracy \u2014 30% structural failure",
    ], size=15)

    add_notes(slide,
        "Layer 2 \u2014 30 seconds.\n\n"
        "Thirteen metabolic pathway agents reach equilibrium through local "
        "feedback alone. Under perturbation, distributed allocation retains "
        "71% of GDP. Centralized retains 53%. The planner's fixed plan is "
        "wrong after the shock, and it has no mechanism to discover the new "
        "optimum. Even FBA with perfect stoichiometric knowledge fails on "
        "30% of real gene knockouts."
    )

    # =========================================================================
    # SLIDE 4: Price System of the Cell (30 sec)
    # =========================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_green_bar(slide)

    add_text(slide, Inches(0.5), Inches(0.2), Inches(12), Inches(0.6),
             "The Price System of the Cell",
             size=28, color=GOLD, bold=True)

    try_pic(slide,
            INDIVIDUAL / "panel_a_price_tiers.png",
            FIGURES / "price_system_of_the_cell.png",
            left=Inches(0.3), top=Inches(1.0),
            width=Inches(6.2), height=Inches(2.8))

    try_pic(slide,
            INDIVIDUAL / "panel_b_shadow_prices.png",
            left=Inches(0.3), top=Inches(4.0),
            width=Inches(6.2), height=Inches(2.8))

    add_bullets(slide, Inches(6.8), Inches(1.0), Inches(6.2), Inches(6.0), [
        "Three tiers of molecular prices:",
        "  Intracellular ratios (ATP/ADP, NAD+/NADH) = cost of capital",
        "  Intercellular signals (cytokines, O\u2082 tension) = market prices",
        "  mTOR integrates all prices \u2192 grow-or-conserve = the entrepreneur",
        "",
        "Shadow prices shift with context (FBA, iML1515):",
        "  NADH: cheap on glucose, expensive on acetate",
        "  Oxygen: free aerobically, most valuable anaerobically",
        "Same molecule, different value = Menger\u2019s subjective value",
    ], size=14)

    add_notes(slide,
        "Price System \u2014 30 seconds.\n\n"
        "Living systems have a real price system at three tiers. Intracellular "
        "metabolite ratios are cost of capital. Intercellular signals are market "
        "prices. mTOR integrates all of these and makes the grow-or-conserve "
        "decision \u2014 the entrepreneur reading the market. Shadow prices from "
        "FBA prove Menger's subjective value: the same molecule has different "
        "value depending on context."
    )

    # =========================================================================
    # SLIDE 5: Cancer Breaks the Price System (30 sec)
    # =========================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_green_bar(slide)

    add_text(slide, Inches(0.5), Inches(0.2), Inches(12), Inches(0.6),
             "Cancer Breaks the Price System",
             size=28, color=RED, bold=True)

    try_pic(slide,
            INDIVIDUAL / "panel_d_cancer_price_system.png",
            FIGURES / "price_system_of_the_cell.png",
            left=Inches(0.3), top=Inches(1.0),
            width=Inches(6.2), height=Inches(5.8))

    add_bullets(slide, Inches(6.8), Inches(1.0), Inches(6.2), Inches(6.0), [
        "TCGA PanCancerAtlas: 10,967 samples, 32 cancer types",
        "Most mutated genes map to price system components:",
        "  TP53 \u2014 damage price (apoptosis signal)",
        "  PIK3CA, PTEN, mTOR \u2014 price integrators",
        "  EGFR, ERBB2 \u2014 growth factor receptors (price readers)",
        "",
        "The variation IS the finding:",
        "  TP53: 96% ovarian, 1% thyroid",
        "  PIK3CA: 52% uterine, 3% ovarian",
        "Different tissues break different price components",
        "",
        "The disease is not the variant \u2014 it\u2019s the broken price system",
    ], size=14)

    add_text(slide, Inches(0.5), Inches(6.8), Inches(12), Inches(0.5),
             "Cancer is the calculation problem at the cellular level.",
             size=16, color=GREEN_MID, bold=True, align=PP_ALIGN.CENTER)

    add_notes(slide,
        "Cancer \u2014 30 seconds.\n\n"
        "Using TCGA data \u2014 10,967 samples, 32 cancer types \u2014 the most "
        "mutated genes all map to price system components. TP53 is the "
        "damage price. PIK3CA and PTEN are price integrators. EGFR is "
        "a price reader. The tissue-specific variation is the finding: "
        "TP53 is 96% in ovarian but 1% in thyroid. Different tissues "
        "break different price components. Cancer is not a gene disease. "
        "It's a price system disease."
    )

    # =========================================================================
    # SLIDE 6: Layer 3 — Trade + Cells as Agents (30 sec)
    # =========================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_green_bar(slide)

    add_text(slide, Inches(0.5), Inches(0.2), Inches(12), Inches(0.6),
             "Comparative Advantage Across the Tree of Life",
             size=28, color=GOLD, bold=True)

    try_pic(slide,
            INDIVIDUAL / "L3_trade_network_graph.png",
            FIGURES / "layer3_trade_network.png",
            left=Inches(0.3), top=Inches(1.0),
            width=Inches(6.0), height=Inches(2.8))

    try_pic(slide,
            INDIVIDUAL / "L3_trade_cost_heatmap.png",
            FIGURES / "layer3_trade_network_annotated.png",
            left=Inches(0.3), top=Inches(4.0),
            width=Inches(6.0), height=Inches(2.8))

    add_bullets(slide, Inches(6.8), Inches(1.0), Inches(6.2), Inches(3.0), [
        "Each organism specializes: coral \u2192 biomineralization, spider \u2192 silk",
        "Trade cost scales with evolutionary distance:",
        "  Within-kingdom: 0.17\u20130.38 | Cross-kingdom: 0.65\u20130.83",
        "Forced exchange (full codon optimization) destroys local knowledge",
        "Trade blocs emerge spontaneously \u2014 no designer needed",
    ], size=15)

    # Lower half — cells as agents
    add_text(slide, Inches(6.8), Inches(4.2), Inches(6.2), Inches(0.4),
             "Cells as Economic Agents",
             size=18, color=GREEN_LIGHT, bold=True)

    add_bullets(slide, Inches(6.8), Inches(4.7), Inches(6.2), Inches(2.5), [
        "Same genome, 8 cell types \u2014 specialization by choice, not command",
        "Communication Gini = 0.0 \u2014 no gatekeeper, perfectly distributed",
        "Remove any cell type \u2192 70\u201390% of communication intact",
        "Same ligand, different meaning \u2014 subjective value at the molecular level",
    ], size=14, color=GREEN_LIGHT)

    add_notes(slide,
        "Trade + Cells \u2014 30 seconds.\n\n"
        "Each organism specializes in capabilities others lack. Trade cost "
        "correlates with evolutionary distance. Forced exchange destroys "
        "information. Trade blocs emerge spontaneously. At the cellular level, "
        "same genome produces 8 cell types. Communication Gini is 0.0 \u2014 "
        "perfectly distributed, no gatekeeper. Same ligand received differently "
        "by different cells \u2014 subjective value at the molecular level."
    )

    # =========================================================================
    # SLIDE 7: Throughline (15 sec)
    # =========================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_green_bar(slide)

    add_text(slide, Inches(1), Inches(0.5), Inches(11), Inches(0.8),
             "Every Layer, One Answer",
             size=36, color=GREEN_MID, bold=True, align=PP_ALIGN.CENTER)

    # Summary rows
    rows = [
        ("Topology",    "No master node \u2014 19:1 robustness ratio"),
        ("Single-Cell", "No master cell \u2014 Gini = 0.0, subjective value"),
        ("Economy",     "Distributed 71% vs centralized 53% under stress"),
        ("FBA",         "Omniscient planner: 70% \u2014 30% structural failure"),
        ("Price System","Three-tier prices \u2014 mTOR is the entrepreneur"),
        ("Cancer",      "Disease = broken price system, not broken parts"),
        ("Trade",       "Voluntary exchange succeeds, forced exchange destroys"),
    ]
    for i, (layer, finding) in enumerate(rows):
        y = Inches(1.6) + Inches(0.55) * i
        add_text(slide, Inches(1.5), y, Inches(2.8), Inches(0.5),
                 layer, size=17, color=GOLD, bold=True)
        add_text(slide, Inches(4.5), y, Inches(7.5), Inches(0.5),
                 finding, size=16, color=OFF_WHITE)

    # Closing box
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(2), Inches(5.8), Inches(9.3), Inches(1.3))
    shape.fill.solid()
    shape.fill.fore_color.rgb = GREEN_DARK
    shape.line.color.rgb = GREEN_MID
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = (
        "Life is the original decentralized network. The self-assembling internet. "
        "The adaptive computer that built itself, runs itself, repairs itself, "
        "and has been doing so for four billion years."
    )
    p.font.size = Pt(15)
    p.font.color.rgb = GREEN_LIGHT
    p.font.name = "Calibri"
    p.alignment = PP_ALIGN.CENTER

    p2 = tf.add_paragraph()
    p2.text = "The question is not how to program life. It is how to join the computation."
    p2.font.size = Pt(15)
    p2.font.color.rgb = GOLD
    p2.font.bold = True
    p2.font.name = "Calibri"
    p2.alignment = PP_ALIGN.CENTER
    p2.space_before = Pt(6)

    add_notes(slide,
        "Conclusion \u2014 15 seconds.\n\n"
        "Every layer answers the same question: does biology operate like "
        "a centrally planned economy or a free market? The answer is market. "
        "At every scale. Life is the original decentralized network. The question "
        "is not how to program life. It is how to join the computation."
    )

    prs.save(str(OUT))
    print(f"Saved: {OUT}")
    print(f"  {len(prs.slides)} slides")


if __name__ == "__main__":
    build()
