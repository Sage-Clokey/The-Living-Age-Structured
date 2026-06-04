"""
Generate 3-Minute Progress Report Slides with Speaker Notes.
BME 129C — Sage Clokey — Spring 2026

Structure (per assignment):
  - Goal & approach (0.5 min) — 1 slide
  - Accomplishments to date (1 min) — 2 slides
  - Progress this week (1.5 min) — 2-3 slides

Usage: python paper/generate_progress_slides.py
"""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

PROJECT_ROOT = Path(__file__).resolve().parent.parent.parent
FIG_DIR = PROJECT_ROOT / "paper" / "figures" / "individual"

# Colors
BG = RGBColor(0x0d, 0x11, 0x17)
TEXT = RGBColor(0xe6, 0xed, 0xf3)
GREEN = RGBColor(0x52, 0xb7, 0x88)
GOLD = RGBColor(0xe9, 0xc4, 0x6a)
DIM = RGBColor(0x8b, 0x94, 0x9e)
RED = RGBColor(0xe6, 0x39, 0x46)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)


def add_slide(title, bullets, notes, fig_name=None, fig_right=False):
    """Add a slide with dark background, title, bullets, optional figure, and speaker notes."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

    # Background
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG

    # Title
    tx = slide.shapes.add_textbox(Inches(0.4), Inches(0.25), Inches(12.5), Inches(0.7))
    tf = tx.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(26)
    p.font.color.rgb = GREEN
    p.font.bold = True

    if fig_name and fig_right:
        # Layout: bullets left, figure right
        bullet_left = Inches(0.4)
        bullet_width = Inches(6.5)
        bullet_top = Inches(1.2)
        bullet_height = Inches(5.5)

        img_path = FIG_DIR / f"{fig_name}.png"
        if img_path.exists():
            slide.shapes.add_picture(str(img_path), Inches(7.2), Inches(1.0), Inches(5.8), Inches(5.5))
        else:
            # Try combined figures dir
            img_path2 = PROJECT_ROOT / "paper" / "figures" / f"{fig_name}.png"
            if img_path2.exists():
                slide.shapes.add_picture(str(img_path2), Inches(7.2), Inches(1.0), Inches(5.8), Inches(5.5))
    elif fig_name:
        # Layout: bullets top, figure bottom or figure only
        bullet_left = Inches(0.4)
        bullet_width = Inches(12.5)
        bullet_top = Inches(1.2)
        bullet_height = Inches(2.5)

        img_path = FIG_DIR / f"{fig_name}.png"
        if not img_path.exists():
            img_path = PROJECT_ROOT / "paper" / "figures" / f"{fig_name}.png"
        if img_path.exists():
            slide.shapes.add_picture(str(img_path), Inches(1.5), Inches(3.8), Inches(10), Inches(3.5))
    else:
        # No figure — full-width bullets
        bullet_left = Inches(0.4)
        bullet_width = Inches(12.5)
        bullet_top = Inches(1.2)
        bullet_height = Inches(5.5)

    # Bullets
    bx = slide.shapes.add_textbox(bullet_left, bullet_top, bullet_width, bullet_height)
    bf = bx.text_frame
    bf.word_wrap = True

    for i, bullet in enumerate(bullets):
        if i == 0:
            bp = bf.paragraphs[0]
        else:
            bp = bf.add_paragraph()

        # Check for sub-bullets (indented with "  -")
        if bullet.startswith("  -"):
            bp.text = f"    {bullet.strip('- ').strip()}"
            bp.font.size = Pt(14)
            bp.font.color.rgb = DIM
        elif bullet.startswith(">>"):
            # Highlight line
            bp.text = f"  {bullet.strip('> ').strip()}"
            bp.font.size = Pt(16)
            bp.font.color.rgb = GOLD
            bp.font.bold = True
        else:
            bp.text = f"\u2022  {bullet}"
            bp.font.size = Pt(16)
            bp.font.color.rgb = TEXT

        bp.space_after = Pt(6)

    # Speaker notes
    notes_slide = slide.notes_slide
    notes_tf = notes_slide.notes_text_frame
    notes_tf.text = notes


# ===================================================================
# SLIDE 1: Title
# ===================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
fill = slide.background.fill
fill.solid()
fill.fore_color.rgb = BG

tx = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(11), Inches(1.5))
tf = tx.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Living Systems as Decentralized Economies"
p.font.size = Pt(36)
p.font.color.rgb = TEXT
p.font.bold = True
p.alignment = PP_ALIGN.CENTER

tx2 = slide.shapes.add_textbox(Inches(1), Inches(3.2), Inches(11), Inches(0.6))
tf2 = tx2.text_frame
tf2.word_wrap = True
p2 = tf2.paragraphs[0]
p2.text = "Progress Report — BME 129C Capstone"
p2.font.size = Pt(20)
p2.font.color.rgb = GREEN
p2.alignment = PP_ALIGN.CENTER

tx3 = slide.shapes.add_textbox(Inches(1), Inches(4.0), Inches(11), Inches(0.5))
tf3 = tx3.text_frame
tf3.word_wrap = True
p3 = tf3.paragraphs[0]
p3.text = "Sage Clokey  |  Spring 2026  |  Advisor: R. Dubois"
p3.font.size = Pt(16)
p3.font.color.rgb = DIM
p3.alignment = PP_ALIGN.CENTER

notes_slide = slide.notes_slide
notes_slide.notes_text_frame.text = (
    "Title slide. This is a solo project — no group members.\n\n"
    "The central question: are living systems literally decentralized economies? "
    "Not as metaphor — as measurable, quantifiable, structurally equivalent systems. "
    "I'm using biological data at three scales to test whether Austrian economic "
    "principles (Hayek, Mises, Menger, Rothbard, Kirzner) describe molecular biology.\n\n"
    "The title comes from Hayek's The Fatal Conceit (1988): 'The curious task of "
    "economics is to demonstrate to men how little they really know about what they "
    "imagine they can design.' This project applies that insight to the most ancient "
    "economy: the cell."
)

# ===================================================================
# SLIDE 2: Goal & Approach (~30 seconds)
# ===================================================================
add_slide(
    "Goal & Approach",
    [
        "Thesis: living systems operate as decentralized economies",
        "  - Not metaphor — measurable, testable",
        "Hayek's Fatal Conceit (1988): you can't",
        "  - redesign what evolved spontaneously",
        "Three layers of evidence:",
        "  - L1: Network topology",
        "  - L2: Economic modeling + FBA",
        "  - L3: Cross-species trade",
        ">> Does distributed coordination outperform central planning?",
    ],
    notes="GOAL (30 seconds):\n\n"
    "My thesis is that living systems — from gene regulatory networks inside a single "
    "bacterium to cross-species symbioses — operate as decentralized economies. Not as "
    "a metaphor, but as measurable, structurally equivalent systems.\n\n"
    "I'm grounding this in Austrian economics — Hayek's knowledge problem, Mises' "
    "calculation problem, Menger's spontaneous order. Hayek's final book, The Fatal "
    "Conceit (1988), names the error directly: the belief that human reason can redesign "
    "evolved order. His 'curious task of economics' line — showing men how little they "
    "know about what they imagine they can design — is the thesis of this project applied "
    "to molecular biology.\n\n"
    "Layer 1 looks at network structure. Layer 2 models metabolic pathways as economic "
    "agents. Layer 3 maps gene transferability as a trade network.\n\n"
    "The core question: does the decentralized architecture outperform centralized "
    "alternatives — and is this testable with real data?",
    fig_name="L1_robustness_curves",
    fig_right=True,
)

# ===================================================================
# SLIDE 3: Accomplishments — Layers 1 & 1b (~30 seconds)
# ===================================================================
add_slide(
    "Accomplished: Network Topology & Single-Cell Economy",
    [
        "Layer 1 — No Master Node:",
        "  - 5 bio networks vs 5 reference architectures",
        "  - 19:1 robustness ratio under targeted attack",
        "  - Feed-forward loops = Hayekian price signals",
        "  - WBPA: biology resists centralization",
        "",
        "Layer 1b — Cells as Economic Agents:",
        "  - 2,638 PBMCs, 8 cell types (scRNA-seq)",
        "  - Communication Gini = 0.0 — no gatekeeper",
        "  - 75% communication survives any removal",
    ],
    notes="ACCOMPLISHMENTS — LAYERS 1 & 1b (30 seconds):\n\n"
    "Layer 1 is fully complete. I pulled real gene regulatory networks from RegulonDB, "
    "metabolic networks from KEGG, and protein interaction networks from STRING. I compared "
    "them against five reference architectures — star graph, random, lattice, hub-and-spoke, "
    "and Barabasi-Albert.\n\n"
    "Key finding: biological networks show a 19-to-1 robustness advantage over centralized "
    "architectures under targeted attack. The star graph collapses after removing 2% of nodes; "
    "biological PPI networks survive removing 37%.\n\n"
    "I also did motif analysis following Alon 2002 — the feed-forward loop is massively "
    "over-represented, which I interpret as a Hayekian price signal motif. And self-regulation "
    "analysis following Topirceanu 2018 shows biological networks actively resist centralization "
    "through a betweenness-preferential attachment mechanism.\n\n"
    "Layer 1b analyzed single-cell RNA-seq data from human PBMCs — 2,638 cells, 8 cell types. "
    "The betweenness Gini of the communication network is 0.0 — perfectly distributed, no "
    "master cell. And removing any single cell type still leaves 75% of communication intact.",
    fig_name="L1_self_reg_spearman",
    fig_right=True,
)

# ===================================================================
# SLIDE 4: Accomplishments — Layers 2 & 3 (~30 seconds)
# ===================================================================
add_slide(
    "Accomplished: Economic Modeling & Trade Network",
    [
        "Layer 2 — Pathways as Agents:",
        "  - 13 pathways with ATP budgets",
        "  - Distributed: 2.9x lower GDP variance",
        "  - 4 perturbation tests (Hayek, Mises, Kirzner)",
        "",
        "Layer 3 — Cross-Species Trade:",
        "  - Codon usage distance = trade friction",
        "  - 18 published transfers validated",
        "  - Trade blocs emerge spontaneously",
    ],
    notes="ACCOMPLISHMENTS — LAYERS 2 & 3 (30 seconds):\n\n"
    "Layer 2 models 13 metabolic pathways from my genome design system as economic agents. "
    "Each has inputs, outputs, and an ATP budget. I simulate distributed allocation — where "
    "each agent adjusts based on local metabolite concentrations — versus centralized "
    "allocation, where a planner assigns rates.\n\n"
    "The distributed regime shows 2.9x lower GDP variance and retains 71% of output under "
    "perturbation versus 53% for the planner. I also built a smart planner with real-time "
    "metabolite pool access — Mises' strongest opponent — and it still fails under perturbation.\n\n"
    "The perturbation suite tests four specific Austrian predictions: substrate shock tests "
    "Hayek's knowledge problem, ATP crisis tests Mises' calculation problem, and two Kirzner "
    "tests for entrepreneurial discovery and alertness.\n\n"
    "Layer 3 maps gene transferability across species as a trade network. Codon distance is "
    "trade friction. I validated against 18 real published gene transfers — GFP, insulin, "
    "silk fibroin — and lower trade cost correlates with higher success. Trade blocs emerge "
    "spontaneously from the compatibility data.",
    fig_name="L2_perturbation_robustness",
    fig_right=True,
)

# ===================================================================
# SLIDE 5: This Week's Progress — FBA (45 seconds)
# ===================================================================
add_slide(
    "This Week: FBA — The Omniscient Planner vs Reality",
    [
        "Flux Balance Analysis with iML1515:",
        "  - 2,712 reactions, 1,877 metabolites, 1,516 genes",
        "  - FBA = strongest possible central planner",
        "",
        "Knockout predictions vs Keio collection:",
        ">> 70% accuracy — planner fails on 12 genes",
        "  - False negatives = Hayek's knowledge gap",
        "",
        "Shadow prices = Hayekian price signals",
        "  - FBA must compute prices to solve",
    ],
    notes="THIS WEEK — FBA (45 seconds):\n\n"
    "The biggest addition this week is Flux Balance Analysis using the iML1515 genome-scale "
    "model of E. coli. This is the strongest possible central planner: a linear program that "
    "simultaneously optimizes across all 2,712 reactions with perfect knowledge of every "
    "stoichiometric constraint.\n\n"
    "I compared its gene knockout predictions against the Keio collection — systematic "
    "single-gene deletions with measured growth phenotypes. The planner achieves 70% accuracy. "
    "That's good but not perfect.\n\n"
    "The failures are the interesting part. 12 false negatives — genes where FBA says the cell "
    "will survive the knockout but it actually dies. These are genes that are essential because "
    "of REGULATORY requirements that the linear program can't encode — allosteric feedback, "
    "protein folding dependencies, gene expression timing. This is exactly Hayek's knowledge "
    "problem: the planner has all the stoichiometric data but misses the local, tacit knowledge.\n\n"
    "I also extracted shadow prices — the dual variables of the LP. These ARE Hayekian price "
    "signals. The irony is that FBA must compute these prices internally to solve its "
    "optimization, which proves Hayek's point that prices contain essential information.",
    fig_name="L2_fba_knockout_scatter",
    fig_right=True,
)

# ===================================================================
# SLIDE 6: This Week's Progress — Figures & Presentation (45 seconds)
# ===================================================================
add_slide(
    "This Week: FBA Perturbation & Methods",
    [
        "FBA perturbation analysis:",
        "  - Glucose->Acetate: 24% growth (Hayek)",
        "  - Aerobic->Anaerobic: 18% growth (Mises)",
        "  - Nitrogen limitation: 11% growth (Kirzner)",
        "",
        "Planner re-solves instantly;",
        "biology must discover via local feedback",
        "",
        "Methods section written (docx, Arial 11pt)",
        "  - 9 subsections, all accession numbers",
        "  - Data Availability: GitHub + public DBs",
    ],
    notes="THIS WEEK — FIGURES & METHODS (45 seconds):\n\n"
    "I also split all the multi-panel figures into individual PNGs — 21 separate figures, "
    "each one graph. This makes the presentation much cleaner.\n\n"
    "The FBA perturbation analysis tests three environmental shifts. When I switch the carbon "
    "source from glucose to acetate, FBA instantly re-optimizes — but real E. coli shows a "
    "diauxic lag phase. The planner assumes instant knowledge; biology has to discover the "
    "change through distributed sensing. Same pattern for anaerobic growth and nitrogen "
    "limitation — FBA computes the optimal answer in one step, but the cell's distributed "
    "regulatory network grows toward it through iterative local feedback.\n\n"
    "I wrote the methods section as a formal document — Arial 11pt, single-spaced, 9 "
    "subsections covering every database, accession number, and analytical method. Includes "
    "a table of all databases with URLs and accession numbers, and a Data Availability "
    "section pointing to GitHub.\n\n"
    "Everything is automated — one command generates all figures and the PowerPoint.",
    fig_name="L2_fba_pert_carbon",
    fig_right=True,
)

# ===================================================================
# SLIDE 7: Next Steps
# ===================================================================
add_slide(
    "Next Steps",
    [
        "Run full pipeline with live API data",
        "  - Full RegulonDB, KEGG, STRING networks",
        "  - 1,000 randomizations for motif analysis",
        "",
        "Finalize paper with FBA results",
        "  - Shadow prices + conditional essentiality",
        "",
        "Statistical rigor:",
        "  - Mann-Whitney U, bootstrap CIs",
        "",
        "Polish figures for final poster",
    ],
    notes="NEXT STEPS:\n\n"
    "Run the full pipeline with live API data instead of the built-in quick-mode data. "
    "This means pulling the full RegulonDB network, all KEGG pathways, and STRING PPI "
    "networks. The motif analysis will use 1,000 randomizations instead of 100.\n\n"
    "Finalize the paper text — the FBA results need to be integrated into the Results "
    "and Discussion sections. The shadow price analysis and conditional essentiality data "
    "add a whole new dimension to the Austrian argument.\n\n"
    "Add proper statistical tests — Mann-Whitney U for comparing distributed vs centralized "
    "GDP distributions, bootstrap confidence intervals on robustness metrics.\n\n"
    "And polish the figures for the final poster and presentation.",
    fig_name="L3_trade_network_graph",
    fig_right=True,
)


# Save
out_path = PROJECT_ROOT / "paper" / "progress_report_slides.pptx"
prs.save(str(out_path))
print(f"Progress report slides saved to: {out_path}")
