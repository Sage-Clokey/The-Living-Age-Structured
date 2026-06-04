"""
Build PowerPoint progress report slides for BME 129C Capstone — Week 6.
3 minutes: 0.5 min goal/approach, 1 min accomplishments, 1.5 min recent progress.

Week 6 (May 5–9, 2026):
  - Completed oral exam (Week 5) — presented all 6 questions conversationally
  - Built "The Price System of the Cell" — 4-panel synthesis figure
  - Created comprehensive figure guide mapping every claim to a falsifiable visualization
  - Added price system layer to run_all.py pipeline

Usage: python paper/build_slides_week6.py
"""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

FIGURES = Path(__file__).resolve().parent.parent / "figures"
INDIVIDUAL = FIGURES / "individual"
OUT = Path(__file__).resolve().parent.parent / "progress_report_slides_week6.pptx"

# Colors
GREEN_MID = RGBColor(0x52, 0xb7, 0x88)
GREEN_LIGHT = RGBColor(0x95, 0xd5, 0xb2)
GOLD = RGBColor(0xe9, 0xc4, 0x6a)
RED = RGBColor(0xe6, 0x39, 0x46)
WHITE = RGBColor(0xff, 0xff, 0xff)
OFF_WHITE = RGBColor(0xf5, 0xf5, 0xf5)
DARK_BG = RGBColor(0x0d, 0x11, 0x17)
GRAY = RGBColor(0x8b, 0x94, 0x9e)
BLUE = RGBColor(0x4e, 0xa8, 0xde)


def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text_box(slide, left, top, width, height, text, font_size=18,
                 color=WHITE, bold=False, alignment=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = "Calibri"
    p.alignment = alignment
    return txBox


def add_bullets(slide, bullets, left, top, width, height,
                font_size=16, color=WHITE, spacing=Pt(8)):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, (text, subs) in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = text
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Calibri"
        p.space_after = spacing
        p.level = 0
        for sub in subs:
            p2 = tf.add_paragraph()
            p2.text = sub
            p2.font.size = Pt(font_size - 2)
            p2.font.color.rgb = GRAY
            p2.font.name = "Calibri"
            p2.space_after = Pt(4)
            p2.level = 1
    return txBox


def add_green_bar(slide):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.33), Inches(0.06))
    shape.fill.solid()
    shape.fill.fore_color.rgb = GREEN_MID
    shape.line.fill.background()


def try_add_picture(slide, *candidates, left=0, top=0, width=None, height=None):
    for path in candidates:
        p = Path(path)
        if p.exists():
            slide.shapes.add_picture(str(p), left, top, width, height)
            return True
    return False


def build():
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    # =========================================================================
    # SLIDE 1: Title
    # =========================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)
    add_green_bar(slide)

    add_text_box(slide, Inches(1), Inches(1.2), Inches(11), Inches(1.2),
                 "The Price System of the Cell",
                 font_size=44, color=GREEN_MID, bold=True)
    add_text_box(slide, Inches(1), Inches(2.6), Inches(11), Inches(0.8),
                 "Life as a Decentralized Economy",
                 font_size=24, color=GOLD)
    add_text_box(slide, Inches(1), Inches(3.8), Inches(11), Inches(0.5),
                 "BME 129C Capstone  |  Week 6 Progress Report",
                 font_size=16, color=GRAY)
    add_text_box(slide, Inches(1), Inches(4.4), Inches(11), Inches(0.5),
                 "Sage Clokey   |   Spring 2026   |   UC Santa Cruz",
                 font_size=18, color=GREEN_LIGHT)
    add_text_box(slide, Inches(1), Inches(5.5), Inches(11), Inches(0.8),
                 "7 layers of evidence  |  40+ figures  |  Price system synthesis",
                 font_size=14, color=GRAY, alignment=PP_ALIGN.CENTER)

    notes = slide.notes_slide
    notes.notes_text_frame.text = (
        "Title slide. Week 6 progress report.\n\n"
        "Oral exam completed last week. This week's central accomplishment: "
        "the Price System of the Cell figure, a 4-panel synthesis that shows "
        "living systems have a real price system — not as metaphor but as "
        "measurable molecular mechanism. Also built the comprehensive figure guide."
    )

    # =========================================================================
    # SLIDE 2: Panel A — Price Tiers
    # =========================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)
    add_green_bar(slide)

    add_text_box(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.5),
                 "Panel A: The Cell Has a Real Price System",
                 font_size=28, color=GOLD, bold=True)

    try_add_picture(slide,
                    INDIVIDUAL / "panel_a_price_tiers.png",
                    FIGURES / "price_system_of_the_cell.png",
                    left=Inches(0.3), top=Inches(1.0),
                    width=Inches(6.5), height=Inches(5.8))

    bullets = [
        ("Three tiers of molecular prices:", [
            "Intracellular ratios = cost of capital",
            "  ATP/ADP, NAD+/NADH, AMP/ATP",
            "Intercellular signals = market prices",
            "  Cytokines, morphogens, growth factors, O2 tension",
        ]),
        ("mTOR = the entrepreneur", [
            "Reads all price inputs simultaneously",
            "Makes a grow-or-conserve decision",
            "No central planner tells it what to do",
        ]),
        ("Not metaphor — measurable molecular signals", [
            "Every ratio is quantifiable",
            "Every signal has a known receptor",
        ]),
    ]
    add_bullets(slide, bullets, Inches(7.2), Inches(1.0), Inches(5.8), Inches(5.8),
                font_size=14, color=OFF_WHITE)

    notes = slide.notes_slide
    notes.notes_text_frame.text = (
        "PANEL A — Price Tiers:\n\n"
        "The cell reads prices at three levels. Intracellular metabolite ratios "
        "are the cost of capital — they emerge from the cell's own activity. "
        "Intercellular signals are market prices — tissue-level information. "
        "mTOR integrates all of these and makes a grow-or-conserve decision. "
        "It is the entrepreneur reading the market."
    )

    # =========================================================================
    # SLIDE 4: Panel B — Shadow Prices (Subjective Value)
    # =========================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)
    add_green_bar(slide)

    add_text_box(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.5),
                 "Panel B: Menger's Subjective Value — Measured",
                 font_size=28, color=GOLD, bold=True)

    try_add_picture(slide,
                    INDIVIDUAL / "panel_b_shadow_prices.png",
                    FIGURES / "price_system_of_the_cell.png",
                    left=Inches(0.3), top=Inches(1.0),
                    width=Inches(6.5), height=Inches(5.8))

    bullets = [
        ("Same metabolite, different value in different conditions", [
            "NADH: low price on glucose, high price on acetate",
            "Oxygen: near zero aerobically, spikes under anaerobic",
            "Ammonium: irrelevant normally, critical under nitrogen limitation",
        ]),
        ("FBA shadow prices from iML1515 genome-scale model", [
            "2,712 reactions, 1,877 metabolites, 1,516 genes",
            "Shadow prices = dual variables from linear program",
            "Even the omniscient planner must compute prices to solve",
        ]),
        ("Menger (1871): value is subjective, not intrinsic", [
            "The molecule doesn't change — the context does",
            "The cell has known this for 4 billion years",
        ]),
    ]
    add_bullets(slide, bullets, Inches(7.2), Inches(1.0), Inches(5.8), Inches(5.8),
                font_size=14, color=OFF_WHITE)

    notes = slide.notes_slide
    notes.notes_text_frame.text = (
        "PANEL B — Shadow Prices:\n\n"
        "FBA dual variables from the full E. coli genome-scale model under four "
        "conditions. The same metabolite has completely different value depending "
        "on context. This IS Menger's subjective value theory, measured in a "
        "genome-scale model. The molecule doesn't change — the context does."
    )

    # =========================================================================
    # SLIDE 5: Panel C — Price Discovery
    # =========================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)
    add_green_bar(slide)

    add_text_box(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.5),
                 "Panel C: Price Discovery Without a Planner",
                 font_size=28, color=GOLD, bold=True)

    try_add_picture(slide,
                    INDIVIDUAL / "panel_c_price_discovery.png",
                    FIGURES / "price_system_of_the_cell.png",
                    left=Inches(0.3), top=Inches(1.0),
                    width=Inches(6.5), height=Inches(5.8))

    bullets = [
        ("Metabolite pools oscillate then converge", [
            "Early oscillation = discovery phase",
            "Late convergence = equilibrium found",
            "No planner told them where to settle",
        ]),
        ("Supply and demand, happening in a cell", [
            "Oversupply = producers slow down (price drops)",
            "Scarcity = producers speed up (price rises)",
            "Local feedback alone drives coordination",
        ]),
        ("The invisible hand at the molecular level", [
            "Each agent adjusts based on local conditions",
            "Global order emerges from individual decisions",
            "This is Hayek's price discovery, measured in metabolites",
        ]),
    ]
    add_bullets(slide, bullets, Inches(7.2), Inches(1.0), Inches(5.8), Inches(5.8),
                font_size=14, color=OFF_WHITE)

    notes = slide.notes_slide
    notes.notes_text_frame.text = (
        "PANEL C — Price Discovery:\n\n"
        "Metabolite pools oscillate then converge. Early volatility is the market "
        "discovering prices. Late stability is equilibrium found through local "
        "feedback alone. No planner needed. This is the invisible hand at the "
        "molecular level."
    )

    # =========================================================================
    # SLIDE 6: Panel D — Cancer Breaks the Price System
    # =========================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)
    add_green_bar(slide)

    add_text_box(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.5),
                 "Panel D: Cancer Breaks the Price System",
                 font_size=28, color=GOLD, bold=True)

    try_add_picture(slide,
                    INDIVIDUAL / "panel_d_cancer_price_system.png",
                    FIGURES / "price_system_of_the_cell.png",
                    left=Inches(0.3), top=Inches(1.0),
                    width=Inches(6.5), height=Inches(5.8))

    bullets = [
        ("Most mutated genes map to price system components", [
            "TP53 — damage price (apoptosis signal)",
            "PIK3CA, PTEN, mTOR — price integrators",
            "EGFR, ERBB2 — growth factor receptors (price readers)",
            "NF2 — spatial/contact price (Hippo pathway)",
        ]),
        ("The variation IS the signal", [
            "TP53: 96% in ovarian, 1% in thyroid",
            "Different tissues rely on different prices",
            "Same gene, completely different importance by context",
        ]),
        ("Cancer is not a gene disease", [
            "It's a price system disease",
            "The disease is the broken feedback loop",
            "Variants are not diseases — broken readers are",
        ]),
    ]
    add_bullets(slide, bullets, Inches(7.2), Inches(1.0), Inches(5.8), Inches(5.8),
                font_size=14, color=OFF_WHITE)

    add_text_box(slide, Inches(0.8), Inches(6.8), Inches(11.5), Inches(0.5),
                 "The cell has prices. Cancer is what happens when the price system breaks.",
                 font_size=15, color=GREEN_MID, bold=True)

    notes = slide.notes_slide
    notes.notes_text_frame.text = (
        "PANEL D — Cancer Mutations:\n\n"
        "The most mutated genes in cancer all map to specific components of "
        "the price system. Every dot is one cancer type. The variation across "
        "types proves it's context-dependent — TP53 is 96% in ovarian but 1% "
        "in thyroid. Different tissues break different price components. "
        "Cancer isn't a gene disease — it's a price system disease."
    )

    # =========================================================================
    # SLIDE 5: Next Steps (brief)
    # =========================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)
    add_green_bar(slide)

    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.6),
                 "Next Steps", font_size=32, color=GOLD, bold=True)

    nexts = [
        ("Poster layout", []),
        ("Paper finalization", []),
        ("Sensitivity analysis on simulation", []),
        ("Final presentation", []),
    ]
    add_bullets(slide, nexts, Inches(0.8), Inches(1.5), Inches(5), Inches(4),
                font_size=18)

    # Summary box on right
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(7.5), Inches(1.5), Inches(5), Inches(4.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0x1b, 0x43, 0x32)
    shape.line.color.rgb = GREEN_MID
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Week 6 Summary"
    p.font.size = Pt(18)
    p.font.color.rgb = GOLD
    p.font.bold = True
    p.font.name = "Calibri"
    p.alignment = PP_ALIGN.CENTER

    p2 = tf.add_paragraph()
    p2.text = (
        "\nThe cell has a price system.\n\n"
        "ATP/ADP is the cost of capital.\n"
        "Cytokines are market prices.\n"
        "mTOR is the entrepreneur.\n\n"
        "Cancer is what happens when\n"
        "the price system breaks."
    )
    p2.font.size = Pt(14)
    p2.font.color.rgb = GREEN_LIGHT
    p2.font.name = "Calibri"
    p2.alignment = PP_ALIGN.CENTER

    notes = slide.notes_slide
    notes.notes_text_frame.text = "Next steps. Keep it brief — 10 seconds."

    prs.save(str(OUT))
    print(f"Saved: {OUT}")
    print(f"  {len(prs.slides)} slides")


if __name__ == "__main__":
    build()
