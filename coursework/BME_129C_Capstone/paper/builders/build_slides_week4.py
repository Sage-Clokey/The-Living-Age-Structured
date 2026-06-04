"""
Build PowerPoint progress report slides for BME 129C Capstone — Week 4.
3 minutes: 0.5 min goal/approach, 1 min accomplishments, 1.5 min recent progress.

Week 4 (April 14–18, 2026):
  - Integrated Hayek's "The Fatal Conceit" (1988) across paper and slides
  - Wrote Design Principles section — two architectures for sagent engineering
  - Designed experimental spread — 6 violacein pathway experiments
  - Overhauled all figure code — every data point shown, no hidden averages
  - Fixed price_signals bug, regenerated full pipeline

Usage: python paper/build_slides_week4.py
"""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

FIGURES = Path(__file__).resolve().parent.parent / "figures"
INDIVIDUAL = FIGURES / "individual"
OUT = Path(__file__).resolve().parent.parent / "progress_report_slides_week4.pptx"

# Colors
GREEN_MID = RGBColor(0x52, 0xb7, 0x88)
GREEN_LIGHT = RGBColor(0x95, 0xd5, 0xb2)
GOLD = RGBColor(0xe9, 0xc4, 0x6a)
RED = RGBColor(0xe6, 0x39, 0x46)
WHITE = RGBColor(0xff, 0xff, 0xff)
OFF_WHITE = RGBColor(0xf5, 0xf5, 0xf5)
DARK_BG = RGBColor(0x0d, 0x11, 0x17)
GRAY = RGBColor(0x8b, 0x94, 0x9e)


def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text_box(slide, left, top, width, height, text, font_size=18,
                 color=WHITE, bold=False, alignment=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = "Calibri"
    p.alignment = alignment
    return txBox


def add_bullets(slide, bullets, left, top, width, height,
                font_size=16, color=WHITE, spacing=Pt(8)):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, (text, subs) in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = text
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Calibri"
        p.space_after = spacing
        p.level = 0
        for sub in subs:
            p2 = tf.add_paragraph()
            p2.text = sub
            p2.font.size = Pt(font_size - 2)
            p2.font.color.rgb = GRAY
            p2.font.name = "Calibri"
            p2.space_after = Pt(4)
            p2.level = 1
    return txBox


def add_green_bar(slide):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.33), Inches(0.06))
    shape.fill.solid()
    shape.fill.fore_color.rgb = GREEN_MID
    shape.line.fill.background()


def try_add_picture(slide, *candidates, left=0, top=0, width=None, height=None):
    """Try multiple figure paths, add the first one that exists."""
    for path in candidates:
        p = Path(path)
        if p.exists():
            slide.shapes.add_picture(str(p), left, top, width, height)
            return True
    return False


def build():
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    # =========================================================================
    # SLIDE 1: Title
    # =========================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)
    add_green_bar(slide)

    add_text_box(slide, Inches(1), Inches(1.5), Inches(11), Inches(1.2),
                 "Living Systems as Decentralized Economies",
                 font_size=40, color=GREEN_MID, bold=True)
    add_text_box(slide, Inches(1), Inches(2.8), Inches(11), Inches(0.8),
                 "The Curious Task: From Data to Design Principles",
                 font_size=22, color=GOLD)
    add_text_box(slide, Inches(1), Inches(4.2), Inches(11), Inches(0.5),
                 "BME 129C Capstone  |  Week 4 Progress Report",
                 font_size=16, color=GRAY)
    add_text_box(slide, Inches(1), Inches(4.8), Inches(11), Inches(0.5),
                 "Sage Clokey   |   Spring 2026   |   UC Santa Cruz",
                 font_size=18, color=GREEN_LIGHT)

    notes = slide.notes_slide
    notes.notes_text_frame.text = (
        "Title slide. Week 4 progress report.\n\n"
        "The subtitle references Hayek's Fatal Conceit (1988): 'The curious task of "
        "economics is to demonstrate to men how little they really know about what they "
        "imagine they can design.' This week I integrated that work into the paper and "
        "translated the computational findings into concrete design principles for "
        "synthetic biology."
    )

    # =========================================================================
    # SLIDE 2: Goal & Approach (0.5 min)
    # =========================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)
    add_green_bar(slide)

    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(6), Inches(0.6),
                 "Goal & Approach", font_size=32, color=GOLD, bold=True)

    bullets = [
        ("Thesis: living systems operate as decentralized economies", [
            "Not metaphor — measurable structural equivalence"
        ]),
        ("Hayek's Fatal Conceit (1988): 'order generated without design", [
            "can far outstrip plans men consciously contrive' (p. 6)"
        ]),
        ("Three layers of convergent evidence:", [
            "Layer 1: Network topology — no master node (RegulonDB, STRING, KEGG)",
            "Layer 1b: Single-cell economy — cells as specialized agents (pbmc3k)",
            "Layer 2: Economic simulation + FBA — distributed vs centralized",
            "Layer 3: Cross-species trade — comparative advantage via codon distance",
        ]),
        ("NEW: Design principles — how to engineer with this, not against it", []),
    ]
    add_bullets(slide, bullets, Inches(0.8), Inches(1.2), Inches(5.5), Inches(5.5),
                font_size=16)

    try_add_picture(slide,
                    INDIVIDUAL / "L1_robustness_curves.png",
                    FIGURES / "layer1_topology.png",
                    left=Inches(7), top=Inches(1.2),
                    width=Inches(5.8), height=Inches(5.0))

    notes = slide.notes_slide
    notes.notes_text_frame.text = (
        "GOAL (30 seconds):\n\n"
        "Same core thesis as previous weeks — living systems as decentralized economies. "
        "But this week the project crossed from analysis into synthesis. The three layers "
        "of data (topology, economy, trade) are complete. Now I'm translating those "
        "findings into actionable design principles for synthetic biology.\n\n"
        "The key addition: Hayek's Fatal Conceit. His argument that extended orders "
        "'generated without design can far outstrip plans men consciously contrive' maps "
        "directly onto the 19:1 robustness ratio we measured in Layer 1."
    )

    # =========================================================================
    # SLIDE 3: Accomplishments — All Layers Complete (1 min)
    # =========================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)
    add_green_bar(slide)

    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.6),
                 "Accomplished — All Layers Complete",
                 font_size=32, color=GOLD, bold=True)

    left = [
        ("Layer 1: Gene regulatory, metabolic, PPI networks", [
            "5 bio networks vs 5 reference architectures",
            "19:1 robustness ratio (distributed vs centralized)",
            "Feed-forward loops = Hayekian price signals (Z > 10)",
            "WBPA: hubs that grow too powerful get eroded",
        ]),
        ("Layer 1b: Single-cell economy (PBMC 3k)", [
            "2,638 cells, 8 types — division of labor",
            "Communication Gini = 0.0 (no gatekeeper)",
            "75% communication survives any cell type removal",
        ]),
    ]
    add_bullets(slide, left, Inches(0.8), Inches(1.2), Inches(5.5), Inches(5.0),
                font_size=13, color=OFF_WHITE)

    right = [
        ("Layer 2: Economic simulation + FBA", [
            "13 pathway agents, distributed: 2.9x lower variance",
            "71% GDP retained under perturbation (vs 53%)",
            "FBA (omniscient planner): 70% knockout accuracy",
            "Shadow prices = Hayekian price signals",
        ]),
        ("Layer 3: Cross-species trade (8 organisms)", [
            "Codon distance = trade friction",
            "18 published gene transfers validated",
            "Trade blocs emerge spontaneously (Louvain)",
        ]),
    ]
    add_bullets(slide, right, Inches(7), Inches(1.2), Inches(5.8), Inches(5.0),
                font_size=13, color=OFF_WHITE)

    add_text_box(slide, Inches(0.8), Inches(6.5), Inches(11.5), Inches(0.5),
                 "Full paper draft  |  21 figures  |  Reproducible pipeline (run_all.py)  |  Methods section complete",
                 font_size=13, color=GREEN_MID)

    notes = slide.notes_slide
    notes.notes_text_frame.text = (
        "ACCOMPLISHMENTS (1 minute — split across slides 3 and 4):\n\n"
        "All four layers are complete and producing figures. The key numbers:\n"
        "- 19:1 robustness ratio favoring distributed networks\n"
        "- Communication Gini of 0.0 in PBMC cell economy\n"
        "- 2.9x lower GDP variance in distributed vs centralized allocation\n"
        "- FBA (the strongest possible planner) achieves 70% accuracy on knockouts\n"
        "- Trade blocs emerge spontaneously from codon compatibility\n\n"
        "Every layer tells the same story from a different angle: biology is a "
        "decentralized economy, not a centrally planned machine."
    )

    # =========================================================================
    # SLIDE 4: This Week — Fatal Conceit & Design Principles (45 sec)
    # =========================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)
    add_green_bar(slide)

    add_text_box(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.5),
                 "This Week: The Fatal Conceit → Design Principles",
                 font_size=28, color=GOLD, bold=True)

    bullets = [
        ("Integrated Hayek's Fatal Conceit (1988) into paper", [
            "'The curious task of economics is to demonstrate to men how little they "
            "really know about what they imagine they can design' (p. 76)",
            "Applied to synthetic biology: redesigning evolved order destroys information",
        ]),
        ("Wrote Section 4.2: Design Principles — Engineering Economies, Not Machines", [
            "Architecture 1: Single organism designed FOR evolution",
            "Architecture 2: Microbial consortium — division of labor",
            "The core shift: read the economy before intervening",
        ]),
        ("Designed Section 4.7: Experimental Spread — 6 Violacein Experiments", [
            "Violacein (vioABCDE): purple pigment, 5 genes, quantifiable at 575nm",
            "Each experiment tests one Austrian prediction with wet-lab data",
        ]),
    ]
    add_bullets(slide, bullets, Inches(0.8), Inches(1.0), Inches(11.5), Inches(5.5),
                font_size=15)

    notes = slide.notes_slide
    notes.notes_text_frame.text = (
        "THIS WEEK — FATAL CONCEIT & DESIGN PRINCIPLES (45 seconds):\n\n"
        "The biggest intellectual addition this week is Hayek's Fatal Conceit — his "
        "final book, published in 1988. The core argument: 'The curious task of economics "
        "is to demonstrate to men how little they really know about what they imagine they "
        "can design.' I integrated this into the paper's introduction, discussion, and "
        "conclusion.\n\n"
        "Then I translated the computational findings into two concrete design architectures:\n\n"
        "Architecture 1: Single organism designed FOR evolution. Instead of codon-optimizing "
        "a pathway, you read the shadow prices from FBA, build feedback loops so the cell "
        "self-regulates expression, distribute control across multiple promoters, harmonize "
        "codons instead of optimizing them, and then let the cell evolve.\n\n"
        "Architecture 2: Microbial consortium. Split a pathway across multiple strains, "
        "each specializing in part of the production. This mirrors the division of labor "
        "we measured in the PBMC data.\n\n"
        "Finally, I designed 6 concrete experiments using the violacein pathway (vioABCDE "
        "from Chromobacterium violaceum). Violacein is a purple pigment quantifiable by "
        "absorbance at 575nm — cheap to measure, clear signal. Each experiment tests one "
        "specific Austrian prediction."
    )

    # =========================================================================
    # SLIDE 5: Experimental Spread Detail (45 sec)
    # =========================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)
    add_green_bar(slide)

    add_text_box(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.5),
                 "Experimental Spread: 6 Tests of Distributed Design",
                 font_size=28, color=GOLD, bold=True)

    left_exp = [
        ("1. Fixed vs Feedback Promoters (Hayek)", [
            "Constitutive vs tryptophan-responsive — knowledge problem",
        ]),
        ("2. Star vs Distributed Control (Mises)", [
            "Single master promoter vs 5 independent — calculation problem",
        ]),
        ("3. Single Strain vs Consortium (Menger)", [
            "Monolithic vs 2-strain split — spontaneous order",
        ]),
    ]
    add_bullets(slide, left_exp, Inches(0.5), Inches(1.0), Inches(5.8), Inches(5.5),
                font_size=14, color=OFF_WHITE)

    right_exp = [
        ("4. Codon Optimization vs Harmonization (Layer 3)", [
            "CAI-maximized vs host-matched — trade barriers",
        ]),
        ("5. Adaptive Evolution: Fixed vs Feedback (Kirzner)", [
            "100 generations serial passage — entrepreneurial discovery",
        ]),
        ("6. Perturbation Gauntlet (Rothbard)", [
            "Temperature, carbon, antibiotic — fragility test",
        ]),
    ]
    add_bullets(slide, right_exp, Inches(6.8), Inches(1.0), Inches(5.8), Inches(5.5),
                font_size=14, color=OFF_WHITE)

    # Bottom summary
    add_text_box(slide, Inches(0.8), Inches(5.8), Inches(11.5), Inches(1.2),
                 "Prediction: distributed designs show higher VARIATION but equal or better "
                 "MEDIAN yield — and dramatically better robustness under perturbation.\n"
                 "Measurement: violacein absorbance at 575nm across n=8+ biological replicates per condition.",
                 font_size=13, color=GREEN_LIGHT)

    notes = slide.notes_slide
    notes.notes_text_frame.text = (
        "EXPERIMENTAL SPREAD (45 seconds):\n\n"
        "Six experiments, each testing one Austrian economic prediction:\n\n"
        "1. Fixed vs Feedback (Hayek's knowledge problem): Compare constitutive promoters "
        "(the planner sets expression levels once) against tryptophan-responsive promoters "
        "(the cell reads its own metabolite levels and adjusts). Prediction: feedback "
        "variants show wider yield distributions but better recovery after perturbation.\n\n"
        "2. Star vs Distributed (Mises' calculation problem): Put all 5 vio genes under "
        "one master promoter (star topology) vs independent promoters. Prediction: star "
        "fails catastrophically under any single perturbation.\n\n"
        "3. Monoculture vs Consortium (Menger's spontaneous order): All genes in one "
        "strain vs split across two co-cultured strains. Prediction: consortium shows "
        "emergent ratio optimization without tuning.\n\n"
        "4. Codon Optimization vs Harmonization: CAI-maximized (planner rewrites all codons) "
        "vs host-matched harmonization (respects the host's economy). Tests Layer 3 trade "
        "barrier predictions.\n\n"
        "5. Adaptive Evolution: Serial passage both architectures for 100 generations. "
        "Feedback variants should improve faster — Kirznerian discovery.\n\n"
        "6. Perturbation Gauntlet: Hit all constructs with temperature shift, carbon "
        "source switch, and sub-MIC antibiotic. Distributed designs should retain more "
        "production. Rothbard's fragility argument.\n\n"
        "All measured by violacein absorbance at 575nm, n=8+ biological replicates."
    )

    # =========================================================================
    # SLIDE 6: Visualization Overhaul (brief)
    # =========================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)
    add_green_bar(slide)

    add_text_box(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.5),
                 "This Week: Data-Point-Forward Visualization",
                 font_size=28, color=GOLD, bold=True)

    bullets = [
        ("New principle: show every individual data point", [
            "Variation IS the data — averages hide what matters",
            "Bar charts → semi-transparent bars with scatter overlay",
            "Robustness curves → individual trial lines behind average",
            "Hub erosion → per-hub jitter plots with median lines",
        ]),
        ("Updated 6 analysis files across all layers", [
            "topology_analysis.py, self_regulation.py, cell_economy_figures.py",
            "metabolic_economy.py, single_cell_economy.py, fba_analysis.py",
        ]),
        ("Regenerated all 21 figures through full pipeline", []),
    ]
    add_bullets(slide, bullets, Inches(0.8), Inches(1.0), Inches(5.5), Inches(5.5),
                font_size=15)

    # Show example figures: robustness curves + perturbation
    try_add_picture(slide,
                    INDIVIDUAL / "L1_robustness_curves.png",
                    FIGURES / "layer1_topology.png",
                    left=Inches(7), top=Inches(1.0),
                    width=Inches(5.8), height=Inches(2.8))

    try_add_picture(slide,
                    INDIVIDUAL / "L2_perturbation_robustness.png",
                    FIGURES / "layer2_economy.png",
                    left=Inches(7), top=Inches(4.0),
                    width=Inches(5.8), height=Inches(2.8))

    notes = slide.notes_slide
    notes.notes_text_frame.text = (
        "VISUALIZATION OVERHAUL:\n\n"
        "Overhauled every plotting function across the entire codebase. The core "
        "principle: variation is the signal, not noise. When we show that distributed "
        "systems produce more variation that converges through local feedback, hiding "
        "that variation behind averages undermines the thesis.\n\n"
        "Every bar chart now shows semi-transparent bars with individual data points "
        "scattered on top. Robustness curves show each random trial as a thin line "
        "behind the average. Hub erosion plots show per-hub betweenness changes as "
        "jittered scatter with median lines.\n\n"
        "This affected 6 files and all 21 figures. The pipeline regenerates everything "
        "from one command."
    )

    # =========================================================================
    # SLIDE 7: Next Steps
    # =========================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)
    add_green_bar(slide)

    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.6),
                 "Next Steps", font_size=32, color=GOLD, bold=True)

    nexts = [
        ("Statistical validation", [
            "Mann-Whitney U for distributed vs centralized GDP distributions",
            "Bootstrap confidence intervals on robustness metrics",
            "KS test for power-law degree distribution fit",
        ]),
        ("Run full pipeline with live API data", [
            "Full RegulonDB, KEGG, STRING networks",
            "1,000 randomizations for motif analysis",
        ]),
        ("Finalize paper for advisor review", [
            "Integrate Fatal Conceit quotes into remaining sections",
            "Polish experimental spread section for feasibility feedback",
        ]),
        ("Begin poster layout", [
            "Three-layer visual with key figures and one-line findings",
        ]),
    ]
    add_bullets(slide, nexts, Inches(0.8), Inches(1.3), Inches(5.5), Inches(5.5),
                font_size=16)

    # Thesis box on right
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(7.5), Inches(1.5), Inches(5), Inches(4.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0x1b, 0x43, 0x32)
    shape.line.color.rgb = GREEN_MID
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Week 4 Summary"
    p.font.size = Pt(18)
    p.font.color.rgb = GOLD
    p.font.bold = True
    p.font.name = "Calibri"
    p.alignment = PP_ALIGN.CENTER

    p2 = tf.add_paragraph()
    p2.text = (
        "\nThe data says biology is an economy.\n\n"
        "This week: what that means for\n"
        "how we engineer.\n\n"
        "Not optimize — cultivate.\n"
        "Not command — steward.\n"
        "Not design — grow.\n\n"
        "The Fatal Conceit is thinking\n"
        "we can do better than\n"
        "4 billion years of distributed design."
    )
    p2.font.size = Pt(14)
    p2.font.color.rgb = GREEN_LIGHT
    p2.font.name = "Calibri"
    p2.alignment = PP_ALIGN.CENTER

    notes = slide.notes_slide
    notes.notes_text_frame.text = (
        "NEXT STEPS:\n\n"
        "1. Statistical validation — add proper nonparametric tests to every comparison. "
        "Mann-Whitney U for distributed vs centralized GDP. Bootstrap CIs on robustness "
        "metrics. KS test for power-law fit.\n\n"
        "2. Full pipeline run with live API data. Currently using built-in quick-mode "
        "data for speed. The full run pulls from RegulonDB, KEGG, and STRING directly.\n\n"
        "3. Finalize paper for advisor review. Todd Lowe and Josh Stuart need to see "
        "the experimental spread — it needs to be feasible within the resources available.\n\n"
        "4. Begin poster layout. The three-layer structure maps naturally to a visual "
        "hierarchy: topology → economy → trade."
    )

    prs.save(str(OUT))
    print(f"Saved: {OUT}")


if __name__ == "__main__":
    build()
